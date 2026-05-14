#!/usr/bin/env python3
"""High-density PE deck v2 — target 35-40 dense pages with charts."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
import re

# ============================================================
# Theme · Midnight Executive
# ============================================================
class T:
    NAVY = RGBColor(0x0F, 0x27, 0x40)
    NAVY_DARK = RGBColor(0x08, 0x18, 0x2C)
    NAVY_LIGHT = RGBColor(0x19, 0x37, 0x6D)
    ICE = RGBColor(0xCA, 0xDC, 0xFC)
    GOLD = RGBColor(0xD4, 0xA5, 0x74)
    GOLD_LIGHT = RGBColor(0xF6, 0xB1, 0x7A)
    TEXT = RGBColor(0x1A, 0x1A, 0x1A)
    MUTED = RGBColor(0x6B, 0x72, 0x80)
    LIGHT_MUTED = RGBColor(0x9C, 0xA3, 0xAF)
    INVERT = RGBColor(0xFF, 0xFF, 0xFF)
    BG = RGBColor(0xFA, 0xFB, 0xFC)
    CARD = RGBColor(0xFF, 0xFF, 0xFF)
    SUBTLE = RGBColor(0xF1, 0xF4, 0xF8)
    DIVIDER = RGBColor(0xE5, 0xE7, 0xEB)
    SUCCESS = RGBColor(0x16, 0xA3, 0x4A)
    WARN = RGBColor(0xEA, 0x58, 0x0C)
    DANGER = RGBColor(0xDC, 0x26, 0x26)
    INFO = RGBColor(0x25, 0x63, 0xEB)
    PURPLE = RGBColor(0x7C, 0x3A, 0xED)
    W = Inches(13.333)
    H = Inches(7.5)
    F = "PingFang SC"

# ============================================================
# Primitive helpers
# ============================================================
def bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color

def rect(slide, left, top, width, height, color, line=False, line_color=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    if line and line_color:
        s.line.color.rgb = line_color
        s.line.width = Pt(0.5)
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    return s

def rrect(slide, left, top, width, height, color, radius=0.04, line_color=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.adjustments[0] = radius
    s.fill.solid()
    s.fill.fore_color.rgb = color
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = Pt(0.5)
    else:
        s.line.fill.background()
    s.shadow.inherit = False
    return s

def tx(slide, left, top, width, height, text, *,
       size=12, bold=False, color=None, italic=False,
       align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, name=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    render(p, text, size, color or T.TEXT, bold, italic, name or T.F)
    return tb

def add_run(p, text, size, color, bold=False, italic=False, name=None):
    r = p.add_run()
    r.text = text
    f = r.font
    f.name = name or T.F
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color

def render(paragraph, text, size, color, bold=False, italic=False, name=None):
    """Inline markup: {bold}/{blue}/{gold}/{red}/{green}/{italic}/{navy}/{small}"""
    pattern = re.compile(r'\{(bold|blue|navy|italic|gold|red|green|orange|purple|small|big)\}(.*?)\{/\1\}', re.DOTALL)
    pos = 0
    for m in pattern.finditer(text):
        if m.start() > pos:
            add_run(paragraph, text[pos:m.start()], size, color, bold, italic, name)
        tag = m.group(1); content = m.group(2)
        is_bold = bold or tag == 'bold'
        is_italic = italic or tag == 'italic'
        c = color; sz = size
        if tag == 'blue' or tag == 'navy': c = T.NAVY
        elif tag == 'gold': c = T.GOLD
        elif tag == 'red': c = T.DANGER
        elif tag == 'green': c = T.SUCCESS
        elif tag == 'orange': c = T.WARN
        elif tag == 'purple': c = T.PURPLE
        elif tag == 'small': sz = max(8, size - 2)
        elif tag == 'big': sz = size + 2
        add_run(paragraph, content, sz, c, is_bold, is_italic, name)
        pos = m.end()
    if pos < len(text):
        add_run(paragraph, text[pos:], size, color, bold, italic, name)

def multi_tx(slide, left, top, width, height, items, size=11, color=None, spacing=3, bullet="▸ "):
    """Add multi-paragraph text box with bullets."""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        if i == 0 and tf.paragraphs[0].text == "":
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(spacing)
        prefix = bullet if item and not item.startswith('§') else ""
        if item.startswith('§'):  # subhead
            render(p, item[1:].strip(), size + 1, T.NAVY, bold=True)
        elif item == "":
            pass
        else:
            render(p, prefix + item, size, color or T.TEXT)
    return tb

# ============================================================
# Page chrome
# ============================================================
def chrome(slide, page, total, section="", on_dark=False):
    accent = T.GOLD
    rect(slide, Inches(0), Inches(0), T.W, Inches(0.06), accent)
    if section:
        c = T.INVERT if on_dark else T.MUTED
        tx(slide, Inches(8), Inches(0.18), Inches(5), Inches(0.3),
           section, size=9, color=c, align=PP_ALIGN.RIGHT, bold=True)
    c = T.INVERT if on_dark else T.MUTED
    tx(slide, Inches(11.7), Inches(7.1), Inches(1.5), Inches(0.3),
       f"{page} / {total}", size=9, color=c, align=PP_ALIGN.RIGHT)
    tx(slide, Inches(0.4), Inches(7.1), Inches(7), Inches(0.3),
       "STRICTLY PRIVATE & CONFIDENTIAL", size=8, color=c)

def page_title(slide, title, subtitle=""):
    """Compact title bar (taking less vertical space)."""
    tx(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.55),
       title, size=22, bold=True, color=T.NAVY)
    if subtitle:
        tx(slide, Inches(0.5), Inches(0.85), Inches(12.3), Inches(0.35),
           subtitle, size=11, color=T.MUTED, italic=True)
    rect(slide, Inches(0.5), Inches(1.22), Inches(0.5), Inches(0.04), T.GOLD)

# ============================================================
# Native chart helpers
# ============================================================
def add_pie(slide, left, top, width, height, categories, values, colors=None, title="", show_pct=True, show_legend=True):
    """Pie chart with custom colors."""
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series("Series 1", values)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, left, top, width, height, chart_data
    ).chart
    chart.has_title = bool(title)
    if title:
        chart.chart_title.text_frame.text = title
        for run in chart.chart_title.text_frame.paragraphs[0].runs:
            run.font.size = Pt(11)
            run.font.bold = True
            run.font.name = T.F
            run.font.color.rgb = T.NAVY
    chart.has_legend = show_legend
    if show_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(9)
        chart.legend.font.name = T.F

    plot = chart.plots[0]
    plot.has_data_labels = show_pct
    if show_pct:
        dl = plot.data_labels
        dl.number_format = '0%'
        dl.show_percentage = True
        dl.show_category_name = False
        dl.show_value = False
        dl.font.size = Pt(9)
        dl.font.name = T.F
        dl.font.bold = True
        dl.font.color.rgb = T.INVERT

    # Color slices
    if colors:
        for i, point in enumerate(plot.series[0].points):
            if i < len(colors):
                point.format.fill.solid()
                point.format.fill.fore_color.rgb = colors[i]
                point.format.line.color.rgb = T.INVERT
                point.format.line.width = Pt(1.5)
    return chart

def add_bar_h(slide, left, top, width, height, categories, values,
              colors=None, title="", show_values=True, max_value=None):
    """Horizontal bar chart."""
    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series("Series 1", values)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, left, top, width, height, chart_data
    ).chart
    chart.has_title = bool(title)
    if title:
        chart.chart_title.text_frame.text = title
        for run in chart.chart_title.text_frame.paragraphs[0].runs:
            run.font.size = Pt(11)
            run.font.bold = True
            run.font.name = T.F
            run.font.color.rgb = T.NAVY
    chart.has_legend = False

    plot = chart.plots[0]
    plot.has_data_labels = show_values
    if show_values:
        dl = plot.data_labels
        dl.show_value = True
        dl.font.size = Pt(9)
        dl.font.name = T.F
        dl.font.bold = True
        dl.font.color.rgb = T.NAVY
        dl.position = XL_LABEL_POSITION.OUTSIDE_END

    # Color bars
    series = plot.series[0]
    if colors:
        for i, point in enumerate(series.points):
            if i < len(colors):
                point.format.fill.solid()
                point.format.fill.fore_color.rgb = colors[i]
                point.format.line.fill.background()
    else:
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = T.NAVY

    # Format axes
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(9)
    cat_axis.tick_labels.font.name = T.F
    val_axis = chart.value_axis
    val_axis.tick_labels.font.size = Pt(8)
    val_axis.tick_labels.font.name = T.F
    if max_value:
        val_axis.maximum_scale = max_value

    return chart

def add_bar_v(slide, left, top, width, height, categories, series_data,
              colors=None, title="", stacked=False):
    """Vertical bar chart. series_data: list of (name, values)."""
    chart_data = CategoryChartData()
    chart_data.categories = categories
    for name, values in series_data:
        chart_data.add_series(name, values)

    ctype = XL_CHART_TYPE.COLUMN_STACKED if stacked else XL_CHART_TYPE.COLUMN_CLUSTERED
    chart = slide.shapes.add_chart(
        ctype, left, top, width, height, chart_data
    ).chart
    chart.has_title = bool(title)
    if title:
        chart.chart_title.text_frame.text = title
        for run in chart.chart_title.text_frame.paragraphs[0].runs:
            run.font.size = Pt(11)
            run.font.bold = True
            run.font.name = T.F
            run.font.color.rgb = T.NAVY
    chart.has_legend = len(series_data) > 1
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(9)
        chart.legend.font.name = T.F

    # Color series
    plot = chart.plots[0]
    if colors:
        for i, s in enumerate(plot.series):
            if i < len(colors):
                s.format.fill.solid()
                s.format.fill.fore_color.rgb = colors[i]
                s.format.line.fill.background()

    chart.category_axis.tick_labels.font.size = Pt(9)
    chart.category_axis.tick_labels.font.name = T.F
    chart.value_axis.tick_labels.font.size = Pt(8)
    chart.value_axis.tick_labels.font.name = T.F
    return chart

def add_line(slide, left, top, width, height, categories, series_data,
             colors=None, title=""):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    for name, values in series_data:
        chart_data.add_series(name, values)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, left, top, width, height, chart_data
    ).chart
    chart.has_title = bool(title)
    if title:
        chart.chart_title.text_frame.text = title
        for run in chart.chart_title.text_frame.paragraphs[0].runs:
            run.font.size = Pt(11); run.font.bold = True; run.font.name = T.F
            run.font.color.rgb = T.NAVY
    chart.has_legend = len(series_data) > 1
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(9)
        chart.legend.font.name = T.F
    plot = chart.plots[0]
    if colors:
        for i, s in enumerate(plot.series):
            if i < len(colors):
                s.format.line.color.rgb = colors[i]
                s.format.line.width = Pt(2.5)
                s.marker.format.fill.solid()
                s.marker.format.fill.fore_color.rgb = colors[i]
                s.marker.size = 6
    chart.category_axis.tick_labels.font.size = Pt(9)
    chart.category_axis.tick_labels.font.name = T.F
    chart.value_axis.tick_labels.font.size = Pt(8)
    chart.value_axis.tick_labels.font.name = T.F
    return chart

def add_radar(slide, left, top, width, height, categories, series_data, colors=None, title=""):
    chart_data = CategoryChartData()
    chart_data.categories = categories
    for name, values in series_data:
        chart_data.add_series(name, values)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.RADAR, left, top, width, height, chart_data
    ).chart
    chart.has_title = bool(title)
    if title:
        chart.chart_title.text_frame.text = title
        for run in chart.chart_title.text_frame.paragraphs[0].runs:
            run.font.size = Pt(11); run.font.bold = True; run.font.name = T.F
            run.font.color.rgb = T.NAVY
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(9)
    chart.legend.font.name = T.F
    plot = chart.plots[0]
    if colors:
        for i, s in enumerate(plot.series):
            if i < len(colors):
                s.format.line.color.rgb = colors[i]
                s.format.line.width = Pt(2)
    chart.category_axis.tick_labels.font.size = Pt(8)
    chart.category_axis.tick_labels.font.name = T.F
    return chart

# ============================================================
# Table helper (dense)
# ============================================================
def add_table(slide, left, top, width, height, headers, rows, *,
              col_widths=None, header_color=None, alt_rows=True,
              header_font_size=9, body_font_size=8, row_height=None):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = shape.table
    if col_widths:
        total = sum(col_widths)
        for j, w in enumerate(col_widths):
            table.columns[j].width = Emu(int(width * (w / total)))
    if row_height:
        for r in table.rows:
            r.height = row_height

    # Header
    hc = header_color or T.TABLE_HEAD if hasattr(T, 'TABLE_HEAD') else T.NAVY
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = T.NAVY
        tf = cell.text_frame
        tf.margin_left = Inches(0.06); tf.margin_right = Inches(0.06)
        tf.margin_top = Inches(0.03); tf.margin_bottom = Inches(0.03)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        render(p, str(h), header_font_size, T.INVERT, bold=True)

    # Body
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row[:n_cols]):
            cell = table.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = T.SUBTLE if alt_rows and i % 2 == 0 else T.CARD
            tf = cell.text_frame
            tf.margin_left = Inches(0.06); tf.margin_right = Inches(0.06)
            tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            render(p, str(val), body_font_size, T.TEXT)
    return shape

# ============================================================
# Section accent + KPI chip
# ============================================================
def section_label(slide, left, top, text, color=None):
    """Small section signpost above content."""
    c = color or T.GOLD
    rect(slide, left, top + Inches(0.06), Inches(0.25), Inches(0.04), c)
    tx(slide, left + Inches(0.3), top, Inches(4), Inches(0.25),
       text, size=10, bold=True, color=c)

def kpi_chip(slide, left, top, width, height, value, label, sub="", accent=None):
    """Single KPI as compact chip."""
    rrect(slide, left, top, width, height, T.CARD, radius=0.05)
    rect(slide, left, top, Inches(0.06), height, accent or T.GOLD)
    tx(slide, left + Inches(0.15), top + Inches(0.1), width - Inches(0.2), Inches(0.55),
       value, size=22, bold=True, color=T.NAVY)
    tx(slide, left + Inches(0.15), top + Inches(0.7), width - Inches(0.2), Inches(0.3),
       label, size=9, bold=True, color=T.TEXT)
    if sub:
        tx(slide, left + Inches(0.15), top + Inches(0.95), width - Inches(0.2), Inches(0.3),
           sub, size=7, color=T.MUTED, italic=True)

# ============================================================
# Slide templates · HIGH DENSITY
# ============================================================
def slide_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, T.NAVY_DARK)
    rect(s, Inches(0), Inches(0), Inches(0.3), T.H, T.GOLD)
    tx(s, Inches(1), Inches(0.6), Inches(11), Inches(0.4),
       "PRIVATE EQUITY · STRICTLY PRIVATE AND CONFIDENTIAL",
       size=11, bold=True, color=T.GOLD)
    tx(s, Inches(1), Inches(1.6), Inches(11.5), Inches(2.2),
       "全球具身智能产业链格局分析", size=48, bold=True, color=T.INVERT)
    tx(s, Inches(1), Inches(3.0), Inches(11.5), Inches(0.8),
       "Global Embodied AI Industry Chain Analysis", size=18, color=T.ICE_BLUE if hasattr(T, 'ICE_BLUE') else T.ICE)
    tx(s, Inches(1), Inches(4.2), Inches(11.5), Inches(0.6),
       "人形机器人 与 具身智能 产业 PE 视角研究", size=20, color=T.ICE)
    rect(s, Inches(1), Inches(5.0), Inches(2), Inches(0.04), T.GOLD)
    tx(s, Inches(1), Inches(5.2), Inches(11), Inches(0.4),
       "Version 10 · 76 → 35 高密度页", size=13, color=T.GOLD_LIGHT, italic=True)
    tx(s, Inches(1), Inches(6.6), Inches(11), Inches(0.4),
       "数据截止：2026 年 5 月  |  机密文档  |  仅供 PE 内部决策使用",
       size=11, color=T.ICE, italic=True)
    return s

def slide_section(prs, label, title, sub, page, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, T.NAVY)
    rect(s, Inches(1), Inches(3.2), Inches(0.6), Inches(0.08), T.GOLD)
    tx(s, Inches(1), Inches(2.5), Inches(11), Inches(0.5),
       label, size=14, bold=True, color=T.GOLD)
    tx(s, Inches(1), Inches(3.5), Inches(11.5), Inches(1.5),
       title, size=44, bold=True, color=T.INVERT)
    tx(s, Inches(1), Inches(5.0), Inches(11.5), Inches(1.0),
       sub, size=18, color=T.ICE)
    chrome(s, page, total, on_dark=True)
    return s

# ============================================================
# Main build
# ============================================================
def build(prs):
    TOTAL = 35

    # ===== 1. COVER =====
    slide_cover(prs)

    # ===== 2. TOC + KPI Dashboard 合并 =====
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, T.BG)
    chrome(s, 2, TOTAL, "EXECUTIVE OVERVIEW")
    page_title(s, "目录与 8 个关键数字", "Table of Contents & KPI Dashboard")

    # 左侧 TOC
    tx(s, Inches(0.5), Inches(1.5), Inches(5.5), Inches(0.3), "目录", size=13, bold=True, color=T.NAVY)
    toc = [
        ("0", "执行摘要", "当前格局 / 推演 / 5 行动"),
        ("1", "行业本质", "4 拐点 + 顶尖人物矩阵"),
        ("2", "技术路线之争", "VLA vs 世界模型 / 数据 / Scaling Law"),
        ("3", "商业化悬崖 (DD)", "Autonomy / 订单 Tier / 替代方案"),
        ("4", "6 种玩家路径", "Tesla / Figure / 1X / PI / 中硬件 / 中场景"),
        ("5", "产业链与硬件", "BOM / Top10 / 市场规模"),
        ("6", "中美格局与政策", "16 环节 / 三情景 / 估值"),
        ("A", "PE 内部工具", "IC Memo / Q&A / 决策树"),
        ("B", "数据来源 (206 refs)", "Tier 1/2/3 + 置信度"),
    ]
    for i, (num, name, sub) in enumerate(toc):
        y = 1.85 + i * 0.5
        # Number badge
        ci = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(y), Inches(0.32), Inches(0.32))
        ci.fill.solid(); ci.fill.fore_color.rgb = T.NAVY; ci.line.fill.background(); ci.shadow.inherit = False
        tx(s, Inches(0.5), Inches(y), Inches(0.32), Inches(0.32), num, size=11, bold=True,
           color=T.INVERT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        tx(s, Inches(0.95), Inches(y - 0.02), Inches(5.2), Inches(0.3), name, size=12, bold=True, color=T.NAVY)
        tx(s, Inches(0.95), Inches(y + 0.2), Inches(5.2), Inches(0.25), sub, size=9, color=T.MUTED, italic=True)

    # 右侧 KPI Dashboard
    tx(s, Inches(6.5), Inches(1.5), Inches(6.5), Inches(0.3), "8 个关键数字 · KPI Dashboard", size=13, bold=True, color=T.NAVY)
    kpis = [
        ("1.8 万", "2025 全球出货", "YoY +508% IDC", T.GOLD),
        ("80-95%", "中国出货占比", "IDC vs GGII 口径", T.DANGER),
        ("2.8×", "中美 BOM 剪刀差", "$46K vs $130K MS", T.WARN),
        ("$39.5B", "Figure 估值", "美系 Top1", T.INFO),
        ("$75 亿", "2030E 年度 TAM", "中位; 区间 $30-144亿", T.PURPLE),
        ("69-113%", "出货 CAGR 25-30", "1.8万→25-80万", T.SUCCESS),
        ("¥373 亿", "中国具身融资", "YTD 日均 ¥2.5亿", T.DANGER),
        ("¥550+ 亿", "中国地方基金", "累计承诺", T.GOLD),
    ]
    cw = 1.55; ch = 1.3; gap = 0.1
    for i, (val, lab, sub, ac) in enumerate(kpis):
        r, c = i // 4, i % 4
        x = 6.5 + c * (cw + gap); y = 1.85 + r * (ch + gap)
        kpi_chip(s, Inches(x), Inches(y), Inches(cw), Inches(ch), val, lab, sub, ac)

    # 底部 takeaway
    rrect(s, Inches(0.5), Inches(5.0), Inches(12.3), Inches(1.9), T.SUBTLE, radius=0.02)
    tx(s, Inches(0.7), Inches(5.1), Inches(12), Inches(0.4),
       "核心一句话", size=10, bold=True, color=T.GOLD)
    tx(s, Inches(0.7), Inches(5.4), Inches(12), Inches(1.5),
       '{bold}"美强软件、中强硬件" 的非对称稳态已成型。{/bold} 中国主导出货（80-95%）+ 硬件 BOM 优势（2.8×），美国主导 VLA 模型 + 估值（Figure $39.5B、Skild $14B 在谈）；资本市场把"软件大脑"和"硬件本体"按 5-10× 估值分层。2027-2028 是关键验证窗口（Tesla Optimus 量产 + 中系 IPO + 中国 550 亿补贴见效 + VLA Scaling Law）。',
       size=11, color=T.TEXT)

    # ===== 3. 当前格局 + 三情景（合并 = 1 页）=====
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, T.BG)
    chrome(s, 3, TOTAL, "0.1 当前格局核心结论 + 0.2 三情景推演")
    page_title(s, "当前格局 · 中美非对称稳态 + 三种未来情景推演")

    # 左：中美对比（占 50%）
    section_label(s, Inches(0.5), Inches(1.5), "0.1 · 当前格局核心结论", T.NAVY)

    # 中国
    rrect(s, Inches(0.5), Inches(1.9), Inches(3.0), Inches(5.0), T.CARD, radius=0.02)
    rect(s, Inches(0.5), Inches(1.9), Inches(3.0), Inches(0.4), T.DANGER)
    tx(s, Inches(0.65), Inches(1.95), Inches(2.85), Inches(0.3),
       "🇨🇳 中国 · 主导硬件 + 出货", size=11, bold=True, color=T.INVERT)
    multi_tx(s, Inches(0.65), Inches(2.45), Inches(2.8), Inches(4.4), [
        "全球出货 80-95%（IDC vs GGII）",
        "BOM $32-46K vs 美 $130K",
        "出货榜前六悉数中国厂商",
        "智元/宇树/优必选 ~70% 中国",
        "2025 YTD 融资 ¥373 亿",
        "纯人形（剔机器狗）70-85%",
    ], size=10)

    # 美国
    rrect(s, Inches(3.7), Inches(1.9), Inches(3.0), Inches(5.0), T.CARD, radius=0.02)
    rect(s, Inches(3.7), Inches(1.9), Inches(3.0), Inches(0.4), T.INFO)
    tx(s, Inches(3.85), Inches(1.95), Inches(2.85), Inches(0.3),
       "🇺🇸 美国 · 主导模型 + 估值", size=11, bold=True, color=T.INVERT)
    multi_tx(s, Inches(3.85), Inches(2.45), Inches(2.8), Inches(4.4), [
        "Figure $39.5B / Skild $14B 在谈",
        "Physical Intelligence $11B 在谈",
        "Apptronik $5B / 1X $10B+ 在谈",
        "5 家估值合计 ~$79.5B",
        "2025 营收估 ~$105M",
        "{bold}加权 PSR ~757×{/bold}",
        "{italic}软硬件估值分层 5-10×{/italic}",
    ], size=10)

    # 右：三情景饼图 + 文字
    section_label(s, Inches(7.0), Inches(1.5), "0.2 · 三种未来情景（5 年视角）", T.GOLD)

    add_pie(s, Inches(7.0), Inches(1.85), Inches(2.7), Inches(2.7),
            ["Bifurcation 基准", "China-led", "US-led"],
            [42.5, 35, 22.5],
            colors=[T.GOLD, T.DANGER, T.INFO],
            title="概率分布", show_pct=True, show_legend=True)

    # 右侧情景文字描述
    multi_tx(s, Inches(9.9), Inches(1.85), Inches(3.4), Inches(5.1), [
        "§ {bold}China-led（35%）{/bold}",
        "中国硬件成本+场景密度赢得全球",
        "2030 累计：中 100万/美 12万",
        "",
        "§ {bold}US-led（15-25%）{/bold}",
        "humanoid IRA + Tesla $20K 兑现",
        "2030 累计：中 30万/美 80万",
        "",
        "§ {bold}Bifurcation 基准（40-50%）{/bold}",
        "{gold}双轨平行，BOM 重合 < 30%{/gold}",
        "2030 累计：中 60万/美 25万 = 90万",
        "ASP $12K (中) / $40K (美)",
        "",
        "{bold}PE 行动{/bold}：双线对冲（中硬件 + 美软件）",
    ], size=9, bullet="")

    # 三情景柱图 (底部)
    add_bar_v(s, Inches(7.0), Inches(4.7), Inches(6.2), Inches(2.3),
              ["China-led", "US-led", "Bifurcation"],
              [("中国", [100, 30, 60]), ("美国", [12, 80, 25]), ("其他", [8, 10, 5])],
              colors=[T.DANGER, T.INFO, T.LIGHT_MUTED],
              title="2030 累计保有量（万台 installed base）",
              stacked=True)

    # ===== 4. 12 判断 + 5 行动（一页）=====
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, T.BG)
    chrome(s, 4, TOTAL, "0.2 关键判断 + 0.3 PE 行动")
    page_title(s, "12 条关键判断 + 5 个最值得做的决定",
               "对内交付材料 · 每条 ≤ 2 行")

    # 左：12 判断
    section_label(s, Inches(0.5), Inches(1.5), "0.2 · 12 条关键判断", T.NAVY)
    judg = [
        "① {bold}中硬件 + 美软件稳态{/bold} — 未来 5 年难反转",
        "② 双轨平行是基准情景（40-50%）",
        "③ {bold}2027-2028 关键验证窗口{/bold} — Tesla 量产 + 中系 IPO + 补贴见效",
        "④ 美系一线估值泡沫 — 2027-2028 回调 50-70%",
        "⑤ {bold}中系核心零部件是 Alpha 主战场{/bold}（绿的/双环/步科）",
        "⑥ humanoid IRA 是 watchlist 项（区间 20-50%）",
        "⑦ 仿真/数据是双轨稀缺资产（NVIDIA Isaac 接近独占）",
        "⑧ {bold}中国 550 亿补贴占 GDP ~0.08%{/bold}（vs EV 高峰 ~0.51%，仅 15-20%）",
        "⑨ 出口管制双向反制：美→中 算力；中→美 稀土/镁/PEEK",
        "⑩ {bold}年度 TAM $30-144 亿{/bold}（中位 $75 亿）/ CAGR 69-113%",
        "⑪ 稀土永磁低估卡点：占 BOM 0.5% 但中国 92% 全球",
        "⑫ 宁德×智元×千寻\"中州基地\" = 电池厂+机器人协同范式",
    ]
    multi_tx(s, Inches(0.5), Inches(1.85), Inches(6.4), Inches(5.1), judg, size=9, spacing=2, bullet="")

    # 右：5 行动
    section_label(s, Inches(7.0), Inches(1.5), "0.3 · 5 个最值得做的决定", T.GOLD)

    actions = [
        ("① 2026 H1", "中系卡点 + 原材料", T.SUCCESS,
         "未上市丝杠/灵巧手/六维力/IMU\n上游金力永磁/中研股份/中复神鹰\n{bold}定价锚：2030 产能 × PE 15-25×{/bold}"),
        ("② 2026 IPO 窗口", "中系一线整机跟投", T.SUCCESS,
         "宇树(2026.3.20 受理)/智元/银河通用/星动\n{bold}上限 ¥150-250 亿{/bold}"),
        ("③ 2027 后", "美系一线 · 限观察", T.WARN,
         "Figure/Skild/π0 系统性泡沫\n{bold}当前权重 ≤ 5%，2027 H2 加至 10-15%{/bold}"),
        ("④ 2026 起", "仿真/数据/铍铜种子", T.INFO,
         "国内对标 NVIDIA Isaac / AgiBot World\n铍铜国产替代候选\n{bold}早期 + AI 估值溢价{/bold}"),
        ("⑤ 2026 H2", "7 关键监测点", T.DANGER,
         "Optimus V3 / Tesla 100万 / 中国补贴 / IPO 估值\n humanoid IRA / VLA Scaling / 稀土反制"),
    ]
    ay = 1.85
    for label, ttl, c, body in actions:
        rrect(s, Inches(7.0), Inches(ay), Inches(6.3), Inches(1.0), T.CARD, radius=0.03)
        rect(s, Inches(7.0), Inches(ay), Inches(0.08), Inches(1.0), c)
        tx(s, Inches(7.2), Inches(ay + 0.08), Inches(1.4), Inches(0.25),
           label, size=9, bold=True, color=c)
        tx(s, Inches(8.6), Inches(ay + 0.08), Inches(4.6), Inches(0.3),
           ttl, size=11, bold=True, color=T.NAVY)
        tb = s.shapes.add_textbox(Inches(7.2), Inches(ay + 0.4), Inches(6.0), Inches(0.55))
        tf = tb.text_frame; tf.word_wrap = True
        tf.margin_left = Emu(0); tf.margin_right = Emu(0); tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
        for li, line in enumerate(body.split('\n')):
            p = tf.add_paragraph() if li else tf.paragraphs[0]
            p.space_after = Pt(0)
            render(p, line, 8, T.TEXT)
        ay += 1.05

    # =========== CHAPTER 1: 行业本质 ===========
    slide_section(prs, "CHAPTER 1", "行业本质",
                  "4 拐点 · 顶尖人物 10×5 矩阵 · 5 分歧线", 5, TOTAL)

    # ===== 6. 4 拐点 + 顶尖人物矩阵（一页）=====
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, T.BG)
    chrome(s, 6, TOTAL, "1.1 拐点驱动 + 1.2 人物矩阵")
    page_title(s, "为什么是现在 · 4 拐点驱动 + 顶尖人物 10×5 观点矩阵")

    # 左：4 拐点
    section_label(s, Inches(0.5), Inches(1.5), "1.1 · 4 个拐点驱动", T.NAVY)
    inflections = [
        ("①", "技术拐点", T.INFO, "VLA 范式成熟 · 2023 RT-2 → 2025 Helix/π0/GR00T\nSOTA OpenVLA-OFT @ LIBERO 97.1%"),
        ("②", "算力拐点", T.SUCCESS, "训 7B VLA ≈ 23K GPU-h < $100K\nJetson Thor 端侧 2,070 FP4 TFLOPS"),
        ("③", "数据拐点", T.WARN, "Open X 100 万+ / AgiBot 100 万+ / 217 任务\n{italic}但 Jang: 数据稀缺仍是 Physical AGI 瓶颈{/italic}"),
        ("④", "经济拐点", T.DANGER, "美制造业 2030 缺口 210 万 / 中国劳动人口年减 800 万\n{italic}但 ROI：美 6.5 年 / 中 12 年{/italic}"),
    ]
    iy = 1.85
    for label, ttl, c, body in inflections:
        rrect(s, Inches(0.5), Inches(iy), Inches(4.0), Inches(1.2), T.CARD, radius=0.03)
        rect(s, Inches(0.5), Inches(iy), Inches(0.07), Inches(1.2), c)
        tx(s, Inches(0.65), Inches(iy + 0.08), Inches(0.3), Inches(0.4),
           label, size=22, bold=True, color=c)
        tx(s, Inches(1.05), Inches(iy + 0.12), Inches(3), Inches(0.35),
           ttl, size=13, bold=True, color=T.NAVY)
        tb = s.shapes.add_textbox(Inches(1.05), Inches(iy + 0.5), Inches(3.4), Inches(0.7))
        tf = tb.text_frame; tf.word_wrap = True
        tf.margin_left = Emu(0); tf.margin_right = Emu(0); tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
        for li, line in enumerate(body.split('\n')):
            p = tf.add_paragraph() if li else tf.paragraphs[0]
            p.space_after = Pt(0)
            render(p, line, 8, T.TEXT)
        iy += 1.27

    # 右：顶尖人物矩阵
    section_label(s, Inches(4.7), Inches(1.5), "1.2 · 顶尖人物 10×5 观点矩阵", T.GOLD)
    matrix_rows = [
        ["a16z (Casado/Hsu/Polovets)", "中立", "同意", "中立", "中立", "{red}不{/red}"],
        ["Sequoia", "{green}强{/green}", "同意", "同意", "中立", "{red}不{/red}"],
        ["Brett Adcock (Figure)", "{green}强{/green}", "{green}强{/green}", "{green}强{/green}", "{green}强{/green}", "{red}不{/red}"],
        ["Eric Jang (1X)", "同意", "同意", "同意", "同意", "{red}不{/red}"],
        ["Hausman/Levine (PI)", "{green}强{/green}", "同意", "同意", "中立", "{red}不{/red}"],
        ["王兴兴 (宇树)", "同意", "同意", "同意", "同意", "{red}不{/red}"],
        ["王鹤 (银河通用)", "中立", "同意", "{red}不{/red}", "中立", "{red}不{/red}"],
        ["Vanhoucke (Waymo)", "同意", "同意", "中立", "中立", "{red}不{/red}"],
        ["Yann LeCun (AMI Labs)", "{red}强反{/red}", "{red}反{/red}", "{red}反{/red}", "中立", "{red}强反{/red}"],
        ["Marc Raibert (BD)", "同意", "同意", "{red}反{/red}", "{red}反{/red}", "中立"],
    ]
    add_table(s, Inches(4.7), Inches(1.85), Inches(8.6), Inches(4.0),
              ["人物 / 机构", "拐点已到", "工业 PMF 2026", "C 端 2030", "必须双足", "Tesla 会赢"],
              matrix_rows, body_font_size=7, header_font_size=8)

    # 3 个隐藏共识
    rrect(s, Inches(4.7), Inches(6.0), Inches(8.6), Inches(1.0), T.SUBTLE, radius=0.02)
    tx(s, Inches(4.85), Inches(6.05), Inches(8.4), Inches(0.3),
       "3 个隐藏共识 · 矩阵 takeaway", size=9, bold=True, color=T.GOLD)
    tx(s, Inches(4.85), Inches(6.3), Inches(8.4), Inches(0.7),
       '① {bold}"Tesla 会赢"样本内 0 支持票{/bold}（注：不含 Musk/Jonas/Wood 等 bulls；带后分歧 5:5）  ② {bold}时间表激进+全乐观仅 Adcock 一人{/bold}  ③ {bold}必须双足唯一公开反对：Marc Raibert{/bold}',
       size=8, color=T.TEXT)

    # 注释
    tx(s, Inches(4.7), Inches(7.0), Inches(8.6), Inches(0.2),
       "注：非随机抽样；样本不含 Tesla bulls；定义：拐点已到=GPT 时刻；强/不=强同意/不同意",
       size=7, color=T.MUTED, italic=True)

    # ===== 7. 5 分歧线 + 10 引语 + Tesla 没人押注（一页）=====
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, T.BG)
    chrome(s, 7, TOTAL, "1.3-1.4 分歧线 + 引语")
    page_title(s, "5 条核心分歧线 + 最具穿透力的引语")

    section_label(s, Inches(0.5), Inches(1.5), "1.3 · 5 条核心分歧", T.NAVY)
    splits = [
        ("①", "GPT 时刻", T.SUCCESS,
         "{bold}乐观{/bold}：Adcock/王兴兴/Levine\n{bold}悲观{/bold}：LeCun/Karpathy/Pinto\n{italic}a16z 内 Casado: 估值疯狂{/italic}"),
        ("②", "VLA vs 世界模型", T.INFO,
         "{bold}VLA{/bold}：PI π0/Figure Helix/Tesla\n{bold}WM{/bold}：1X/NVIDIA GEAR/LeCun AMI/王兴兴"),
        ("③", "数据来源", T.WARN,
         "{bold}真机{/bold}（美）：PI/Figure\n{bold}仿真{/bold}（中）：银河/NVIDIA\n{bold}视频{/bold}（学术）：Pinto/Jang"),
        ("④", "形态", T.PURPLE,
         "{bold}必须双足{/bold}：Adcock/Bornich/王兴兴\n{bold}形态多元{/bold}：Raibert/Skild\n{bold}非人形{/bold}：Symbotic/Locus 已赢"),
        ("⑤", "谁会赢", T.DANGER,
         "{bold}通用 brain{/bold}：PI/Skild\n{bold}垂直整机{/bold}：Figure/Apptronik\n{bold}中国低成本{/bold}：宇树/智元"),
    ]
    sy = 1.85
    for label, ttl, c, body in splits:
        rrect(s, Inches(0.5), Inches(sy), Inches(6.0), Inches(1.0), T.CARD, radius=0.03)
        rect(s, Inches(0.5), Inches(sy), Inches(0.07), Inches(1.0), c)
        tx(s, Inches(0.65), Inches(sy + 0.05), Inches(0.4), Inches(0.4),
           label, size=18, bold=True, color=c)
        tx(s, Inches(1.05), Inches(sy + 0.1), Inches(2), Inches(0.3),
           ttl, size=11, bold=True, color=T.NAVY)
        tb = s.shapes.add_textbox(Inches(3.0), Inches(sy + 0.1), Inches(3.4), Inches(0.85))
        tf = tb.text_frame; tf.word_wrap = True
        tf.margin_left = Emu(0); tf.margin_right = Emu(0); tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
        for li, line in enumerate(body.split('\n')):
            p = tf.add_paragraph() if li else tf.paragraphs[0]
            p.space_after = Pt(0)
            render(p, line, 8, T.TEXT)
        sy += 1.06

    # 右：5 引语
    section_label(s, Inches(6.8), Inches(1.5), "1.4 · 最具穿透力引语 (Top 8)", T.GOLD)
    quotes = [
        ('"Within 10 years, every home will have a humanoid."', "Brett Adcock (Figure CEO), Time 10/2025"),
        ('"More like the Apollo program than a science experiment."', "Sergey Levine (PI), Dwarkesh 9/2025"),
        ('"LLMs will become useless within 5 years."', "Yann LeCun (AMI Labs), 10/27/2025"),
        ('"Humanoids are most hyped — valuations crazy before any revenue."', "Leo Polovets (Humba VC), a16z 播客"),
        ('"数据关注度有点太高，最大问题在模型；VLA 是傻瓜式架构。"', "王兴兴 (宇树), 21 经济 8/9/2025"),
        ('"工厂先行，家用还要等；让机器像人的不是腿和手，是智能。"', "Marc Raibert (BD), Digitimes 4/2025"),
        ('"5-10 年才能什么活儿都干；特别不建议讲具身 AGI。"', "王鹤 (银河通用), 36Kr"),
        ('"Real-world deployments still largely confined to demos and pilots."', "Oliver Hsu, a16z, 1/13/2026"),
    ]
    qy = 1.85
    for q, attr in quotes:
        rrect(s, Inches(6.8), Inches(qy), Inches(6.5), Inches(0.62), T.SUBTLE, radius=0.02)
        rect(s, Inches(6.8), Inches(qy), Inches(0.05), Inches(0.62), T.GOLD)
        tx(s, Inches(6.95), Inches(qy + 0.05), Inches(6.3), Inches(0.3),
           q, size=8.5, italic=True, color=T.TEXT, bold=True)
        tx(s, Inches(6.95), Inches(qy + 0.36), Inches(6.3), Inches(0.22),
           "— " + attr, size=7, color=T.MUTED)
        qy += 0.66

    # =========== CHAPTER 2: 技术路线 ===========
    slide_section(prs, "CHAPTER 2", "技术路线之争",
                  "VLA vs 世界模型 · 数据三路 · 形态 · Scaling Law", 8, TOTAL)

    # ===== 9. 技术路线全景（一页）=====
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, T.BG)
    chrome(s, 9, TOTAL, "2.1-2.5 技术路线")
    page_title(s, "技术路线之争 · 真正的胜负所在", "VLA vs 世界模型 · 数据三路 · 形态之争 · Scaling Law")

    # 上半左：VLA vs 世界模型
    rrect(s, Inches(0.5), Inches(1.5), Inches(6.2), Inches(2.4), T.CARD, radius=0.02)
    rect(s, Inches(0.5), Inches(1.5), Inches(6.2), Inches(0.3), T.INFO)
    tx(s, Inches(0.65), Inches(1.55), Inches(6), Inches(0.25),
       "路线 ① · 端到端 VLA vs 显式世界模型", size=10, bold=True, color=T.INVERT)
    tx(s, Inches(0.65), Inches(1.85), Inches(2.85), Inches(0.25),
       "{bold}VLA 派 · 主流{/bold}", size=9, bold=True, color=T.NAVY)
    multi_tx(s, Inches(0.65), Inches(2.1), Inches(2.85), Inches(1.7), [
        "PI (π0/π0.5)：Flow-matching",
        "Figure (Helix 02)：S1+S2 双系统",
        "Tesla 垂直整合",
        "Apptronik × Google DM",
        "{bold}PE 投注：$60B+{/bold}",
    ], size=8)
    tx(s, Inches(3.65), Inches(1.85), Inches(2.85), Inches(0.25),
       "{bold}世界模型派{/bold}", size=9, bold=True, color=T.PURPLE)
    multi_tx(s, Inches(3.65), Inches(2.1), Inches(2.85), Inches(1.7), [
        "1X World Model（Eric Jang）",
        "NVIDIA GEAR (GR00T N1)",
        "Yann LeCun → AMI Labs ($1B)",
        "王兴兴 UnifoLM-WMA-0",
        "{italic}学术声誉强 / 融资小{/italic}",
    ], size=8)

    # 上半右：必须双足？
    rrect(s, Inches(6.9), Inches(1.5), Inches(6.4), Inches(2.4), T.CARD, radius=0.02)
    rect(s, Inches(6.9), Inches(1.5), Inches(6.4), Inches(0.3), T.PURPLE)
    tx(s, Inches(7.05), Inches(1.55), Inches(6.2), Inches(0.25),
       "路线 ② · 必须双足吗 · 形态之争", size=10, bold=True, color=T.INVERT)
    forms = [
        ("必须", T.DANGER, "Adcock/Bornich/王兴兴", "$30-130K", "Figure F.03 / NEO / 宇树 G1"),
        ("多元", T.WARN, "Raibert/Skild omni", "$15-50K", "轮式+双臂+灵巧手"),
        ("非人形已赢", T.SUCCESS, "Symbotic/Locus/AutoStore", "$5-50K", "17K AMR / Walmart 40% 降本"),
    ]
    fy = 1.85
    for ttl, c, who, bom, why in forms:
        tx(s, Inches(7.05), Inches(fy), Inches(1.4), Inches(0.22),
           ttl, size=9, bold=True, color=c)
        tx(s, Inches(8.5), Inches(fy), Inches(2.4), Inches(0.22),
           who, size=8, color=T.TEXT)
        tx(s, Inches(10.9), Inches(fy), Inches(1.0), Inches(0.22),
           bom, size=8, bold=True, color=T.NAVY)
        tx(s, Inches(11.9), Inches(fy), Inches(1.5), Inches(0.22),
           why, size=7, color=T.MUTED, italic=True)
        fy += 0.45
    # divider line
    rect(s, Inches(7.05), Inches(3.35), Inches(6.2), Inches(0.02), T.DIVIDER)
    tx(s, Inches(7.05), Inches(3.45), Inches(6.2), Inches(0.4),
       "Marc Raibert (BD): \"让机器人像人不是腿和手，是智能\" · 工业场景轮式 + 双臂 + 灵巧手可能 95% 替代双足，成本低 30-50%",
       size=7.5, italic=True, color=T.TEXT)

    # 下半左：数据三路
    rrect(s, Inches(0.5), Inches(4.05), Inches(6.2), Inches(2.95), T.CARD, radius=0.02)
    rect(s, Inches(0.5), Inches(4.05), Inches(6.2), Inches(0.3), T.WARN)
    tx(s, Inches(0.65), Inches(4.1), Inches(6), Inches(0.25),
       "路线 ③ · 数据来源 · 真机 vs 仿真 vs 视频", size=10, bold=True, color=T.INVERT)
    add_table(s, Inches(0.55), Inches(4.42), Inches(6.1), Inches(2.5),
              ["路径", "代表", "规模", "风险"],
              [["真机 (美)", "PI/Figure/Hausman", "Open X 100万+", "贵+慢"],
               ["仿真 (中)", "银河 GraspVLA/NVIDIA", "10 亿帧/430,000× 实时", "Sim-to-Real Gap"],
               ["视频 (学术)", "Jang/Pinto/V-JEPA2", "Ego4D 3,670h", "视觉→action 桥梁"]],
              col_widths=[1.0, 2.0, 1.5, 1.5], header_font_size=8, body_font_size=7.5)

    # 下半右：Scaling Law
    rrect(s, Inches(6.9), Inches(4.05), Inches(6.4), Inches(2.95), T.CARD, radius=0.02)
    rect(s, Inches(6.9), Inches(4.05), Inches(6.4), Inches(0.3), T.DANGER)
    tx(s, Inches(7.05), Inches(4.1), Inches(6.2), Inches(0.25),
       "路线 ④ · Scaling Law · GPT 时刻到了吗？", size=10, bold=True, color=T.INVERT)
    tx(s, Inches(7.05), Inches(4.4), Inches(3), Inches(0.25),
       "{bold}当前 SOTA 实测{/bold}", size=9, bold=True, color=T.NAVY)
    multi_tx(s, Inches(7.05), Inches(4.7), Inches(3), Inches(2.3), [
        "LIBERO 4: OpenVLA-OFT {bold}97.1%{/bold}",
        "CALVIN: π0 75-85%",
        "SimplerEnv: GR00T 60-70%",
        "{red}真实环境 30-50%{/red}",
        "1X NEO 官方 60-70%",
        "Tesla: \"not in usage\"",
    ], size=8)
    tx(s, Inches(10.15), Inches(4.4), Inches(3.1), Inches(0.25),
       "{bold}GPT 时刻预测{/bold}", size=9, bold=True, color=T.GOLD)
    multi_tx(s, Inches(10.15), Inches(4.7), Inches(3.1), Inches(2.3), [
        "Tedrake：threshold 后非线性",
        "1700h 真机 + 47K 仿真 = 当前最大 LBM",
        "",
        "{bold}预计 2027-2028{/bold}",
        "（数据 1 亿+ / 模型 30B+）",
        "{italic}Finn: scale subordinate to solving{/italic}",
    ], size=8)

    # =========== CHAPTER 3: 商业化悬崖 ===========
    slide_section(prs, "CHAPTER 3", "商业化悬崖 (DD)",
                  "Autonomy · Tier · 替代方案 · RaaS · 投资结论变更", 10, TOTAL)

    # ===== 11. 真实 Autonomy Level 评分卡 (完整 10 行 一页) =====
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, T.BG)
    chrome(s, 11, TOTAL, "3.1 Autonomy Level")
    page_title(s, "真实 Autonomy Level · 行业部署评分卡",
               "基于 SAE 改编 L0-L4 分级 · 10 家公司完整对比")

    # Top: Hero takeaway
    rrect(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.55), T.NAVY, radius=0.02)
    tx(s, Inches(0.7), Inches(1.55), Inches(12), Inches(0.45),
       '{gold}核心结论{/gold}  ·  {bold}媒体宣传 "~$3B 人形订单" 中 Tier A 接近 0；行业普遍 Tier D (PoC)；人形相对 AMR/协作臂 ROI 差 5-16×{/bold}',
       size=11, color=T.INVERT, anchor=MSO_ANCHOR.MIDDLE)

    # Full autonomy table
    aut_rows = [
        ["Tesla Optimus", "Fremont (R&D) 无客户", "L0-L1", "零规模化部署", '{red}Musk Q4 2025: "not in usage in our factories in a material way"{/red}'],
        ["Figure 02", "BMW Spartanburg", "L3（争议）", "1,250h / 90K 钣金 / 30K X3", '{bold}Adcock: "no teleop in market"{/bold}；Scott Walter 等指演示有遥操作'],
        ["Figure 03", "BMW + Munich 分阶段", "声称 L3", "2026 分阶段铺开", "独立验证 [UNDISCLOSED]；任何未来 teleop 都打击 $39B"],
        ["1X NEO", "消费者 pre-order", "L0-L2 混合", "{red}CEO 公开 60-70% 自主率{/red}", '"Expert Mode" 远程操作员；CEO: "running towards a cliff"'],
        ["Apptronik Apollo", "Mercedes / GXO / Jabil", "L2 (lab-trained)", "仍处 pilot", "商业 scale 目标 2026 H2"],
        ["Agility Digit", "GXO / Mercado Libre", "{green}L3 in defined ODD{/green}", "100K+ 货箱 (2025.11) / 2:1 电池", "迄今{bold}唯一公开商业 RaaS{/bold}"],
        ["智元远征 A2", "中国移动数据采集 PoC", "L0-L1", "data-collection service", "实质 R&D 采集合同，不是生产"],
        ["宇树 H1/G1", "demo / 教育 / 科研", "L0-L1 工厂", "无 24×7 生产部署", "G1 占公司 60% 营收（机器狗+教育为主）"],
        ["BD Atlas (新)", "Hyundai / Google DeepMind", "L2-L3", "CES 2026 商业揭幕", "2026 整年订单已锁定"],
        ["UBTech Walker S2", "中国政府 / 工业站点", "L1-L2", "11 月交付节奏", "政府采购投标基本绑定，但端用多为非生产"],
    ]
    add_table(s, Inches(0.5), Inches(2.2), Inches(12.3), Inches(4.5),
              ["公司 / 型号", "部署点", "Autonomy", "运行数据", "关键披露"],
              aut_rows, col_widths=[1.6, 2.0, 1.4, 2.5, 4.8],
              header_font_size=9, body_font_size=7.5)

    # Bottom note
    rrect(s, Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.3), T.SUBTLE, radius=0.02)
    tx(s, Inches(0.7), Inches(6.82), Inches(12), Inches(0.25),
       "⚠ MTBF / MTBI 数字全行业缺失 · \"任何省略 MTBF 的 PE pitch deck 应视为尽调不完整\"",
       size=9, bold=True, color=T.WARN, italic=True)

    # ===== 12. 订单 Tier 分级 + Top 10 订单（一页）=====
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, T.BG)
    chrome(s, 12, TOTAL, "3.2 订单质量分级")
    page_title(s, "订单质量 · Tier A-E 分级框架 + Top 10 订单重新评级",
               "媒体引述 ~$3B 中：真 firm Tier A+B ≈ $1-1.5B；剩余 $1.5-2B 是 C/D/E")

    # 左：Tier 框架（小竖列）
    section_label(s, Inches(0.5), Inches(1.5), "Tier A-E 框架", T.NAVY)
    tier_rows = [
        ("A", "已收款不可撤销", T.SUCCESS, "< 5%"),
        ("B", "Firm 采购合同", T.SUCCESS, "5-15%"),
        ("C", "Master 框架/MSA", T.WARN, "30-60%"),
        ("D", "付费 PoC", T.DANGER, "60-80%"),
        ("E", "LOI / MoU", T.DANGER, "75-95%"),
    ]
    ty = 1.85
    for t, desc, c, risk in tier_rows:
        rrect(s, Inches(0.5), Inches(ty), Inches(2.6), Inches(0.85), T.CARD, radius=0.03)
        rect(s, Inches(0.5), Inches(ty), Inches(0.4), Inches(0.85), c)
        tx(s, Inches(0.5), Inches(ty), Inches(0.4), Inches(0.85),
           t, size=22, bold=True, color=T.INVERT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        tx(s, Inches(1.0), Inches(ty + 0.08), Inches(2.0), Inches(0.3),
           desc, size=10, bold=True, color=T.NAVY)
        tx(s, Inches(1.0), Inches(ty + 0.4), Inches(2.0), Inches(0.3),
           f"取消风险 {risk}", size=8, color=T.MUTED, italic=True)
        ty += 0.95

    # 右：Top 10 订单完整列表
    section_label(s, Inches(3.3), Inches(1.5), "Top 10 公开订单 · 完整重新分级", T.GOLD)
    order_rows = [
        ["1", "UBTech Walker S 2025 全年", "¥11-13 亿", "{green}B/C{/green}", "政府采购授标即绑定；交付进行中"],
        ["2", "UBTech-广西防城港", "¥2.64 亿", "{green}B{/green}", "Walker S2；2025.12 交付"],
        ["3", "UBTech-自贡", "¥1.59 亿", "{green}B{/green}", "Walker S2；数据采集中心"],
        ["4", "宇树 2025 工业订单", "¥12 亿", "{orange}B/C{/orange}", "G1 占 60% 营收；教育+小批量"],
        ["5", "智元+宇树-中移动", "¥1.24 亿", "{red}D (PoC){/red}", "数据采集 R&D 服务合同，非生产"],
        ["6", "星动纪元 2025 累计", "¥5 亿", "{orange}B/C/D{/orange}", "50% 海外大概率 D-tier"],
        ["7", "Apptronik × Merc/GXO/Jabil", "[UNDISCLOSED]", "{red}D{/red}", "商业 scale 目标 2026 H2"],
        ["8", "Agility × GXO / Mercado", "100K+ 货箱", "{green}B (GXO){/green}", "{bold}迄今最经实测的人形商业 datapoint{/bold}"],
        ["9", "Galbot × CATL/Bosch", "\"1,000+ 计划\"", "{red}C/E{/red}", "未公开单价×数量×交付表"],
        ["10", "Figure × BMW", "未披露 / 11 月 PoC", "{orange}B+C{/orange}", "Fortune 报道 BMW 反驳 Adcock"],
        ["—", "Tesla 自厂 1,000 台", "n/a", "{red}剔除{/red}", "Musk: \"not in usage in a material way\""],
    ]
    add_table(s, Inches(3.3), Inches(1.85), Inches(9.7), Inches(4.5),
              ["#", "交易", "金额", "Tier", "caveat"],
              order_rows, col_widths=[0.3, 2.5, 1.5, 1.2, 4.2],
              header_font_size=9, body_font_size=7.5)

    # 底部图：订单分布饼图
    add_pie(s, Inches(0.5), Inches(6.4), Inches(2.8), Inches(0.9),
            ["Tier A+B (firm)", "Tier C+D+E (软)"],
            [40, 60],
            colors=[T.SUCCESS, T.DANGER], show_pct=True, show_legend=False, title="")

    rrect(s, Inches(3.3), Inches(6.4), Inches(9.7), Inches(0.9), T.SUBTLE, radius=0.02)
    tx(s, Inches(3.45), Inches(6.45), Inches(9.5), Inches(0.3),
       "Bottom line · PE 视角订单质量真相",
       size=9, bold=True, color=T.GOLD)
    tx(s, Inches(3.45), Inches(6.7), Inches(9.5), Inches(0.5),
       "Tier A（已收款已交付）{red}≈ 0{/red}；中国端 Tier B 多为政府采购（数据采集/安防/教育，{red}非 ROI 驱动{/red}）；美国端 Tier D（pilot）+ 少量 Tier B（Digit @ GXO RaaS）。{bold}媒体引述 ~$3B 实际 firm 仅 $1-1.5B；剩余 $1.5-2B 是 C/D/E{/bold}",
       size=8, color=T.TEXT)

    # ===== 13. 替代方案 ROI（一页含柱图 + 表格 + 引语）=====
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, T.BG)
    chrome(s, 13, TOTAL, "3.3 替代方案 ROI")
    page_title(s, "替代方案 ROI · 仓储战争已被非人形赢得",
               "人形相对成熟方案 ROI 差 5-16× · McKinsey: 人形 payback 仍是 Locus 3-5×")

    # 左：回本周期柱图
    add_bar_h(s, Inches(0.5), Inches(1.5), Inches(6.0), Inches(3.2),
              ["Locus AMR", "UR 协作臂(典型)", "Symbotic 报", "AutoStore", "人形 Agility naive"],
              [0.5, 0.83, 1.2, 1.5, 8.0],
              colors=[T.SUCCESS, T.SUCCESS, T.SUCCESS, T.SUCCESS, T.DANGER],
              title="回本周期对比（年）", show_values=True)

    # 右：详细对比表
    add_table(s, Inches(6.8), Inches(1.5), Inches(6.5), Inches(3.5),
              ["方案", "代表厂商", "单台价格", "部署规模", "关键 KPI"],
              [
                  ["AMR", "Locus Robotics", "RaaS", "{bold}17K+ AMR / 70 亿+ 拣选{/bold}", "—"],
                  ["Cube ASRS", "AutoStore", "$1-5M", "1,950+ 系统", "{green}99.8% uptime{/green}"],
                  ["G2P 仓库", "Symbotic", "$10M+", "Walmart 42 DC / 400 系统", "FY25 营收 $2.25B"],
                  ["协作臂 6 轴", "UR", "$11-60K", "{bold}> 100K cobot{/bold}", "—"],
                  ["{red}人形{/red}", "Figure/Apptronik/Agility/Tesla", "$25-130K+", "{red}~几千台 (多为 demo){/red}", "MTBF 未披露"],
              ], col_widths=[1.0, 1.8, 1.0, 1.7, 1.0],
              header_font_size=8, body_font_size=7)

    # 底部：专家引语 + Locus 直接喊话
    rrect(s, Inches(0.5), Inches(5.1), Inches(12.3), Inches(0.85), T.SUBTLE, radius=0.02)
    rect(s, Inches(0.5), Inches(5.1), Inches(0.06), Inches(0.85), T.GOLD)
    tx(s, Inches(0.7), Inches(5.18), Inches(12), Inches(0.3),
       "专家观点 · 非人形已赢的论据", size=10, bold=True, color=T.GOLD)
    tx(s, Inches(0.7), Inches(5.5), Inches(12), Inches(0.45),
       '{bold}Marc Raibert{/bold}（Lex Fridman, BD 创始人）：\"A dynamic, wheeled system is far more efficient than a bipedal model.\"  ·  '
       '{bold}Locus 官博{/bold}：\"Humanoids may not be ready for the warehouse; Locus Array is.\"  ·  '
       '{bold}McKinsey 2024{/bold}：人形 payback 5.3 → 2.8 年，仍是 Locus 部署 3-5×',
       size=8, italic=True, color=T.TEXT)

    # 底部 takeaway
    rrect(s, Inches(0.5), Inches(6.05), Inches(12.3), Inches(1.0), T.NAVY, radius=0.02)
    rect(s, Inches(0.5), Inches(6.05), Inches(0.06), Inches(1.0), T.GOLD)
    tx(s, Inches(0.7), Inches(6.15), Inches(12), Inches(0.3),
       "Takeaway · PE 行动重定位", size=10, bold=True, color=T.GOLD)
    tx(s, Inches(0.7), Inches(6.45), Inches(12), Inches(0.55),
       "{bold}降级人形\"工业通用替代\"叙事{/bold}：Agility / Apptronik 在仓储权重 ≤ 3% (vs Symbotic 等 incumbent)；人形定位重新表述为\"少数非结构化任务补位\"——而非替代 AMR。"
       "{bold}估值锚向真正赢家倾斜{/bold}：Symbotic / Locus / AutoStore / UR 才是仓储战争的赢家。",
       size=9, color=T.INVERT)

    # ===== 14. 安全合规 + RaaS 资产负债表（一页）=====
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, T.BG)
    chrome(s, 14, TOTAL, "3.4-3.5 合规 + RaaS")
    page_title(s, "安全合规 · 法律护城河 + RaaS 资产负债表 · Figure $39B 隐含假设")

    # 左：安全合规
    section_label(s, Inches(0.5), Inches(1.5), "3.4 · 安全 / 合规 · 工业部署法律护城河", T.NAVY)
    rrect(s, Inches(0.5), Inches(1.85), Inches(6.3), Inches(5.1), T.CARD, radius=0.02)
    rect(s, Inches(0.5), Inches(1.85), Inches(0.06), Inches(5.1), T.DANGER)
    multi_tx(s, Inches(0.7), Inches(2.0), Inches(6.1), Inches(4.9), [
        "{bold}ISO 10218 / ANSI R15.06{/bold} 主要覆盖工业机器人本体与系统集成；动态稳定腿式 / 双足移动机器人风险仍处于新标准化项目（A3 / IEC TC 184）与解释适用阶段",
        "",
        "A3 已识别 humanoid 属于 \"dynamically stable industrial mobile robot\"；{bold}现有 R15.06 / R15.08 并未完全覆盖这一形态的特定危险{/bold}",
        "",
        "协作机器人（cobots）有 {bold}ISO/TS 15066{/bold} 力 / 压等级标准；人形机器人在工业用途的安全框架仍在制定中",
        "",
        "{bold}责任划分{/bold}：自主水平越高，责任越倾向 OEM + 软件商（与传统\"集成商承担风险\"模型反转）",
        "",
        "{red}家用 / 医疗几乎完全空白{/red}：1X NEO、Optimus 进家庭的 liability framework 全行业空白",
        "",
        "远程操作员合规：1X \"Expert Mode\" 涉及 GDPR / 美国州法律未澄清",
        "",
        "{bold}PE 行动{/bold}：工业部署现实是\"caged-cell\" 仍是路径最短选择 — cobots 和 AMR 在工业场景压制人形的法律基础",
    ], size=9, bullet="")

    # 右：RaaS 资产负债表
    section_label(s, Inches(7.0), Inches(1.5), "3.5 · RaaS 资产负债表 · Figure $39B 隐含", T.GOLD)
    rrect(s, Inches(7.0), Inches(1.85), Inches(6.3), Inches(5.1), T.CARD, radius=0.02)
    rect(s, Inches(7.0), Inches(1.85), Inches(0.06), Inches(5.1), T.WARN)
    tx(s, Inches(7.2), Inches(2.0), Inches(6.0), Inches(0.3),
       "RaaS 模式真实经济", size=10, bold=True, color=T.NAVY)
    multi_tx(s, Inches(7.2), Inches(2.3), Inches(6.0), Inches(1.8), [
        "Agility 公开：$30/hr ($10-12/hr opex)",
        "Figure RaaS 3 种 inferred（未官方披露）：",
        "  (a) $1,000/月 = $12K/年",
        "  {bold}(b) $25/h × 8h × 250d = $50K/年 [本报告 base]{/bold}",
        "  (c) 2030 量产 blended $25K/台",
        "Digit 单台 $250K capex → naive payback ~8 年",
    ], size=9, bullet="")
    tx(s, Inches(7.2), Inches(4.2), Inches(6.0), Inches(0.3),
       "Figure $39B 估值隐含", size=10, bold=True, color=T.DANGER)
    multi_tx(s, Inches(7.2), Inches(4.5), Inches(6.0), Inches(2.4), [
        "本报告 base case ($50K/年 单班 hourly):",
        "  · 2030 出货 5 万 × $50K = $2.5B 年收入",
        "  · $39.5B / $2.5B = {bold}15.8× 2030 收入{/bold}",
        "  · 但 2025 营收估 $15-30M [INFERRED]",
        "  · {bold}2025 PSR 仍 ~790×{/bold}",
        "5 万台部署需 capex ~$5-7B：",
        "  {red}capex-light 资产负债表假设未公开确认{/red}",
    ], size=9, bullet="")

    # ===== 15. 投资结论变更表（一页完整 10 行）=====
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, T.BG)
    chrome(s, 15, TOTAL, "3.6 投资结论变更")
    page_title(s, "投资结论与决策映射表 · 10 处修订全景",
               "DD 修订 → 估值锚下移 → 投资动作调整 · 7/10 指向美系下行风险")

    cv_rows = [
        ["RaaS 5 年现金流 170×", "Figure $25/h×8h×250d ≈ $50K/年；LTV:CAC 仅 2-5×", "{bold}RaaS 没有数量级优势{/bold}", "美系一线整机 ≤ 5%（vs v1 15-25%）"],
        ["2030 出货 150-300 万 / TAM $300-800 亿", "混淆年度新增 vs 累计；中位 ASP $15K", "年度 25-80 万 / 累计 60-150 万 / TAM $30-144 亿", "{bold}估值锚下移 2-4×{/bold}：Figure → $0.87B sanity"],
        ["仓储是人形最先 PMF", "Symbotic $2.25B / Walmart 40% 降本 / Locus < 6 月", "仓储已被非人形赢得", "Agility / Apptronik ≤ 3%；优先 cobot+AMR+ASRS"],
        ["Tesla 部署 1,000 台 = 量产前奏", 'Musk Q4 2025: "not in usage"', "Tesla 未规模化使用", "US-led 15-25% (vs v1 25%)；Bifurcation 40-50%"],
        ["1X NEO = 真自主家用", "CEO 公开 60-70% 自主率；目标 2028 95%+", "1X 当前是远程操作员+AI 混合", "1X 仓位 ≤ 3%（binary 押 2028 自主率突破）"],
        ["中移动 1.24 亿 = 全球最大订单", "合同实质是 R&D 数据采集服务", "降级 Tier D PoC", "智元/宇树估值剔除；Tier A+B 真实 ~$1.5B (vs $3B)"],
        ["补贴超 EV 早期占 GDP", "年化 / 年 GDP ≈ 0.08%（vs EV 高峰 0.51%）", "量级仅 EV 高峰 15-20%", "补贴权重 ≤ 10%；核心论据回归 BOM+场景+数据"],
        ["金力永磁/中研股份独家供应 Tesla", "公司单方披露；Tesla 未公开确认排他性", '"披露为供应商，未获 Tesla 确认独家"', "上游 PE 估值上限剪 30%"],
        ["Figure $39B 估值合理", "100 万台需 capex $10-13B；未公开 capex-light", "隐含未公开假设；sanity check $0.87B", "2027 H2 回调 30% 后加至 10-15%"],
        ["humanoid IRA 概率 40%", "单点概率仅基于 ITIF 单篇白皮书", "降级 policy watchlist", "改为 watchlist（20-50% × Tesla 兑现度）"],
    ]
    add_table(s, Inches(0.5), Inches(1.5), Inches(12.3), Inches(4.9),
              ["原结论", "新证据", "修订后结论", "对投资动作影响"],
              cv_rows, col_widths=[2.5, 3.8, 2.8, 3.2],
              header_font_size=9, body_font_size=7.5)

    # 底部 元 takeaway
    rrect(s, Inches(0.5), Inches(6.55), Inches(12.3), Inches(0.5), T.NAVY, radius=0.02)
    tx(s, Inches(0.7), Inches(6.65), Inches(12), Inches(0.35),
       "{gold}变更元 takeaway{/gold}  ·  10 处修订中 {bold}7/10 指向美系估值下行风险{/bold} · 中系修订相对中性但板块估值锚整体下移 2-4× · 无一处修订让原结论\"更乐观\"——印证 v1 偏 narrative-driven 而非 DD-driven · 核心架构（5 层产业链 / 技术拐点 / 6 玩家路径）仍立得住",
       size=8.5, color=T.INVERT)

    # =========== CHAPTER 4: 6 玩家路径 ===========
    slide_section(prs, "CHAPTER 4", "6 种玩家路径",
                  "不是\"中 vs 美\"二分 · 是 6 种本质不同的打法", 16, TOTAL)

    # ===== 17. 6 路径全景 + 失败概率柱图（一页）=====
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, T.BG)
    chrome(s, 17, TOTAL, "4.1-4.6 六种打法对比")
    page_title(s, "6 种玩家路径分化 · 核心 moat / 失败概率 / PE 视角",
               "PE 应押\"路径组合\"而非单一公司 · 仅路径 ⑤（中国硬件派）当前已盈利验证")

    # 路径 1-3 (左)
    paths_l = [
        ("①", "Tesla", "垂直整合 + 工厂数据飞轮", T.INFO,
         "Musk Q4 2025: \"not in usage\" · 2026.7-8 Fremont V3 100 万/年", "40-50%"),
        ("②", "Figure / Apptronik", "VLA + B 端 RaaS", T.INFO,
         "Figure $39.5B / Apptronik $5B · BMW 11 月 PoC / BotQ 爬产", "30-40%"),
        ("③", "1X", "消费级 + 远程操作员", T.PURPLE,
         "$10B+ 在谈 · NEO 60-70% 自主率 · 押 2028 95%+", "50-60%"),
    ]
    paths_r = [
        ("④", "PI / Skild", "纯模型层", T.PURPLE,
         "PI $11B / Skild $14B 在谈 · 押 brain 通用化跨形态", "40-50%"),
        ("⑤", "中国硬件派", "极致 BOM + 走量", T.DANGER,
         "{bold}宇树 ¥17 亿营收 / 扣非 ¥6 亿 / 毛利 60%{/bold} · 唯一盈利", "{green}15-25%{/green}"),
        ("⑥", "中国场景派", "国资 + 政府订单 + 仿真", T.DANGER,
         "智元 ¥200亿+/银河 ¥225亿/星海图 ¥200亿/星动 ¥100亿+", "{green}25-35%{/green}"),
    ]
    py = 1.5
    for paths_set, x in [(paths_l, 0.5), (paths_r, 6.95)]:
        py = 1.5
        for label, ttl, mot, c, body, fail in paths_set:
            rrect(s, Inches(x), Inches(py), Inches(6.3), Inches(1.45), T.CARD, radius=0.03)
            rect(s, Inches(x), Inches(py), Inches(0.08), Inches(1.45), c)
            tx(s, Inches(x + 0.2), Inches(py + 0.1), Inches(0.4), Inches(0.5),
               label, size=24, bold=True, color=c)
            tx(s, Inches(x + 0.7), Inches(py + 0.1), Inches(3.5), Inches(0.3),
               ttl, size=12, bold=True, color=T.NAVY)
            tx(s, Inches(x + 0.7), Inches(py + 0.4), Inches(4.5), Inches(0.25),
               mot, size=9, italic=True, color=T.MUTED)
            tx(s, Inches(x + 4.5), Inches(py + 0.12), Inches(1.8), Inches(0.3),
               "失败概率", size=8, color=T.MUTED, align=PP_ALIGN.RIGHT)
            tx(s, Inches(x + 4.5), Inches(py + 0.32), Inches(1.8), Inches(0.35),
               fail, size=14, bold=True, color=c, align=PP_ALIGN.RIGHT)
            tb = s.shapes.add_textbox(Inches(x + 0.2), Inches(py + 0.78), Inches(6.0), Inches(0.6))
            tf = tb.text_frame; tf.word_wrap = True
            tf.margin_left = Emu(0); tf.margin_right = Emu(0); tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
            p = tf.paragraphs[0]
            render(p, body, 8, T.TEXT)
            py += 1.55

    # 底部失败概率柱图
    add_bar_h(s, Inches(0.5), Inches(6.2), Inches(8.0), Inches(1.0),
              ["⑤ 中国硬件派", "⑥ 中国场景派", "② Figure/Apptronik", "① Tesla", "④ PI/Skild", "③ 1X"],
              [20, 30, 35, 45, 45, 55],
              colors=[T.SUCCESS, T.SUCCESS, T.WARN, T.DANGER, T.DANGER, T.DANGER],
              title="失败概率对比（中点估算 %）", show_values=True, max_value=70)

    # 右下：takeaway
    rrect(s, Inches(8.8), Inches(6.2), Inches(4.5), Inches(1.0), T.NAVY, radius=0.02)
    tx(s, Inches(9.0), Inches(6.27), Inches(4.3), Inches(0.3),
       "核心 takeaway", size=10, bold=True, color=T.GOLD)
    tx(s, Inches(9.0), Inches(6.55), Inches(4.3), Inches(0.65),
       "美系 4 路径 5 家估值合计 ~$79.5B / 2025 营收 ~$105M = {bold}PSR 757× 加权{/bold} · ①②③④⑥ 总暴露 $90-100B {bold}全部依赖 Scaling Law 2027-28 兑现{/bold}",
       size=8.5, color=T.INVERT)

    # =========== CHAPTER 5: 产业链 ===========
    slide_section(prs, "CHAPTER 5", "产业链与硬件（补充）",
                  "BOM · 国产化率 · 玩家 · 市场规模", 18, TOTAL)

    # ===== 19. BOM + 国产化率（一页）=====
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, T.BG)
    chrome(s, 19, TOTAL, "5.0-5.2 BOM + 国产化率")
    page_title(s, "单台 BOM 拆解 + 上游硬件国产化率三梯队",
               "中美 BOM 剪刀差 2.8× · 关键卡点：行星滚柱丝杠（仅 19-22%）")

    # 左：BOM 表
    section_label(s, Inches(0.5), Inches(1.5), "5.0 · 单台 BOM 拆解", T.NAVY)
    bom_rows = [
        ["执行器总成", "40-55%", "$14-19K", "50-60%", "$60-80K"],
        ["├ 谐波减速器", "8-12%", "$3-4K", "5-8%", "$6-10K"],
        ["├ 行星滚柱丝杠", "8-15%", "$3-5K", "10-15%", "$12-18K"],
        ["├ 无框力矩电机", "5-8%", "$2-3K", "6-10%", "$8-12K"],
        ["灵巧手", "15-32%", "$5-11K", "17.2%", "$9.5K"],
        ["传感器", "8-15%", "$3-5K", "8-12%", "$10-15K"],
        ["AI 芯片", "5-8%", "$2-3K", "5-8%", "$5-10K"],
        ["电池", "3-6%", "$1-2K", "3-5%", "$4-6K"],
        ["结构 / 装配", "8-12%", "$3-4K", "12-18%", "$15-25K"],
        ["{bold}合计 BOM{/bold}", "100%", "{bold}$32-46K{/bold}", "100%", "{bold}$90-150K{/bold}"],
    ]
    add_table(s, Inches(0.5), Inches(1.85), Inches(6.3), Inches(4.0),
              ["子系统", "中 %", "中金额", "美 %", "美金额"],
              bom_rows, col_widths=[2.2, 0.9, 1.1, 0.9, 1.2],
              header_font_size=9, body_font_size=8)

    # 中美剪刀差 KPI
    rrect(s, Inches(0.5), Inches(6.0), Inches(6.3), Inches(0.95), T.NAVY, radius=0.02)
    tx(s, Inches(0.7), Inches(6.05), Inches(6.1), Inches(0.3),
       "中美 BOM 剪刀差", size=9, bold=True, color=T.GOLD)
    tx(s, Inches(0.7), Inches(6.3), Inches(6.1), Inches(0.6),
       "{big}{bold}2.8×{/bold}{/big}  ·  $32-46K（中）vs $90-150K（美，Morgan Stanley Humanoid 100）", size=11, color=T.INVERT)

    # 右：国产化率柱图
    section_label(s, Inches(7.0), Inches(1.5), "5.2 · 国产化率三梯队", T.GOLD)
    add_bar_h(s, Inches(7.0), Inches(1.85), Inches(6.3), Inches(5.1),
              ["稀土永磁", "高功率电池", "激光雷达", "深度相机",
               "六维力传感", "谐波减速器", "灵巧手",
               "无框电机", "AI 端侧芯片", "RV 减速器",
               "光电编码器", "IMU 导航级",
               "AI 训练芯片", "行星滚柱丝杠", "铍铜"],
              [92, 90, 95, 70, 70, 62, 60, 45, 35, 18, 28, 25, 30, 20, 0],
              colors=[T.SUCCESS, T.SUCCESS, T.SUCCESS, T.SUCCESS, T.SUCCESS, T.SUCCESS, T.SUCCESS,
                      T.WARN, T.WARN, T.WARN, T.WARN, T.WARN,
                      T.DANGER, T.DANGER, T.DANGER],
              title="国产化率 % (2025)", show_values=True, max_value=100)

    # ===== 20. Top 10 估值 / 出货 / 订单（一页三表）=====
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, T.BG)
    chrome(s, 20, TOTAL, "5.4 下游 · Top 10 三视图")
    page_title(s, "下游整机 · 估值 / 出货 / 订单 Top 10 三视图",
               "中系出货榜前六全部包揽 · 美系估值远高于中系")

    # 估值 Top 10 柱图
    add_bar_h(s, Inches(0.4), Inches(1.5), Inches(4.4), Inches(5.5),
              ["Figure AI 🇺🇸", "Skild AI 🇺🇸", "Physical Intel 🇺🇸", "1X 🇺🇸/🇳🇴", "Apptronik 🇺🇸",
               "Neura 🇩🇪", "银河通用 🇨🇳", "智元 🇨🇳", "星海图 🇨🇳", "星动 🇨🇳", "宇树 🇨🇳"],
              [39.5, 14, 11, 10, 5, 4.4, 3, 2.8, 2.8, 1.4, 1.7],
              colors=[T.INFO, T.INFO, T.INFO, T.INFO, T.INFO, T.WARN, T.DANGER, T.DANGER, T.DANGER, T.DANGER, T.DANGER],
              title="估值 Top 10（$B / IPO 预计）", show_values=True)

    # 出货 Top 10
    add_bar_h(s, Inches(4.9), Inches(1.5), Inches(4.0), Inches(5.5),
              ["宇树 (含机器狗)", "智元", "优必选", "众擎", "加速进化", "傅利叶",
               "Agility", "Tesla*", "1X", "Apptronik", "Figure"],
              [5500, 5168, 1500, 1200, 800, 700, 500, 1000, 300, 300, 150],
              colors=[T.DANGER, T.DANGER, T.DANGER, T.DANGER, T.DANGER, T.DANGER,
                      T.INFO, T.LIGHT_MUTED, T.INFO, T.INFO, T.INFO],
              title="2025 出货 Top 10（台）", show_values=True)

    # 订单 Top 10
    add_bar_h(s, Inches(9.0), Inches(1.5), Inches(4.2), Inches(5.5),
              ["优必选 Walker S 全年", "宇树科技 全年",
               "优必选-广西防城港", "优必选-工业客户",
               "优必选-自贡", "智元+宇树-中移动",
               "星动纪元 累计"],
              [12, 12, 2.64, 2.5, 1.59, 1.24, 5],
              colors=[T.SUCCESS, T.SUCCESS, T.SUCCESS, T.SUCCESS, T.SUCCESS, T.WARN, T.WARN],
              title="2025 订单 Top（¥亿）", show_values=True)

    # 底部 takeaway
    rrect(s, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.3), T.SUBTLE, radius=0.02)
    tx(s, Inches(0.7), Inches(7.05), Inches(12), Inches(0.25),
       "* Tesla 自厂 95% 部署但 Musk 承认 \"not in usage\" · 宇树/智元含机器狗+工业版 · 中移动 1.24 亿是数据采集 R&D 合同（Tier D）",
       size=8, italic=True, color=T.MUTED)

    # ===== 21. 市场规模预测 + 渗透率 + 投资项目（一页）=====
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, T.BG)
    chrome(s, 21, TOTAL, "5.5-5.6 市场规模 + 渗透率 + 投资项目")
    page_title(s, "市场规模预测 · 渗透率曲线 vs EV · 2024-2026 投资项目",
               "本报告自算 TAM 中位 $75 亿 (50 万 × $15K) · 比 EV 早期更陡的渗透曲线")

    # 左上：渗透率曲线
    add_line(s, Inches(0.5), Inches(1.5), Inches(6.4), Inches(3.5),
             ["2024", "2025", "2027", "2030", "2035", "2040", "2050"],
             [("人形（累计 / 万台）", [0.3, 2.3, 27, 100, 2200, 11000, 65000]),
              ("EV（累计销量 / 万辆）", [200, 800, 5000, 25000, 80000, 200000, 800000])],
             colors=[T.DANGER, T.INFO],
             title="累计保有量曲线 · 人形 vs EV（对数表示）")

    # 右上：投行预测对比 + 本报告自算
    add_bar_h(s, Inches(7.1), Inches(1.5), Inches(6.2), Inches(3.5),
              ["Citi (2050)", "MS (2050)", "McKinsey (2040)", "Goldman (2035)",
               "BofA (2030)", "BCC (2030)", "MarketsM (2030)",
               "{bold}本报告自算{/bold} (2030)"],
              [7000, 4700, 370, 38, 30, 11, 15.26, 7.5],
              colors=[T.LIGHT_MUTED, T.LIGHT_MUTED, T.LIGHT_MUTED, T.LIGHT_MUTED,
                      T.WARN, T.WARN, T.WARN, T.GOLD],
              title="2030+ 市场规模预测对比（$ 十亿）", show_values=True)

    # 下半：投资项目表
    section_label(s, Inches(0.5), Inches(5.2), "5.6 · 2024-2026 全球投资项目 Top 10", T.GOLD)
    inv_rows = [
        ["1", "智元", "上海临港数据工厂 4,000㎡", "~¥3 亿", "2025.1"],
        ["2", "宇树", "杭州总部+研发基地 IPO 募投", "{bold}¥42 亿{/bold}", "2027"],
        ["3", "Tesla", "Fremont Optimus V3 量产线", "—", "{bold}2026.7-8{/bold}"],
        ["4", "Apptronik", "Austin 扩建", "{bold}$520M{/bold}", "2026.H2"],
        ["5", "1X", "Hayward NEO 工厂（垂直一体化）", "—", "2026.5"],
        ["6", "银河通用", "北京亦庄 + Galbot S1 产线", "¥25 亿（大基金）", "2026"],
        ["7", "Figure AI", "BotQ Austin 工厂 1.2 万台/年", "—", "已投产"],
        ["8", "宁德×智元×千寻", "中州基地电池厂+人形产线", "—", "已投产"],
        ["9", "优必选", "柳州工业人形机器人工厂", "—", "2026 Q1"],
        ["10", "NEOM × 沙特 PIF", "一期 10 万台部署", "—", "2025-2027"],
    ]
    add_table(s, Inches(0.5), Inches(5.55), Inches(12.3), Inches(1.55),
              ["#", "公司", "项目", "金额", "投产"],
              inv_rows, col_widths=[0.3, 1.5, 5.5, 2.0, 1.5],
              header_font_size=8, body_font_size=7)

    # =========== CHAPTER 6: 中美格局 ===========
    slide_section(prs, "CHAPTER 6", "中美格局与政策",
                  "16 环节能力 · 三情景 · 政策 · 估值方法论", 22, TOTAL)

    # ===== 23. 中美 16 环节雷达图 + 双维度三角（一页）=====
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, T.BG)
    chrome(s, 23, TOTAL, "6.1-6.2 中美能力 + 不可能三角")
    page_title(s, "中美 16 环节能力对比 + 双维度不可能三角",
               "美强软件、中强硬件；任何路线都只能优化\"安全/技术/经济\"或\"数据/泛化/成本\"中的两个")

    # 左：16 环节雷达
    add_radar(s, Inches(0.3), Inches(1.4), Inches(7.0), Inches(5.7),
              ["端侧芯片", "训练算力", "谐波", "丝杠", "无框电机", "灵巧手", "力/触觉",
               "稀土永磁", "PEEK/碳纤维", "铍铜", "VLA 模型", "仿真平台",
               "训练数据", "整机量产", "客户场景", "资本市场"],
              [("中国", [3, 2, 4, 2, 3, 4, 4, 5, 4, 1, 3, 2, 4, 5, 5, 3]),
               ("美国", [5, 5, 3, 4, 4, 4, 4, 1, 3, 5, 5, 5, 4, 3, 4, 5])],
              colors=[T.DANGER, T.INFO],
              title="中美 16 环节能力（1-5 ★）")

    # 右：不可能三角 + 文字
    section_label(s, Inches(7.5), Inches(1.5), "6.2 · 双维度不可能三角", T.NAVY)

    # 维度 A
    rrect(s, Inches(7.5), Inches(1.85), Inches(5.7), Inches(2.5), T.CARD, radius=0.02)
    tx(s, Inches(7.65), Inches(1.95), Inches(5.5), Inches(0.3),
       "维度 A · PE 视角", size=10, bold=True, color=T.GOLD)
    tx(s, Inches(7.65), Inches(2.25), Inches(5.5), Inches(0.3),
       "国家安全 × 技术领先 × 经济效益", size=9, italic=True, color=T.MUTED)
    multi_tx(s, Inches(7.65), Inches(2.6), Inches(5.5), Inches(1.65), [
        "美国偏好：{blue}安全 + 技术{/blue}（出口管制 + 资本人才）",
        "中国偏好：{red}安全 + 经济{/red}（自主可控 + 成本敏感）",
        "欧洲偏好：经济 + 部分技术",
        "{italic}类比电池产业链\"国家安全 × 环境保护 × 经济效益\"{/italic}",
    ], size=9)

    # 维度 B
    rrect(s, Inches(7.5), Inches(4.45), Inches(5.7), Inches(2.65), T.CARD, radius=0.02)
    tx(s, Inches(7.65), Inches(4.55), Inches(5.5), Inches(0.3),
       "维度 B · 技术视角", size=10, bold=True, color=T.INFO)
    tx(s, Inches(7.65), Inches(4.85), Inches(5.5), Inches(0.3),
       "数据规模 × 泛化能力 × 商业化成本", size=9, italic=True, color=T.MUTED)
    multi_tx(s, Inches(7.65), Inches(5.2), Inches(5.5), Inches(1.8), [
        "Tesla / Figure：{bold}海量真机数据 + 高泛化 → 高成本（$130K）{/bold}",
        "银河通用：{bold}全仿真合成 + 中等泛化 → 低成本（$30K）{/bold}",
        "宇树 / 众擎：{bold}低成本 + 低泛化（专用场景）→ 走量{/bold}",
        "{italic}任何路线只能优化两个{/italic}",
    ], size=9)

    # ===== 24. 三情景对照 + 政策 + 补贴（一页）=====
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, T.BG)
    chrome(s, 24, TOTAL, "6.3-6.4 三情景 + 政策 + 补贴")
    page_title(s, "三种情景详细对照 + 政策环境 + 中国补贴省级分布")

    # 上半：三情景表格
    sc_rows = [
        ["触发条件", "中国 2027 ramp + 美无 IRA", "Tesla 100 万兑现 + IRA 推出", "双方推动本土化但未脱钩"],
        ["2030 累计保有量", "中 100万 / 美 12万 / 其他 8万", "中 30万 / 美 80万 / 其他 10万", "{bold}中 60万 / 美 25万 / 其他 5万 = 90万{/bold}"],
        ["2030 ASP", "$10K (中) / $35K (美)", "$20K（双方收敛）", "$12K (中) / $40K (美)"],
        ["主要受益方", "中系全链 + 海外硬件依赖中国", "美系一线 + FTA + NVIDIA", "各国\"双轨\" + NVIDIA+CATL+稀土"],
        ["最大风险", "中国本土洗牌过度", "中国失出口 + 算力管制深化", "全球效率损失"],
        ["PE 行动", "重押中系 + 海外硬件", "押美系 + FTA + NVIDIA", "{bold}双线对冲{/bold}"],
    ]
    add_table(s, Inches(0.5), Inches(1.4), Inches(12.3), Inches(2.5),
              ["维度", "China-led (35%)", "US-led (15-25%)", "{bold}Bifurcation (40-50% 基准){/bold}"],
              sc_rows, col_widths=[2.2, 3.4, 3.4, 3.3],
              header_font_size=9, body_font_size=7.5)

    # 下半左：政策对照
    section_label(s, Inches(0.5), Inches(4.1), "6.4 · 政策环境对照", T.NAVY)
    rrect(s, Inches(0.5), Inches(4.45), Inches(6.0), Inches(2.6), T.CARD, radius=0.02)
    rect(s, Inches(0.5), Inches(4.45), Inches(0.06), Inches(2.6), T.INFO)
    tx(s, Inches(0.7), Inches(4.5), Inches(5.8), Inches(0.25),
       "🇺🇸 美国", size=10, bold=True, color=T.INFO)
    multi_tx(s, Inches(0.7), Inches(4.78), Inches(5.8), Inches(2.2), [
        "2025.1 Trump EO 14179 取代 Biden 14110",
        "2025.5 撤销 AI Diffusion Rule",
        "2025.12 H200 case-by-case + 25% 关税",
        "B200/GB300 维持 presumption of denial",
        "2025.9 CA SB 53 / NY RAISE Act",
        "{bold}ITIF 呼吁推 \"humanoid IRA\"{/bold}",
        "{italic}policy watchlist：20-50% × Optimus 兑现{/italic}",
    ], size=8)

    rrect(s, Inches(6.7), Inches(4.45), Inches(6.1), Inches(2.6), T.CARD, radius=0.02)
    rect(s, Inches(6.7), Inches(4.45), Inches(0.06), Inches(2.6), T.DANGER)
    tx(s, Inches(6.9), Inches(4.5), Inches(5.8), Inches(0.25),
       "🇨🇳 中国", size=10, bold=True, color=T.DANGER)
    multi_tx(s, Inches(6.9), Inches(4.78), Inches(5.8), Inches(2.2), [
        "工信部 2023.11 人形机器人创新发展指导意见",
        "{bold}大基金三期 2026.3 首投具身（银河 ¥25 亿）{/bold}",
        "北京 / 上海 / 深圳 / 合肥 / 武汉 各 ¥100 亿+",
        "{bold}累计承诺 ¥550+ 亿{/bold}",
        "5 年年化 ¥110 亿/年 ≈ 0.08% GDP",
        "{red}仅为 EV 2015 高峰 0.51% GDP 的 15-20%{/red}",
        "2025-2027 三年可调用资金 ¥800-1000 亿",
    ], size=8)

    # ===== 25. 估值方法论 + Sanity check + Comps（一页）=====
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, T.BG)
    chrome(s, 25, TOTAL, "6.5 估值方法论")
    page_title(s, "估值 sanity check · Figure / 智元 / 宇树 + Comps Quartile",
               "三家合计 $3.88B / TAM 中位 $7.5B = 52% 集中度合理 · 美系真泡沫 / 中系接近 EV IPO 中位")

    # 左：DCF Sanity check 表
    section_label(s, Inches(0.5), Inches(1.5), "6.5.1 · DCF Sanity check（简化）", T.NAVY)
    sc2_rows = [
        ["2030 出货", "5 万台（10%）", "12 万台（24%）", "14 万（含机器狗）"],
        ["2030 ASP", "$25K", "¥85K ($12K)", "¥60K ($8.5K)"],
        ["2030 收入", "$1.25B", "$1.44B", "$1.19B"],
        ["2030 EBIT", "13% = $163M", "3% = $43M", "8% = $95M"],
        ["WACC", "15%", "18%", "16%"],
        ["{bold}Sanity 估值{/bold}", "{bold}$0.87B{/bold}", "{bold}$0.17B{/bold}", "{bold}$0.46B{/bold}"],
        ["当前估值", "$39.5B", "$2.8B", "$5.6B IPO"],
        ["{red}市场/sanity 比值{/red}", "{red}45×{/red}", "16×", "12×"],
    ]
    add_table(s, Inches(0.5), Inches(1.85), Inches(6.3), Inches(3.3),
              ["参数", "Figure 🇺🇸", "智元 🇨🇳", "宇树 🇨🇳"],
              sc2_rows, col_widths=[1.7, 1.7, 1.7, 1.7],
              header_font_size=9, body_font_size=8)

    # 右：Comps Quartile 柱图
    section_label(s, Inches(7.0), Inches(1.5), "6.5.2 · Comps EV/Sales 倍数", T.GOLD)
    add_bar_h(s, Inches(7.0), Inches(1.85), Inches(6.3), Inches(3.3),
              ["Figure (790×)", "Intuitive (14×)", "Anthropic (40×)", "OpenAI (50×)",
               "蔚来 IPO (53×)", "宇树 (23×)", "智元 (20×)", "小鹏 IPO (24×)",
               "Symbotic (4.5×)", "ABB (2.8×)"],
              [790, 14, 40, 50, 53, 23, 20, 24, 4.5, 2.8],
              colors=[T.DANGER, T.INFO, T.INFO, T.INFO,
                      T.WARN, T.SUCCESS, T.SUCCESS, T.WARN,
                      T.LIGHT_MUTED, T.LIGHT_MUTED],
              title="EV / Sales LTM 倍数对比")

    # 下半：4 套估值方法论 + 关键定位
    rrect(s, Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.7), T.SUBTLE, radius=0.02)
    tx(s, Inches(0.7), Inches(5.5), Inches(12), Inches(0.3),
       "估值方法论 4 套 · 关键定位 PSR 重新校准后", size=10, bold=True, color=T.NAVY)
    multi_tx(s, Inches(0.7), Inches(5.8), Inches(12), Inches(1.3), [
        "① {bold}VC Scorecard{/bold}（Seed-A）：团队 30%+市场 25%+产品 15%+客户 15%+竞争 15%；中位 pre-money $50-200M",
        "② {bold}Berkus{/bold}（Pre-rev）：5 维度 × 各 $0.5-2M",
        "③ {bold}EV/Sales{/bold}（B+）：订单 × 倍数；中系 10-20×，美系 30-80×",
        "④ {bold}DCF on 2030 命中产能{/bold}（C+）：产能 × ASP × 毛利 × 折现",
        "{red}关键定位{/red}：Figure 790× 是真泡沫；{green}智元 20× / 宇树 23× 接近 EV IPO 中位 22×{/green}，不是泡沫；系统性泡沫只在美系一线",
    ], size=8, spacing=1, bullet="")

    # ===== 26. 2026 Q1-Q2 事件年表（一页）=====
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, T.BG)
    chrome(s, 26, TOTAL, "6.6 事件年表")
    page_title(s, "2026 Q1-Q2 关键事件年表",
               "中国具身累计融资 ¥373 亿 · YTD 日均 ¥2.5 亿 · 14 起 ≥¥10 亿单笔")

    timeline = [
        ("2026.1", T.INFO, "商业化 + CES", [
            "优必选 Walker S2 进入比亚迪/吉利 · 累计订单 ¥14-15 亿",
            "北京经开区《具身智能产业新城》行动计划印发",
            "CES 2026 — Hyundai AI 战略 · Atlas 接入 Gemini Robotics",
        ]),
        ("2026.2", T.WARN, "Apptronik $520M 大融资", [
            "{bold}Apptronik Series A-X $520M（估值 $5B / 累计 $935M）{/bold}",
            "国内具身月度融资 ¥160+ 亿（月度新高）",
            "1X 公布 NEO 量产路线：$20K 售卖 / $499 月租",
        ]),
        ("2026.3", T.DANGER, "中国最大单月 + 宇树 IPO 受理", [
            "{bold}银河通用 ¥25 亿 B+，国家大基金三期首投具身{/bold}",
            "星动纪元 ¥10 亿 A+（三星 / 中金 / 国新）· 逐际动力 $2 亿 B（京东 / 上汽 / 蔚来）",
            "{bold}宇树科创板 IPO 受理（3.20）· 募 ¥42.02 亿 · 预计市值 > ¥400 亿{/bold}",
        ]),
        ("2026.4", T.GOLD, "Q1 数据 + 标委会 + 1X 工厂", [
            "{bold}Q1 国内具身融资 ¥200+ 亿（YoY +60%）· 14 起 ≥ ¥10 亿单笔{/bold}",
            "工信部具身智能标委会正式成立（中国电子学会承担秘书处）",
            "{bold}1X Hayward NEO 工厂投产（4.30）· 美国首个垂直一体化人形工厂{/bold}",
        ]),
        ("2026.5 当前", T.SUCCESS, "里程碑式商业化", [
            "优必选 × 空中客车签 Walker S2 采购合作（工业人形 2026 产能 > 1 万台）",
            "{bold}Figure 03 量产爬坡：产线效率 24× 跃升（1 台/天 → 1 台/小时）已下线 350+ 台{/bold}",
            "{bold}Tesla Fremont 终止 Model S/X · 7-8 月启动 Optimus V3 量产线（目标 100 万/年）{/bold}",
        ]),
    ]
    ty = 1.5
    for date, c, title_str, items in timeline:
        rrect(s, Inches(0.5), Inches(ty), Inches(12.3), Inches(1.05), T.CARD, radius=0.02)
        rect(s, Inches(0.5), Inches(ty), Inches(0.08), Inches(1.05), c)
        # Date badge
        rrect(s, Inches(0.7), Inches(ty + 0.13), Inches(1.3), Inches(0.55), c, radius=0.1)
        tx(s, Inches(0.7), Inches(ty + 0.13), Inches(1.3), Inches(0.55),
           date, size=11, bold=True, color=T.INVERT, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        tx(s, Inches(2.1), Inches(ty + 0.1), Inches(10.5), Inches(0.3),
           title_str, size=11, bold=True, color=T.NAVY)
        tb = s.shapes.add_textbox(Inches(2.1), Inches(ty + 0.38), Inches(10.6), Inches(0.65))
        tf = tb.text_frame; tf.word_wrap = True
        tf.margin_left = Emu(0); tf.margin_right = Emu(0); tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
        for li, item in enumerate(items):
            p = tf.add_paragraph() if li else tf.paragraphs[0]
            p.space_after = Pt(0)
            render(p, "• " + item, 8, T.TEXT)
        ty += 1.1

    # =========== APPENDIX A: PE 工具 ===========
    slide_section(prs, "APPENDIX A", "PE 内部分析工具",
                  "deal-level · 与正文行业研究分开", 27, TOTAL)

    # ===== 28. 综合判断 + 优先级矩阵（一页）=====
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, T.BG)
    chrome(s, 28, TOTAL, "A.1-A.2 综合判断 + 优先级")
    page_title(s, "综合判断三段论 + PE 投资优先级矩阵",
               "现状 / 2027-2028 验证 / 2030-2050 长期 + ★★★★★ 标的清单")

    # 上：三段论
    section_label(s, Inches(0.5), Inches(1.5), "A.1 · 综合判断三段论", T.NAVY)
    seg_data = [
        ("现状", T.INFO, "非对称稳态格局",
         "中国主导出货（80-95%）+ BOM 优势（2.8×）；美国主导模型 + 估值\n双向卡脖子：美→中 H100/B200+铍铜；中→美 稀土 92%+镁 87%+电池 90%+PEEK"),
        ("2027-2028", T.WARN, "5 个关键变量验证窗口",
         "1) Tesla 100 万/$20K 是否兑现 · 2) 中系 IPO 兑现 · 3) 550 亿是否催生洗牌\n4) humanoid IRA 是否推出 · 5) VLA Scaling Law 时刻"),
        ("2030-2050", T.SUCCESS, "世代级机会",
         "全球累计：60-150 万 → 5-8 亿 · 年度 TAM：$30-144 亿 → $5-7 万亿\n蓝领劳动力替代 = EV / 光伏 / AI 同等量级"),
    ]
    sy = 1.85
    for ttl, c, sub, body in seg_data:
        rrect(s, Inches(0.5), Inches(sy), Inches(12.3), Inches(0.85), T.CARD, radius=0.02)
        rect(s, Inches(0.5), Inches(sy), Inches(0.08), Inches(0.85), c)
        tx(s, Inches(0.7), Inches(sy + 0.08), Inches(1.3), Inches(0.3),
           ttl, size=10, bold=True, color=c)
        tx(s, Inches(2.1), Inches(sy + 0.08), Inches(10.5), Inches(0.3),
           sub, size=11, bold=True, color=T.NAVY)
        tb = s.shapes.add_textbox(Inches(2.1), Inches(sy + 0.4), Inches(10.6), Inches(0.43))
        tf = tb.text_frame; tf.word_wrap = True
        tf.margin_left = Emu(0); tf.margin_right = Emu(0); tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
        for li, line in enumerate(body.split('\n')):
            p = tf.add_paragraph() if li else tf.paragraphs[0]
            p.space_after = Pt(0)
            render(p, line, 8, T.TEXT)
        sy += 0.9

    # 下：优先级矩阵
    section_label(s, Inches(0.5), Inches(4.6), "A.2 · PE 投资优先级矩阵", T.GOLD)
    prio_rows = [
        ["★★★★★", "基础原材料卡点", "金力永磁/中研股份/中复神鹰/宝武镁业", "已上市 PE 30-60×", "2026 H1"],
        ["★★★★★", "中系卡脖子零部件", "丝杠(贝斯特)/六维力(蓝点)/灵巧手(因时)/IMU(芯动联科)", "未上市 pre-money $50M-$1B", "2026 H1"],
        ["★★★★", "仿真 / 数据基础设施", "国内对标 NVIDIA Isaac/数据工厂", "早期 + AI 估值溢价", "2026"],
        ["★★★★", "中系一线整机 (IPO)", "宇树 / 智元 / 银河通用 / 星动纪元", "¥100-400 亿", "2026 IPO 窗口"],
        ["★★★", "美系一线整机 (限观察)", "Figure / Apptronik / 1X / Skild / π0", "$5-40B", "2027 后调整"],
        ["★★", "美系芯片 / 模型生态", "NVIDIA / Skild / π0 / Generalist AI", "高估值", "长期持有"],
        ["★★", "应用场景集成商", "车厂自研子 / 物流 RaaS / 医疗康复(傅利叶)", "多元", "2026-2028"],
        ["★", "全栈大而全（卷死）", "二线整机", "—", "—"],
    ]
    add_table(s, Inches(0.5), Inches(4.95), Inches(12.3), Inches(2.0),
              ["优先级", "投资方向", "代表标的", "估值锚", "时间窗"],
              prio_rows, col_widths=[1.2, 2.4, 4.5, 2.4, 1.8],
              header_font_size=9, body_font_size=7.5)

    # ===== 29. 组合配置 + 决策树（一页）=====
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, T.BG)
    chrome(s, 29, TOTAL, "A.3 组合配置 + 决策树")
    page_title(s, "PE 组合配置建议 + 关键决策树",
               "三种基金类型 × 7 维度配置 · 三情景判断路径")

    # 左：3 个组合饼图
    section_label(s, Inches(0.5), Inches(1.5), "A.3 · 组合配置 · 三种基金类型", T.NAVY)
    portfolios = [
        ("人民币基金 · 中系重权", T.DANGER, [15, 25, 25, 5, 5, 15, 10]),
        ("美元基金 · 美系重权", T.INFO, [5, 10, 10, 30, 25, 15, 5]),
        ("双货币基金 · 平衡", T.GOLD, [10, 20, 20, 15, 15, 15, 5]),
    ]
    px = 0.5
    portfolio_labels = ["原材料卡点", "中系零部件", "中系整机", "美系软件", "美系整机", "仿真/数据", "应用层"]
    portfolio_colors = [T.WARN, T.DANGER, T.DANGER, T.INFO, T.INFO, T.PURPLE, T.GOLD]
    for ttl, c, data in portfolios:
        tx(s, Inches(px), Inches(1.85), Inches(3.7), Inches(0.3),
           ttl, size=10, bold=True, color=c, align=PP_ALIGN.CENTER)
        add_pie(s, Inches(px), Inches(2.15), Inches(3.7), Inches(2.6),
                portfolio_labels, data, colors=portfolio_colors,
                title="", show_pct=True, show_legend=False)
        px += 4.05
    # Legend at right
    tx(s, Inches(0.5), Inches(4.85), Inches(12.3), Inches(0.3),
       f"{portfolio_labels[0]}（橙）· {portfolio_labels[1]}（红）· {portfolio_labels[2]}（红）· {portfolio_labels[3]}（蓝）· {portfolio_labels[4]}（蓝）· {portfolio_labels[5]}（紫）· {portfolio_labels[6]}（金）",
       size=8, color=T.MUTED, italic=True, align=PP_ALIGN.CENTER)

    # 下：决策树
    section_label(s, Inches(0.5), Inches(5.2), "A.3 · 关键决策树", T.GOLD)
    tree = [
        ("判断 1", "是否相信 China-led（35%）？", T.DANGER, "→ 重押中系全产业链 + 海外硬件依赖中国"),
        ("判断 2", "是否相信 Bifurcation 基准（40-50%）？", T.GOLD, "→ 双线对冲：中系硬件+原材料 + 美系软件 + NVIDIA"),
        ("判断 3", "是否相信 US-led（15-25%）？", T.INFO, "→ 重押 Tesla / Figure / Apptronik / NVIDIA / FTA 国家硬件"),
        ("默认", "观望情景", T.LIGHT_MUTED, "→ 保持现金 + 少量种子（丝杠/灵巧手/IMU/仿真）· 等 2027-2028 兑现窗口"),
    ]
    ty = 5.55
    for label, ques, c, ans in tree:
        rrect(s, Inches(0.5), Inches(ty), Inches(12.3), Inches(0.36), T.CARD, radius=0.02)
        rect(s, Inches(0.5), Inches(ty), Inches(0.06), Inches(0.36), c)
        tx(s, Inches(0.7), Inches(ty + 0.05), Inches(1.0), Inches(0.25),
           label, size=9, bold=True, color=c)
        tx(s, Inches(1.7), Inches(ty + 0.05), Inches(4.5), Inches(0.25),
           ques, size=9, bold=True, color=T.NAVY)
        tx(s, Inches(6.2), Inches(ty + 0.05), Inches(7.0), Inches(0.25),
           ans, size=8.5, color=T.TEXT)
        ty += 0.4

    # ===== 30. Q&A Top 5（一页）=====
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, T.BG)
    chrome(s, 30, TOTAL, "A.4 投委会必问 Q&A")
    page_title(s, "投委会必问 · Top 5 Q&A",
               "反方挑战 + 我们答复 · IC 演练版本")

    qa = [
        ("Q1", "出货 85% 是教育/数据采集，工业落地 < 5%，是补贴催生的 GDP 数字？",
         "① 教育采购量大但单价低；数据采集是 VLA 训练真实需求 ② 工业 PoC 已闭环（比亚迪 2026 计划 20K 台、空客签约）③ BMW-Figure / GXO-Agility 真实 RaaS 收入；但承认中国工业 PoC → 规模化仍需 12-18 月验证"),
        ("Q2", "Tesla 部署 1,000 台 useful work 受质疑，2026 V3 100 万产能不兑现，US-led 概率降到多少？",
         "base case Tesla 2030 = 30/200/100 万；如 2026 Q3 < 10K 台/年，US-led 从 25% → 15%，Bifurcation 升至 50%"),
        ("Q3", "LP 大半在美国，投中系标的政治退出风险如何对冲？humanoid IRA + 中概股退市怎么办？",
         "① 美元基金聚焦\"中国制造但海外销售\"零部件 ② QFII/沪深港通 ③ 二级对冲（做空 Figure / 做多中国零部件）④ 跨境 SPV ⑤ humanoid IRA 出台后美系硬件供应商相对受益"),
        ("Q4", "宇树/智元 IPO 估值 ¥120-400 亿 是合理还是泡沫？vs 蔚来 IPO？",
         "宇树 2025 营收 ¥17 亿 + 扣非净利 ¥6 亿 + 毛利 60.27% — 同期蔚来（亏损）没有的；毛利结构类似 Apple 而非蔚来；DCF sanity $0.46B vs IPO $5.6B = 12× 比值；跟投上限 ¥250 亿安全锚"),
        ("Q5", "如果让你只投 3 个标的，押 5 年 IRR > 30%，是哪 3 个？",
         "① 宇树（A 股 IPO 受理 / 毛利 60%）— IRR 5y 35-45%  ② 金力永磁（Tesla 供应商 / 3-5 年窗口 / Tesla 锁定）— IRR 30-40%  ③ 仿真/数据种子轮（国内对标 NVIDIA Isaac / AgiBot World）— IRR 50%+（10× potential）"),
    ]
    qy = 1.5
    for label, ques, ans in qa:
        rrect(s, Inches(0.5), Inches(qy), Inches(12.3), Inches(1.05), T.CARD, radius=0.02)
        rect(s, Inches(0.5), Inches(qy), Inches(0.08), Inches(1.05), T.GOLD)
        tx(s, Inches(0.75), Inches(qy + 0.08), Inches(0.7), Inches(0.3),
           label, size=14, bold=True, color=T.GOLD)
        tx(s, Inches(1.5), Inches(qy + 0.1), Inches(11.5), Inches(0.3),
           ques, size=10, bold=True, color=T.DANGER)
        tb = s.shapes.add_textbox(Inches(1.5), Inches(qy + 0.42), Inches(11.5), Inches(0.6))
        tf = tb.text_frame; tf.word_wrap = True
        tf.margin_left = Emu(0); tf.margin_right = Emu(0); tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
        p = tf.paragraphs[0]
        render(p, "答：" + ans, 8, T.TEXT)
        qy += 1.1

    # ===== 31. IC Memo 智元（一页）=====
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, T.BG)
    chrome(s, 31, TOTAL, "A.5 IC Memo 智元")
    page_title(s, "IC Memo 样本 · 智元机器人 · CONDITIONAL PROCEED",
               "Lead 早期布局，限 Pre-IPO 轮持仓 ≤ 基金规模 3%")

    # 左：Executive Summary + Company + 5 Pillars
    rrect(s, Inches(0.5), Inches(1.5), Inches(6.5), Inches(5.5), T.CARD, radius=0.02)
    rect(s, Inches(0.5), Inches(1.5), Inches(0.08), Inches(5.5), T.GOLD)
    tx(s, Inches(0.75), Inches(1.6), Inches(6.2), Inches(0.3),
       "1. Executive Summary", size=10, bold=True, color=T.NAVY)
    multi_tx(s, Inches(0.75), Inches(1.9), Inches(6.2), Inches(1.2), [
        "中国 2026 估值最高人形整机厂（约 $2.8B+，IPO 锚 $4-7B）",
        "创始人彭志辉；2025 量产 5,168 台（IDC）；占中国 40% / 全球 30%（Top 1）",
        "Top 3 risks：① 量产爬坡未审计 ② 估值 frothy（但 PSR 20× 接近 EV IPO 中位）③ 软件相对 Helix/π0 优势不明显",
    ], size=8)
    tx(s, Inches(0.75), Inches(3.1), Inches(6.2), Inches(0.3),
       "2. Company / 3. Industry", size=10, bold=True, color=T.NAVY)
    multi_tx(s, Inches(0.75), Inches(3.4), Inches(6.2), Inches(0.85), [
        "产品：远征 A2 (¥620K) / A2 青春版 (¥168K) / 灵犀 X1 (¥109K)",
        "投资方：腾讯/上汽/京东/比亚迪/红杉/北京机器人基金/上海国资",
        "Industry：占中国 2025 出货 40%、全球 ~30%，Top 1 by IDC",
    ], size=8)
    tx(s, Inches(0.75), Inches(4.3), Inches(6.2), Inches(0.3),
       "5. Investment Thesis (5 Pillars)", size=10, bold=True, color=T.NAVY)
    multi_tx(s, Inches(0.75), Inches(4.6), Inches(6.2), Inches(2.4), [
        "(a) 创始人技术信誉 + 华为系供应链整合",
        "(b) 国资背书强（北京机器人基金 / 上海国投 / 上汽 / 比亚迪）",
        "(c) 已建立工业 + 通信巡检 + 数据采集 三场景订单",
        "(d) AgiBot World 100 万+ 真机轨迹开源（中国 ImageNet 时刻）",
        "(e) 量产基地落地速度领先国内同行",
    ], size=8)

    # 右：Returns Analysis 表 + Recommendation
    section_label(s, Inches(7.2), Inches(1.5), "7. Returns Analysis · 5 年 hold", T.GOLD)
    ret_rows = [
        ["Bull (China-led 35%)", "2027 量产 5 万 + 工业兑现", "$35B IPO", "5.0×", "{green}38%{/green}"],
        ["Base (Bifurcation 40%)", "2027 量产 2 万 + Tesla 部分兑现", "$18B", "2.6×", "{green}21%{/green}"],
        ["Bear (US-led 25%)", "2027 量产 < 10K + 美 IRA 出台", "$5B", "0.7×", "{red}-6%{/red}"],
        ["{bold}概率加权 (35/40/25){/bold}", "—", "—", "{bold}3.0×{/bold}", "{bold}~20%{/bold}"],
    ]
    add_table(s, Inches(7.2), Inches(1.85), Inches(6.1), Inches(2.3),
              ["情景", "触发条件", "Exit 估值", "MOIC", "IRR"],
              ret_rows, col_widths=[1.7, 2.0, 1.0, 0.7, 0.7],
              header_font_size=9, body_font_size=7.5)

    # 8. Risks
    tx(s, Inches(7.2), Inches(4.3), Inches(6.1), Inches(0.3),
       "8. Risks（排序 + Mitigants）", size=10, bold=True, color=T.NAVY)
    multi_tx(s, Inches(7.2), Inches(4.6), Inches(6.1), Inches(1.8), [
        "① 量产 yield — mitigant: 季度 yield report + 第三方审计权 + ratchet",
        "② 估值泡沫 — mitigant: IPO < $5B 触发反稀释",
        "③ 关键人风险（彭志辉）— mitigant: 4 年 vesting + key-man insurance",
        "④ 政策反复 — mitigant: 海外 GR + FTA hedging",
    ], size=8)

    # 9. Recommendation
    rrect(s, Inches(7.2), Inches(6.5), Inches(6.1), Inches(0.55), T.SUCCESS, radius=0.02)
    tx(s, Inches(7.4), Inches(6.55), Inches(5.9), Inches(0.25),
       "9. CONDITIONAL PROCEED", size=10, bold=True, color=T.INVERT)
    tx(s, Inches(7.4), Inches(6.78), Inches(5.9), Inches(0.25),
       "跟投 $20-30M · 前提：ratchet + Q4'26 量产数据 + 财务尽调 + 4 年 vesting",
       size=8, color=T.INVERT)

    # =========== APPENDIX B + Closing ===========
    s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s, T.BG)
    chrome(s, 32, TOTAL, "B 数据来源 + 置信度")
    page_title(s, "数据来源与引用追溯 · 206 条 + 置信度声明",
               "Tier 1/2/3 分级 · 数据点 210 / 公司 81 / JSON 见 data/references.json")

    # 左：Tier 饼图
    add_pie(s, Inches(0.5), Inches(1.5), Inches(4.0), Inches(3.5),
            ["Tier 1 一手/权威 · 81 条 (39%)",
             "Tier 2 券商/Bloomberg · 84 条 (41%)",
             "Tier 3 行业媒体 · 41 条 (20%)"],
            [39, 41, 20],
            colors=[T.SUCCESS, T.INFO, T.WARN],
            title="206 引用按 Tier 分布")

    # 右：详细说明
    section_label(s, Inches(4.8), Inches(1.5), "Tier 1/2/3 说明", T.NAVY)
    multi_tx(s, Inches(4.8), Inches(1.85), Inches(8.0), Inches(3.0), [
        "{bold}Tier 1（一手 / 权威）{/bold}：Goldman Sachs · Morgan Stanley · Citi · BofA · McKinsey · Bain · IDC · IFR · USGS · 工信部 · 中国信通院 · GGII · 公司公告 · 招股说明书 · arXiv · BIS 文件",
        "",
        "{bold}Tier 2（券商 / Bloomberg）{/bold}：中信建投 / 东吴 / 华泰 / 招商 / 东方 / 安信 / 国信 / 华宝 / 中泰 / 国金 / 民生 · Bloomberg · Reuters · The Information · 第一财经 · 21 经济网 · 华尔街见闻 · 钛媒体",
        "",
        "{bold}Tier 3（行业媒体）{/bold}：36 氪 · 量子位 · 机器之心 · 智东西 · 艾邦机器人 · 机器人大讲堂 · 知乎专栏",
    ], size=9, bullet="")

    # 下：置信度表
    section_label(s, Inches(0.5), Inches(5.2), "数据置信度声明 · 已知 limitations", T.GOLD)
    conf_rows = [
        ["2025 全球出货", "1.8 万台", "{green}High{/green}", "IDC/信通院/GGII 三方交叉"],
        ["中国出货占全球", "80-95%（按口径）", "{orange}Medium{/orange}", "IDC ~95% / GGII ~85%"],
        ["中美 BOM 剪刀差", "$46K vs $130K", "{orange}Medium{/orange}", "MS Humanoid 100；BofA $35K 为压制后"],
        ["Figure $39.5B 估值", "已落定", "{green}High{/green}", "Figure 官方 + The Information"],
        ["智元 5,168 出货", "—", "{orange}Medium{/orange}", "IDC 转引；原报告未公开"],
        ["三情景概率", "35/15-25/40-50", "{red}Analyst{/red}", "非市场共识"],
        ["humanoid IRA", "Policy watchlist", "{red}Watchlist{/red}", "不给单点概率"],
        ["稀土磁材 92%", "—", "{green}High{/green}", "USGS + 中国稀土协会"],
    ]
    add_table(s, Inches(0.5), Inches(5.55), Inches(12.3), Inches(1.5),
              ["数据点", "报告引用值", "置信度", "备注"],
              conf_rows, col_widths=[2.6, 2.4, 2.0, 5.3],
              header_font_size=9, body_font_size=7.5)

    # ===== Closing =====
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, T.NAVY_DARK)
    rect(s, Inches(0), Inches(0), Inches(0.3), T.H, T.GOLD)
    tx(s, Inches(1), Inches(2.0), Inches(11.5), Inches(1.5),
       "本报告对外严格保密", size=40, bold=True, color=T.INVERT)
    tx(s, Inches(1), Inches(3.0), Inches(11.5), Inches(0.8),
       "仅供 PE 内部决策使用", size=22, color=T.GOLD)
    rect(s, Inches(1), Inches(4.2), Inches(2), Inches(0.04), T.GOLD)
    tx(s, Inches(1), Inches(4.4), Inches(11.5), Inches(0.5),
       "数据截止：2026 年 5 月 · Version 10", size=14, color=T.ICE)
    tx(s, Inches(1), Inches(4.9), Inches(11.5), Inches(0.5),
       "下次更新：随关键 catalyst 触发", size=12, color=T.ICE, italic=True)
    tx(s, Inches(1), Inches(5.2), Inches(11.5), Inches(0.5),
       "（Tesla V3 量产 / 中系 IPO / VLA Scaling Law）", size=11, color=T.ICE_BLUE if hasattr(T, 'ICE_BLUE') else T.ICE)
    tx(s, Inches(1), Inches(6.6), Inches(11.5), Inches(0.4),
       "总产出：6 万中文字 · 35 高密度幻灯片 · 206 引用 · 81 家公司档案",
       size=11, color=T.GOLD)

    return prs

if __name__ == "__main__":
    prs = Presentation()
    prs.slide_width = T.W
    prs.slide_height = T.H
    build(prs)
    output = "/Users/bytedance/Downloads/claude-cowork/embodied_ai_report/embodied_ai_deck.pptx"
    prs.save(output)
    print(f"✅ Saved: {output}")
    print(f"Total slides: {len(prs.slides)}")
