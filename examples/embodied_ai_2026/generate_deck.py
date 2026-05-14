#!/usr/bin/env python3
"""Generate professional PE-grade PPTX for embodied AI report v10.
Uses python-pptx directly. 'Midnight Executive' palette + content from outline.md.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy
from lxml import etree

# ============================================================
# DESIGN SYSTEM · Midnight Executive (PE-grade)
# ============================================================
class Theme:
    # Primary palette
    NAVY = RGBColor(0x0F, 0x27, 0x40)           # Deep navy — titles, dark backgrounds
    NAVY_DARK = RGBColor(0x08, 0x18, 0x2C)      # Even darker — cover gradient
    ICE_BLUE = RGBColor(0xCA, 0xDC, 0xFC)       # Light accent — secondary text on dark
    GOLD = RGBColor(0xD4, 0xA5, 0x74)           # Gold — KPI numbers, callouts
    GOLD_LIGHT = RGBColor(0xF6, 0xB1, 0x7A)     # Lighter gold — hover/secondary

    # Text
    TEXT = RGBColor(0x1A, 0x1A, 0x1A)           # Body dark
    TEXT_MUTED = RGBColor(0x6B, 0x72, 0x80)     # Captions / meta
    TEXT_INVERT = RGBColor(0xFF, 0xFF, 0xFF)    # On dark

    # Backgrounds
    BG = RGBColor(0xFA, 0xFB, 0xFC)             # Page background
    CARD = RGBColor(0xFF, 0xFF, 0xFF)           # Card background
    SUBTLE = RGBColor(0xF1, 0xF4, 0xF8)         # Subtle alternation
    DIVIDER = RGBColor(0xE5, 0xE7, 0xEB)        # Dividers
    TABLE_HEAD = RGBColor(0x0F, 0x27, 0x40)     # Table header dark navy

    # Semantic
    SUCCESS = RGBColor(0x16, 0xA3, 0x4A)        # Green
    WARN = RGBColor(0xEA, 0x58, 0x0C)           # Orange
    DANGER = RGBColor(0xDC, 0x26, 0x26)         # Red
    INFO = RGBColor(0x25, 0x63, 0xEB)           # Blue

    # Slide
    W = Inches(13.333)
    H = Inches(7.5)

    # Fonts (中文友好)
    F_TITLE = "PingFang SC"
    F_BODY = "PingFang SC"

# ============================================================
# Helpers
# ============================================================
def set_slide_bg(slide, color):
    """Solid color background."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color, line_color=None, no_line=True):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if no_line:
        shape.line.fill.background()
    elif line_color:
        shape.line.color.rgb = line_color
    shape.shadow.inherit = False
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color, radius=0.05, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    # Set corner radius via adjustment
    shape.adjustments[0] = radius
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_text(slide, left, top, width, height, text, *,
             font_size=14, bold=False, color=None, italic=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font_name=None):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor

    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    f = run.font
    f.name = font_name or Theme.F_BODY
    f.size = Pt(font_size)
    f.bold = bold
    f.italic = italic
    if color:
        f.color.rgb = color
    return tb

def add_paragraphs(tb, items, *, font_size=12, color=None, bullet_char="•", spacing=4):
    """Add multiple paragraphs to existing text frame with bullets."""
    tf = tb.text_frame
    for i, item in enumerate(items):
        if i == 0 and tf.paragraphs[0].text == "":
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        text, *opts = item if isinstance(item, tuple) else (item,)
        opt = opts[0] if opts else {}
        is_bold = opt.get("bold", False)
        item_color = opt.get("color", color)
        size = opt.get("size", font_size)
        p.alignment = opt.get("align", PP_ALIGN.LEFT)
        p.space_after = Pt(spacing)
        run = p.add_run()
        prefix = bullet_char + " " if bullet_char and not opt.get("nobullet") else ""
        run.text = prefix + text
        f = run.font
        f.name = Theme.F_BODY
        f.size = Pt(size)
        f.bold = is_bold
        if item_color:
            f.color.rgb = item_color

# ============================================================
# Page elements: header bar, footer, page number
# ============================================================
def add_page_chrome(slide, page_num, total, section_name="", section_color=None, on_dark=False):
    """Top accent + bottom footer with page #."""
    # Top thin accent bar
    accent_color = section_color or Theme.GOLD
    add_rect(slide, Inches(0), Inches(0), Theme.W, Inches(0.08), accent_color)

    # Top right: section name
    if section_name:
        text_color = Theme.TEXT_INVERT if on_dark else Theme.TEXT_MUTED
        add_text(slide, Inches(9.5), Inches(0.2), Inches(3.5), Inches(0.3),
                 section_name, font_size=9, color=text_color,
                 align=PP_ALIGN.RIGHT, bold=False)

    # Footer page number (bottom right)
    footer_color = Theme.TEXT_INVERT if on_dark else Theme.TEXT_MUTED
    add_text(slide, Inches(11.5), Inches(7.1), Inches(1.7), Inches(0.3),
             f"{page_num} / {total}", font_size=9, color=footer_color,
             align=PP_ALIGN.RIGHT)
    # Footer left: confidential
    add_text(slide, Inches(0.4), Inches(7.1), Inches(6), Inches(0.3),
             "STRICTLY PRIVATE & CONFIDENTIAL", font_size=8,
             color=footer_color, align=PP_ALIGN.LEFT)

# ============================================================
# Slide templates
# ============================================================
def slide_cover(prs, title, subtitle, meta):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(s, Theme.NAVY_DARK)

    # Gradient stripe (gold) on left
    add_rect(s, Inches(0), Inches(0), Inches(0.3), Theme.H, Theme.GOLD)

    # Top: confidential mark
    add_text(s, Inches(1), Inches(0.6), Inches(11), Inches(0.4),
             "PRIVATE EQUITY · STRICTLY PRIVATE AND CONFIDENTIAL",
             font_size=11, bold=True, color=Theme.GOLD, align=PP_ALIGN.LEFT)

    # Big title
    add_text(s, Inches(1), Inches(2.0), Inches(11.5), Inches(2.2),
             title, font_size=48, bold=True, color=Theme.TEXT_INVERT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)

    # Subtitle
    add_text(s, Inches(1), Inches(4.6), Inches(11.5), Inches(1.0),
             subtitle, font_size=20, color=Theme.ICE_BLUE,
             align=PP_ALIGN.LEFT)

    # Meta (date / version)
    add_text(s, Inches(1), Inches(6.7), Inches(11), Inches(0.4),
             meta, font_size=11, color=Theme.ICE_BLUE,
             align=PP_ALIGN.LEFT, italic=True)

    return s

def slide_section_divider(prs, chapter_label, title, subtitle, page_num, total):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, Theme.NAVY)
    # Gold accent line
    add_rect(s, Inches(1), Inches(3.2), Inches(0.6), Inches(0.08), Theme.GOLD)
    # Chapter label
    add_text(s, Inches(1), Inches(2.5), Inches(11), Inches(0.5),
             chapter_label, font_size=14, bold=True, color=Theme.GOLD,
             align=PP_ALIGN.LEFT)
    # Title
    add_text(s, Inches(1), Inches(3.5), Inches(11.5), Inches(1.5),
             title, font_size=44, bold=True, color=Theme.TEXT_INVERT,
             align=PP_ALIGN.LEFT)
    # Subtitle
    add_text(s, Inches(1), Inches(5.0), Inches(11.5), Inches(1.0),
             subtitle, font_size=18, color=Theme.ICE_BLUE,
             align=PP_ALIGN.LEFT)
    add_page_chrome(s, page_num, total, on_dark=True)
    return s

def slide_title_only(prs, title, subtitle=""):
    """Helper: add title + optional subtitle to top of slide."""
    pass  # Used inside content slide builders

def add_content_title(slide, title, subtitle=""):
    """Add title bar at top of content slide."""
    add_text(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.7),
             title, font_size=24, bold=True, color=Theme.NAVY,
             align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, Inches(0.6), Inches(1.0), Inches(12), Inches(0.4),
                 subtitle, font_size=13, color=Theme.TEXT_MUTED,
                 align=PP_ALIGN.LEFT, italic=True)
    # Thin underline
    add_rect(slide, Inches(0.6), Inches(1.4), Inches(0.6), Inches(0.04), Theme.GOLD)

def slide_bullets(prs, title, bullets, page_num, total, section_name="", subtitle=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, Theme.BG)
    add_page_chrome(s, page_num, total, section_name)
    add_content_title(s, title, subtitle)

    # Bullet list area
    top = Inches(1.7)
    box_w = Inches(12.1)
    box_h = Inches(5.3)

    tb = s.shapes.add_textbox(Inches(0.6), top, box_w, box_h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)

    for i, item in enumerate(bullets):
        if i == 0 and tf.paragraphs[0].text == "":
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        # Render with possible inline bold
        render_inline(p, "▸ " + item, base_size=14, base_color=Theme.TEXT)
    return s

def render_inline(paragraph, text, base_size=14, base_color=None, base_bold=False, base_italic=False):
    """Parse {bold}...{/bold}, {blue}...{/blue}, {italic}...{/italic} into runs."""
    import re
    pattern = re.compile(r'\{(bold|blue|italic|gold|red|green)\}(.*?)\{/\1\}', re.DOTALL)
    pos = 0
    for m in pattern.finditer(text):
        if m.start() > pos:
            add_run(paragraph, text[pos:m.start()], base_size, base_color or Theme.TEXT, base_bold, base_italic)
        tag = m.group(1)
        content = m.group(2)
        is_bold = base_bold or tag == 'bold'
        is_italic = base_italic or tag == 'italic'
        col = base_color or Theme.TEXT
        if tag == 'blue': col = Theme.NAVY
        elif tag == 'gold': col = Theme.GOLD
        elif tag == 'red': col = Theme.DANGER
        elif tag == 'green': col = Theme.SUCCESS
        add_run(paragraph, content, base_size, col, is_bold, is_italic)
        pos = m.end()
    if pos < len(text):
        add_run(paragraph, text[pos:], base_size, base_color or Theme.TEXT, base_bold, base_italic)

def add_run(paragraph, text, size, color, bold=False, italic=False):
    run = paragraph.add_run()
    run.text = text
    f = run.font
    f.name = Theme.F_BODY
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color

def slide_kpi_dashboard(prs, title, kpis, page_num, total, section_name=""):
    """4-8 KPI cards in grid."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, Theme.BG)
    add_page_chrome(s, page_num, total, section_name)
    add_content_title(s, title)

    n = len(kpis)
    cols = 4 if n <= 4 else (4 if n == 8 else 5)
    rows = (n + cols - 1) // cols

    margin = 0.6
    gap = 0.25
    card_w = (13.333 - 2 * margin - (cols - 1) * gap) / cols
    card_h = 1.7 if rows == 1 else 2.0
    top_start = 1.8

    for i, (value, label, sub) in enumerate(kpis):
        r = i // cols
        c = i % cols
        x = margin + c * (card_w + gap)
        y = top_start + r * (card_h + gap)
        # Card background
        card = add_rounded_rect(s, Inches(x), Inches(y), Inches(card_w), Inches(card_h),
                                Theme.CARD, radius=0.05)
        # Gold accent stripe
        add_rect(s, Inches(x), Inches(y), Inches(0.08), Inches(card_h), Theme.GOLD)
        # Value (large)
        add_text(s, Inches(x + 0.2), Inches(y + 0.2), Inches(card_w - 0.3), Inches(0.8),
                 value, font_size=32, bold=True, color=Theme.NAVY,
                 align=PP_ALIGN.LEFT)
        # Label
        add_text(s, Inches(x + 0.2), Inches(y + 1.0), Inches(card_w - 0.3), Inches(0.4),
                 label, font_size=11, bold=True, color=Theme.TEXT,
                 align=PP_ALIGN.LEFT)
        # Sublabel
        if sub:
            add_text(s, Inches(x + 0.2), Inches(y + 1.35), Inches(card_w - 0.3), Inches(0.5),
                     sub, font_size=9, color=Theme.TEXT_MUTED,
                     align=PP_ALIGN.LEFT)
    return s

def slide_table(prs, title, headers, rows, page_num, total, section_name="", subtitle="", note=""):
    """Professional table with navy header + alternating rows."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, Theme.BG)
    add_page_chrome(s, page_num, total, section_name)
    add_content_title(s, title, subtitle)

    # Table area
    table_top = 1.7
    table_left = 0.6
    table_w = 12.1
    avail_h = 5.3 - (0.4 if note else 0)

    n_cols = len(headers)
    n_rows = len(rows) + 1  # +1 for header

    table_shape = s.shapes.add_table(
        n_rows, n_cols,
        Inches(table_left), Inches(table_top),
        Inches(table_w), Inches(avail_h)
    )
    table = table_shape.table

    # Header row
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = ""
        cell.fill.solid()
        cell.fill.fore_color.rgb = Theme.TABLE_HEAD
        tf = cell.text_frame
        tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
        tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0.05)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = str(h)
        run.font.name = Theme.F_BODY
        run.font.size = Pt(10)
        run.font.bold = True
        run.font.color.rgb = Theme.TEXT_INVERT

    # Body rows
    for i, row in enumerate(rows, start=1):
        for j, val in enumerate(row[:n_cols]):
            cell = table.cell(i, j)
            cell.text = ""
            # Alternating row colors
            cell.fill.solid()
            cell.fill.fore_color.rgb = Theme.SUBTLE if i % 2 == 0 else Theme.CARD
            tf = cell.text_frame
            tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
            tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.04)
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            render_inline(p, str(val), base_size=9, base_color=Theme.TEXT)

    if note:
        add_text(s, Inches(0.6), Inches(6.85), Inches(12), Inches(0.25),
                 note, font_size=9, italic=True, color=Theme.TEXT_MUTED)
    return s

def slide_cards(prs, title, cards, page_num, total, section_name="", subtitle="", note=""):
    """N cards in horizontal row (N=2-5). Each card: header + body."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, Theme.BG)
    add_page_chrome(s, page_num, total, section_name)
    add_content_title(s, title, subtitle)

    n = len(cards)
    margin = 0.6
    gap = 0.2
    card_w = (13.333 - 2 * margin - (n - 1) * gap) / n
    card_h = 5.0
    top = 1.7

    for i, card in enumerate(cards):
        x = margin + i * (card_w + gap)
        # Card BG
        add_rounded_rect(s, Inches(x), Inches(top), Inches(card_w), Inches(card_h),
                         Theme.CARD, radius=0.04)
        # Top accent
        accent = card.get('accent', Theme.GOLD)
        add_rect(s, Inches(x), Inches(top), Inches(card_w), Inches(0.08), accent)
        # Header label (if any)
        y_cursor = top + 0.25
        if card.get('label'):
            add_text(s, Inches(x + 0.25), Inches(y_cursor), Inches(card_w - 0.5), Inches(0.3),
                     card['label'], font_size=10, bold=True, color=accent,
                     align=PP_ALIGN.LEFT)
            y_cursor += 0.35
        # Title
        if card.get('title'):
            add_text(s, Inches(x + 0.25), Inches(y_cursor), Inches(card_w - 0.5), Inches(0.7),
                     card['title'], font_size=15, bold=True, color=Theme.NAVY,
                     align=PP_ALIGN.LEFT)
            y_cursor += 0.7
        # Body (list of strings/tuples)
        if card.get('body'):
            tb = s.shapes.add_textbox(Inches(x + 0.25), Inches(y_cursor),
                                       Inches(card_w - 0.5), Inches(top + card_h - y_cursor - 0.25))
            tf = tb.text_frame
            tf.word_wrap = True
            tf.margin_left = Emu(0); tf.margin_right = Emu(0)
            tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
            for idx, item in enumerate(card['body']):
                if idx == 0 and tf.paragraphs[0].text == "":
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.alignment = PP_ALIGN.LEFT
                p.space_after = Pt(4)
                render_inline(p, item, base_size=10, base_color=Theme.TEXT)
    if note:
        add_text(s, Inches(0.6), Inches(6.85), Inches(12), Inches(0.25),
                 note, font_size=9, italic=True, color=Theme.TEXT_MUTED)
    return s

def slide_comparison(prs, title, left, right, page_num, total, section_name="", subtitle=""):
    """2 column comparison."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, Theme.BG)
    add_page_chrome(s, page_num, total, section_name)
    add_content_title(s, title, subtitle)

    col_w = 5.95
    col_gap = 0.3
    col_top = 1.7
    col_h = 5.2

    for i, col in enumerate([left, right]):
        x = 0.6 + i * (col_w + col_gap)
        # Card
        add_rounded_rect(s, Inches(x), Inches(col_top), Inches(col_w), Inches(col_h),
                         Theme.CARD, radius=0.03)
        # Top accent (color)
        accent = col.get('accent', Theme.GOLD if i == 0 else Theme.NAVY)
        add_rect(s, Inches(x), Inches(col_top), Inches(col_w), Inches(0.1), accent)
        # Title
        add_text(s, Inches(x + 0.3), Inches(col_top + 0.25), Inches(col_w - 0.5), Inches(0.5),
                 col['title'], font_size=18, bold=True, color=accent, align=PP_ALIGN.LEFT)
        # Body
        if col.get('body'):
            tb = s.shapes.add_textbox(Inches(x + 0.3), Inches(col_top + 0.85),
                                       Inches(col_w - 0.5), Inches(col_h - 1.0))
            tf = tb.text_frame
            tf.word_wrap = True
            tf.margin_left = Emu(0); tf.margin_right = Emu(0); tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
            for idx, item in enumerate(col['body']):
                if idx == 0 and tf.paragraphs[0].text == "":
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.alignment = PP_ALIGN.LEFT
                p.space_after = Pt(5)
                render_inline(p, "▸ " + item, base_size=11, base_color=Theme.TEXT)
    return s

def slide_hero(prs, title, page_num, total, section_name="", attribution=""):
    """Hero statement — large centered text on dark bg."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, Theme.NAVY)
    add_page_chrome(s, page_num, total, section_name, on_dark=True)

    # Big quote mark
    add_text(s, Inches(1.5), Inches(1.5), Inches(2), Inches(2),
             '"', font_size=120, bold=True, color=Theme.GOLD,
             align=PP_ALIGN.LEFT)

    # Statement
    add_text(s, Inches(1.5), Inches(2.5), Inches(10.5), Inches(3),
             title, font_size=28, bold=True, color=Theme.TEXT_INVERT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)

    if attribution:
        add_text(s, Inches(1.5), Inches(5.8), Inches(10.5), Inches(0.6),
                 "— " + attribution, font_size=14, italic=True, color=Theme.GOLD,
                 align=PP_ALIGN.LEFT)
    return s

def slide_two_column_table_kpi(prs, title, kpis, page_num, total, section_name="", subtitle=""):
    """KPI grid layout - alternative."""
    return slide_kpi_dashboard(prs, title, kpis, page_num, total, section_name)

# ============================================================
# MAIN: build all slides
# ============================================================
def build_deck():
    prs = Presentation()
    prs.slide_width = Theme.W
    prs.slide_height = Theme.H

    TOTAL = 76  # actual

    # =============== SLIDE 1: COVER ===============
    slide_cover(prs,
        title="全球具身智能产业链格局分析",
        subtitle="人形机器人 与 具身智能 产业 PE 视角研究",
        meta="数据截止：2026 年 5 月  |  Version 10  |  机密文档")

    # =============== SLIDE 2: TOC ===============
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, Theme.BG)
    add_page_chrome(s, 2, TOTAL, "TABLE OF CONTENTS")
    add_content_title(s, "目录 / Table of Contents")
    toc_items = [
        ("0", "执行摘要", "当前格局 / 未来推演 / PE 5 个行动决定"),
        ("1", "行业本质", "4 拐点 + 顶尖人物矩阵 + 5 条核心分歧"),
        ("2", "技术路线之争", "VLA vs 世界模型 / 数据 / 形态 / Scaling Law"),
        ("3", "商业化悬崖 (DD)", "Autonomy Level / 订单 Tier / 替代方案 / RaaS"),
        ("4", "6 种玩家路径", "Tesla / Figure / 1X / PI-Skild / 中硬件派 / 中场景派"),
        ("5", "产业链与硬件（补充）", "BOM / 玩家 Top10 / 市场规模"),
        ("6", "中美格局与政策", "16 环节能力对比 / 不可能三角 / 三情景 / 补贴"),
        ("A", "附录 · PE 内部工具", "IC Memo / Q&A / 决策树 / 投资优先级"),
        ("B", "附录 · 数据来源 (206 refs)", "Tier 1/2/3 分级 + 置信度声明"),
    ]
    for i, (num, name, sub) in enumerate(toc_items):
        y = 1.7 + i * 0.55
        # Number circle
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), Inches(y), Inches(0.4), Inches(0.4))
        circle.fill.solid()
        circle.fill.fore_color.rgb = Theme.NAVY
        circle.line.fill.background()
        circle.shadow.inherit = False
        add_text(s, Inches(0.7), Inches(y), Inches(0.4), Inches(0.4),
                 num, font_size=14, bold=True, color=Theme.TEXT_INVERT,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Name + sub
        add_text(s, Inches(1.3), Inches(y - 0.02), Inches(11.5), Inches(0.3),
                 name, font_size=15, bold=True, color=Theme.NAVY,
                 align=PP_ALIGN.LEFT)
        add_text(s, Inches(1.3), Inches(y + 0.22), Inches(11.5), Inches(0.3),
                 sub, font_size=10, color=Theme.TEXT_MUTED,
                 align=PP_ALIGN.LEFT, italic=True)

    # =============== CHAPTER 0: EXEC SUMMARY ===============
    slide_section_divider(prs, "CHAPTER 0", "执行摘要",
        "当前格局核心结论 · 未来推演 · PE 5 个最值得做的决定", 3, TOTAL)

    slide_kpi_dashboard(prs, "8 个关键数字 · KPI Dashboard", [
        ("1.8 万", "2025 全球出货量（台）", "YoY +508%（IDC）"),
        ("80-95%", "中国全球出货占比", "(IDC vs GGII 口径区间)"),
        ("2.8×", "中美 BOM 剪刀差", "$46K vs $130K (MS)"),
        ("$39.5B", "Figure AI 估值", "(美系 Top1)"),
        ("$75 亿", "2030E 年度市场 TAM", "(中位; 区间 $30-144 亿)"),
        ("69-113%", "2025-2030 出货 CAGR", "(按 1.8万→25-80万)"),
        ("¥373 亿", "中国具身累计融资 (YTD)", "日均 ¥2.5 亿"),
        ("¥550+ 亿", "中国地方政府基金", "累计承诺"),
    ], 4, TOTAL, "执行摘要")

    slide_comparison(prs, "当前格局核心结论",
        left={"title": "🇨🇳 中国 · 主导硬件 + 出货", "accent": Theme.DANGER, "body": [
            "全球出货 80-95%（IDC vs GGII 口径）",
            "BOM 优势：中系 $32-46K vs 美系 $130K",
            "全球出货榜前六悉数中国厂商",
            "智元 / 宇树 / 优必选 三家合计 ~70% 中国出货",
            "2025 YTD 累计融资 ¥373 亿",
            "{italic}但：纯人形（剔除机器狗）份额估 70-85%{/italic}",
        ]},
        right={"title": "🇺🇸 美国 · 主导模型 + 估值", "accent": Theme.INFO, "body": [
            "Figure $39.5B / Skild $14B 在谈",
            "Physical Intelligence $11B 在谈",
            "Apptronik $5B / 1X $10B+ 在谈 (2025.9)",
            "美系 4 路径下 5 家估值合计 ~$79.5B",
            "2025 营收估 ~$105M → PSR ~757× 加权",
            "{bold}资本市场把软件大脑 vs 硬件本体按 5-10× 估值分层{/bold}",
        ]},
        page_num=5, total=TOTAL, section_name="执行摘要")

    slide_cards(prs, "未来格局推演 · 三种情景",
        cards=[
            {"label": "情景 ①", "title": "China-led · 35%", "accent": Theme.DANGER,
             "body": [
                "{bold}触发条件{/bold}",
                "中国 2027 ramp + 美无 humanoid IRA",
                "",
                "{bold}2030 累计保有量{/bold}",
                "中 100 万 / 美 12 万 / 其他 8 万",
                "{bold}2030 ASP{/bold}：$10K (中) / $35K (美)",
                "",
                "{bold}PE 行动{/bold}：重押中系全链 + 海外硬件依赖中国",
             ]},
            {"label": "情景 ②", "title": "US-led · 15-25%", "accent": Theme.INFO,
             "body": [
                "{bold}触发条件{/bold}",
                "Tesla 100 万兑现 + humanoid IRA 立法",
                "",
                "{bold}2030 累计保有量{/bold}",
                "中 30 万 / 美 80 万 / 其他 10 万",
                "{bold}2030 ASP{/bold}：$20K (双方收敛)",
                "",
                "{bold}PE 行动{/bold}：押美系 + FTA 国家硬件 + NVIDIA",
             ]},
            {"label": "基准 ③", "title": "Bifurcation · 40-50%", "accent": Theme.GOLD,
             "body": [
                "{bold}触发条件{/bold}",
                "双方推动本土化但未脱钩",
                "",
                "{bold}2030 累计保有量{/bold}",
                "中 60 万 / 美 25 万 / 其他 5 万 = {bold}90 万{/bold}",
                "{bold}2030 ASP{/bold}：$12K (中) / $40K (美)",
                "",
                "{bold}PE 行动{/bold}：双线对冲（中硬件 + 美软件 + NVIDIA）",
             ]},
        ], page_num=6, total=TOTAL, section_name="执行摘要")

    slide_bullets(prs, "12 条关键判断 (1/2)", [
        "{bold}1{/bold} · 中国硬件 + 美国软件已是稳态格局，未来 5 年难以反转",
        "{bold}2{/bold} · 双轨平行是基准情景（概率 40-50%），BOM 重合度将降至 < 30%",
        "{bold}3{/bold} · {bold}2027-2028 是关键验证窗口{/bold}（Tesla Optimus 量产 + 中系 IPO + 中国 550 亿补贴见效 + VLA Scaling Law）",
        "{bold}4{/bold} · 美系一线估值系统性泡沫，2026 H2 后保持观望，2027-2028 预计回调 50-70%",
        "{bold}5{/bold} · {bold}中系核心零部件（绿的、五洲新春、步科、双环）是 Alpha 主战场{/bold}",
        "{bold}6{/bold} · humanoid IRA 推出窗口 2026 H2-2027 H1（policy watchlist，区间 20-50%）",
    ], 7, TOTAL, "执行摘要")

    slide_bullets(prs, "12 条关键判断 (2/2)", [
        "{bold}7{/bold} · 仿真 / 数据基础设施是双轨格局下的稀缺资产（NVIDIA Isaac 接近独占）",
        "{bold}8{/bold} · 中国 550+ 亿补贴年化占 GDP ~0.08%（vs EV 2015 高峰 ~0.51%），量级仅 EV 高峰 15-20%",
        "{bold}9{/bold} · 出口管制最大杀伤力在算力（美→中）+ 稀土/镁/PEEK（中→美）双向反制",
        "{bold}10{/bold} · 2030 年度新增出货 25-80 万台，累计保有量 60-150 万台，{bold}年度 TAM $30-144 亿（中位 $75 亿）{/bold}",
        "{bold}11{/bold} · 稀土永磁是被低估的战略卡点：占 BOM 仅 0.5% 但中国 92% 全球控制",
        "{bold}12{/bold} · 宁德时代 + 智元 + 千寻\"中州基地\"（2025.12）= 中国\"电池厂 + 机器人\"协同范式雏形",
    ], 8, TOTAL, "执行摘要")

    slide_cards(prs, "PE 行动 · 5 个最值得做的决定",
        cards=[
            {"label": "① 2026 H1", "title": "中系核心零部件 + 原材料卡点", "accent": Theme.SUCCESS,
             "body": [
                "未上市丝杠/灵巧手/六维力/IMU",
                "上游金力永磁/中研股份/中复神鹰",
                "",
                "{bold}定价锚{/bold}",
                "2030 命中产能 × PE 15-25×",
             ]},
            {"label": "② 2026 IPO 窗口", "title": "中系一线整机跟投", "accent": Theme.SUCCESS,
             "body": [
                "宇树（2026.3.20 受理）",
                "智元（港股 IPO 中）",
                "银河通用（大基金首投 ¥225 亿）",
                "",
                "{bold}上限{/bold}：¥150-250 亿",
             ]},
            {"label": "③ 2027 后", "title": "美系一线 · 限观察", "accent": Theme.WARN,
             "body": [
                "Figure / Skild / π0 系统性泡沫",
                "",
                "{bold}当前权重 ≤ 5%{/bold}",
                "2027 H2 回调 30% 后加至 10-15%",
                "",
                "美系一线 2027 估值回调 50-70%",
             ]},
            {"label": "④ 2026 起", "title": "仿真 / 数据 / 铍铜种子", "accent": Theme.INFO,
             "body": [
                "国内对标 NVIDIA Isaac",
                "对标 AgiBot World 数据工厂",
                "铍铜国产替代候选",
                "",
                "{bold}早期 + AI 估值溢价{/bold}",
             ]},
            {"label": "⑤ 2026 H2", "title": "7 个关键监测点", "accent": Theme.DANGER,
             "body": [
                "Optimus V3 量产数据",
                "Tesla 100 万兑现度",
                "中国 550 亿补贴落地",
                "中系 IPO 估值",
                "humanoid IRA 立法征兆",
                "VLA Scaling Law 论文",
                "中美稀土/算力反制",
             ]},
        ], page_num=9, total=TOTAL, section_name="执行摘要")

    # =============== CHAPTER 1: 行业本质 ===============
    slide_section_divider(prs, "CHAPTER 1", "行业本质",
        "为什么是现在 · 4 拐点驱动 · 顶尖人物矩阵 · 5 条核心分歧", 10, TOTAL)

    slide_cards(prs, "4 个拐点驱动 · 为什么是 2024-2026", [
        {"label": "①", "title": "技术拐点", "accent": Theme.INFO,
         "body": ["{bold}VLA 范式成熟{/bold}",
                  "2023 RT-2 → 2025 Helix/π0/GR00T",
                  "S1 慢推理 + S2 快控制 + flow matching",
                  "{bold}SOTA{/bold}：OpenVLA-OFT @ LIBERO 97.1%"]},
        {"label": "②", "title": "算力拐点", "accent": Theme.SUCCESS,
         "body": ["{bold}训练成本下降{/bold}",
                  "训 7B VLA ≈ 23K GPU-hour < $100K",
                  "Jetson Thor 端侧 2,070 FP4 TFLOPS",
                  "{bold}仿真{/bold}：Genesis 430,000× 实时"]},
        {"label": "③", "title": "数据拐点", "accent": Theme.WARN,
         "body": ["{bold}飞轮启动{/bold}",
                  "Open X-Embodiment 100 万+ 轨迹",
                  "AgiBot World 100 万+ / 217 任务",
                  "{italic}但 Joel Jang 警告：数据稀缺仍是瓶颈{/italic}"]},
        {"label": "④", "title": "经济拐点", "accent": Theme.DANGER,
         "body": ["{bold}劳动力短缺{/bold}",
                  "美制造业 2025 空缺 85 万 → 2030 缺 210 万",
                  "中国劳动人口年减 800 万 / 2025 出生 792 万",
                  "{italic}但 ROI：美 6.5 年 / 中 12 年——拐点未到{/italic}"]},
    ], 11, TOTAL, "行业本质")

    # Top voices matrix as table
    slide_table(prs, "顶尖人物观点矩阵 · 10 × 5",
        headers=["人物 / 机构", "拐点已到", "工业 PMF 2026", "C 端 2030", "必须双足", "Tesla 会赢"],
        rows=[
            ["a16z (Casado, Hsu) + Polovets", "中立", "同意", "中立", "中立", "{red}不同意{/red}"],
            ["Sequoia (Grady, Huang)", "{green}强同意{/green}", "同意", "同意", "中立", "{red}不同意{/red}"],
            ["Brett Adcock (Figure)", "{green}强同意{/green}", "{green}强同意{/green}", "{green}强同意{/green}", "{green}强同意{/green}", "{red}不同意{/red}"],
            ["Eric Jang (1X)", "同意", "同意", "同意", "同意", "{red}不同意{/red}"],
            ["Hausman/Levine (PI)", "{green}强同意{/green}", "同意", "同意", "中立", "{red}不同意{/red}"],
            ["王兴兴 (宇树)", "同意", "同意", "同意", "同意", "{red}不同意{/red}"],
            ["王鹤 (银河通用)", "中立", "同意", "{red}不同意{/red}", "中立", "{red}不同意{/red}"],
            ["V. Vanhoucke (Waymo)", "同意", "同意", "中立", "中立", "{red}不同意{/red}"],
            ["Yann LeCun (AMI Labs)", "{red}强不同意{/red}", "{red}不同意{/red}", "{red}不同意{/red}", "中立", "{red}强不同意{/red}"],
            ["Marc Raibert (BD)", "同意", "同意", "{red}不同意{/red}", "{red}不同意{/red}", "中立"],
        ], page_num=12, total=TOTAL, section_name="行业本质",
        note="注：非随机抽样；样本不含 Musk/Adam Jonas/Cathie Wood 等 Tesla bulls；纳入后业内分歧大致 5:5")

    slide_cards(prs, "矩阵的 3 个隐藏共识", [
        {"label": "共识 ①", "title": "Tesla 会赢 · 0 支持票", "accent": Theme.DANGER,
         "body": ["最显著的\"隐而不宣\"共识",
                  "Adcock 公开说 Figure 比 OpenAI 强",
                  "Skild/PI 走 brain 路线绕开特斯拉硬件",
                  "中国阵营走低成本供应链",
                  "Boston Dynamics 绑定现代汽车",
                  "{italic}注：样本不含 Musk/Jonas/Wood；带 bulls 后分歧 5:5{/italic}"]},
        {"label": "共识 ②", "title": "时间表激进 + 全乐观", "accent": Theme.WARN,
         "body": ["{bold}只有 Brett Adcock 一人{/bold}",
                  "\"10 年内每家都有 humanoid\"",
                  "BMW 已 5 月部署 Figure 02",
                  "这是 Figure $39B 估值的双刃剑",
                  "任何 demo 真实性争议都直接打击估值"]},
        {"label": "共识 ③", "title": "必须双足 · Marc Raibert 反对", "accent": Theme.INFO,
         "body": ["{bold}唯一公开反对{/bold}",
                  "\"让机器人像人不是腿和手，是智能\"",
                  "Boston Dynamics 创始人偏好轮式",
                  "对工业场景：轮式+双臂+灵巧手",
                  "可能 95% 替代双足 + 成本低 30-50%"]},
    ], 13, TOTAL, "行业本质")

    slide_comparison(prs, "5 条核心分歧线 · GPT 时刻 + VLA vs 世界模型",
        left={"title": "分歧 ① · GPT 时刻到了吗", "accent": Theme.SUCCESS, "body": [
            "{bold}乐观派{/bold}",
            "Brett Adcock：前夜",
            "王兴兴：1-2 年到",
            "Levine：Apollo program",
            "",
            "{bold}悲观派{/bold}",
            "LeCun：LLM 5 年内 useless（$1B 押 AMI Labs）",
            "Karpathy：ghosts of behavior",
            "Casado (a16z 内部)：估值疯狂",
        ]},
        right={"title": "分歧 ② · VLA vs 世界模型", "accent": Theme.INFO, "body": [
            "{bold}VLA 派{/bold}",
            "Physical Intelligence (π0/π0.5)",
            "Figure (Helix 02)",
            "Tesla / Apptronik",
            "",
            "{bold}世界模型派{/bold}",
            "1X World Model（Eric Jang）",
            "NVIDIA GEAR (Joel Jang)",
            "LeCun → AMI Labs（$1B）",
            "王兴兴 UnifoLM-WMA-0",
        ]},
        page_num=14, total=TOTAL, section_name="行业本质")

    slide_comparison(prs, "5 条核心分歧线 · 数据 / 形态 / 谁会赢",
        left={"title": "分歧 ③ · 数据来源", "accent": Theme.WARN, "body": [
            "{bold}真机派（美主流）{/bold}",
            "Hausman / Levine / Adcock",
            "Open X 100 万+ 轨迹",
            "",
            "{bold}仿真派（中国押注）{/bold}",
            "银河通用 10 亿帧",
            "NVIDIA Isaac / Genesis",
            "王鹤：\"中国弯道超车\"",
            "",
            "{bold}视频派（学术）{/bold}",
            "Joel Jang / Lerrel Pinto",
            "Ego4D 3,670 小时",
        ]},
        right={"title": "分歧 ④⑤ · 形态 + 谁会赢", "accent": Theme.DANGER, "body": [
            "{bold}必须双足{/bold}：Adcock / Bornich / 王兴兴",
            "{bold}形态多元{/bold}：Marc Raibert / Skild omni-bodied",
            "{bold}非人形派{/bold}：Symbotic / Locus 已跑通商业化",
            "",
            "{bold}美国阵营{/bold}",
            "通用 brain (PI/Skild) vs 垂直整机 (Figure/Apptronik/1X)",
            "a16z 押 Figure/Skild/PI；NVIDIA 全押",
            "",
            "{bold}中国阵营{/bold}",
            "王兴兴：硬件不是瓶颈，模型架构才是",
            "王鹤：仿真+合成是中国弯道超车",
        ]},
        page_num=15, total=TOTAL, section_name="行业本质")

    slide_hero(prs,
        title="Within 10 years, every home will have a humanoid.",
        page_num=16, total=TOTAL, section_name="行业本质 · 引语",
        attribution="Brett Adcock (Figure CEO), Time Magazine Oct 2025")

    slide_bullets(prs, "9 条最具穿透力引语 (剩余)", [
        "{bold}Sergey Levine{/bold} (Dwarkesh 9/2025) — \"It's more like the Apollo program than a science experiment.\"",
        "{bold}Yann LeCun{/bold} (10/27/2025) — \"Large Language Models will become useless within 5 years.\"",
        "{bold}Leo Polovets{/bold} (Humba VC, a16z 播客) — \"Humanoids are probably one of the most hyped areas where valuations get crazy before any revenue.\"",
        "{bold}王兴兴{/bold} (21 经济 8/2025) — \"数据关注度有点太高，最大问题在模型；当前 VLA 是傻瓜式架构。\"",
        "{bold}Marc Raibert{/bold} (Digitimes 4/2025) — \"工厂先行，家用还要等；让机器像人的不是腿和手，是智能。\"",
        "{bold}王鹤{/bold} (36Kr) — \"5-10 年才能什么活儿都干；特别不建议讲具身 AGI。\"",
        "{bold}Oliver Hsu, a16z{/bold} (1/13/2026) — \"Real-world deployments still largely confined to demos and pilot programs.\"",
        "{bold}Andrej Karpathy{/bold} (YC 6/2025) — \"Today's AI is missing hippocampus, amygdala, cerebellum.\"",
        "{bold}Joel Jang, NVIDIA GEAR{/bold} (2/6/2026) — \"Physical AGI's bottleneck is data scarcity, not hardware.\"",
    ], 17, TOTAL, "行业本质 · 引语")

    # =============== CHAPTER 2: 技术路线 ===============
    slide_section_divider(prs, "CHAPTER 2", "技术路线之争",
        "真正的胜负所在 · VLA vs 世界模型 · 数据 · 形态 · Scaling Law", 18, TOTAL)

    slide_comparison(prs, "路线 ① · 端到端 VLA vs 显式世界模型",
        left={"title": "端到端 VLA 派", "accent": Theme.INFO, "body": [
            "{bold}核心假设{/bold}",
            "参数 scale + 数据 scale + flow matching 足以解决泛化",
            "",
            "{bold}代表玩家{/bold}",
            "Physical Intelligence (π0/π0.5)：Flow-matching + PaliGemma",
            "Figure (Helix 02)：S1 慢推理 + S2 快控制",
            "Tesla：垂直整合自家数据",
            "Apptronik：与 Google DeepMind 合作",
            "",
            "{bold}PE 投注{/bold}",
            "Figure $39B / PI $11B / Skild $14B 在谈",
        ]},
        right={"title": "世界模型派", "accent": Theme.SUCCESS, "body": [
            "{bold}核心假设{/bold}",
            "必须显式建模物理世界因果，才能 zero-shot 泛化",
            "",
            "{bold}代表玩家{/bold}",
            "1X World Model：数据生成飞轮",
            "NVIDIA GEAR (GR00T N1/N1.6)：data pyramid",
            "Yann LeCun → AMI Labs：V-JEPA2 → AMI，$1B 反 LLM",
            "王兴兴 (宇树)：UnifoLM-WMA-0 开源",
            "",
            "{bold}融资规模{/bold}",
            "学术声誉强但融资规模相对小",
        ]},
        page_num=19, total=TOTAL, section_name="技术路线")

    slide_table(prs, "路线 ② · 数据来源 · 真机 vs 仿真 vs 视频",
        headers=["路径", "代表玩家", "核心思想", "规模", "关键风险"],
        rows=[
            ["真机数据飞轮 (美主流)", "PI / Figure / Hausman", "Deploy at scale, gather real data", "Open X 100 万+ 轨迹", "真机数据贵且慢"],
            ["仿真合成 (中国押注)", "银河通用 / NVIDIA / Genesis", "10 亿帧仿真训练 / 0 真机", "GraspVLA 10 亿帧 / Genesis 430,000× 实时", "Sim-to-Real Gap 是已知未解"],
            ["人类视频 (学术派)", "Joel Jang / Lerrel Pinto", "Internet-scale video learning", "Ego4D 3,670h / V-JEPA2 100 万小时", "从视觉到 action 仍需 bridge"],
        ], page_num=20, total=TOTAL, section_name="技术路线",
        note="这是中美技术路线根本分歧 — 美押真机飞轮，中押仿真+合成；2026-2028 会有交叉验证")

    slide_cards(prs, "路线 ③ · 必须双足吗 · 形态之争", [
        {"label": "派别 A", "title": "必须双足派", "accent": Theme.DANGER,
         "body": ["{bold}Adcock / Bornich / 王兴兴{/bold}",
                  "\"世界是为人造的\"",
                  "",
                  "代表：Figure F.03 / 1X NEO / Optimus / 宇树 G1",
                  "",
                  "{bold}BOM{/bold}：$30-130K"]},
        {"label": "派别 B", "title": "形态多元派", "accent": Theme.WARN,
         "body": ["{bold}Marc Raibert (BD){/bold}",
                  "\"让机器人像人不是腿和手，是智能\"",
                  "",
                  "Skild AI \"omni-bodied\"",
                  "张巍 (逐际动力) 偏多形态",
                  "",
                  "{bold}BOM{/bold}：$15-50K（轮式 + 双臂 + 灵巧手）"]},
        {"label": "派别 C", "title": "非人形派 · 已赢", "accent": Theme.SUCCESS,
         "body": ["{bold}Symbotic / AutoStore / Locus / UR{/bold}",
                  "用\"非人形\"已跑通商业化",
                  "",
                  "Locus 17K AMR / 70 亿次拣选",
                  "Symbotic Walmart 40% 交付成本下降",
                  "AutoStore 99.8% uptime / 1,950 站点",
                  "",
                  "{bold}BOM{/bold}：$5-50K"]},
    ], 21, TOTAL, "技术路线")

    slide_comparison(prs, "Scaling Law 现状 · GPT 时刻到了吗",
        left={"title": "当前 SOTA 实测", "accent": Theme.INFO, "body": [
            "LIBERO 4 套件：OpenVLA-OFT {bold}97.1%{/bold}（受限基准）",
            "CALVIN：π0 ~75-85%",
            "SimplerEnv：GR00T N1 ~60-70%",
            "{bold}真实环境（光照变化、未见物体）：实测约 30-50%{/bold}",
            "1X NEO 官方承认：60-70% autonomy at launch",
            "Tesla Optimus 工厂：Musk Q4 2025 \"not in usage in a material way\"",
        ]},
        right={"title": "Scaling Law 预测", "accent": Theme.WARN, "body": [
            "{bold}Russ Tedrake (TRI LBM 8/2025){/bold}",
            "\"capabilities scale more than linearly with data once threshold crossed\"",
            "",
            "1700h 真机 + 47,000 仿真 rollout 是当前最大 LBM",
            "",
            "{bold}真正 GPT 时刻预计 2027-2028{/bold}",
            "（数据 1 亿+ 轨迹 / 模型 30B+）",
            "",
            "{italic}Chelsea Finn (PI): scale is necessary, but subordinate to solving the problem{/italic}",
        ]},
        page_num=22, total=TOTAL, section_name="技术路线")

    # =============== CHAPTER 3: 商业化悬崖 ===============
    slide_section_divider(prs, "CHAPTER 3", "商业化悬崖 (DD)",
        "从 Demo 到 Product 的真实距离 · Autonomy / 订单 Tier / 替代方案", 23, TOTAL)

    slide_hero(prs,
        title='媒体宣传 "~$3B 人形订单" 中 Tier A 接近 0；行业普遍 Tier D (PoC)；人形相对 AMR/协作臂 ROI 差 5-16×',
        page_num=24, total=TOTAL, section_name="商业化悬崖",
        attribution="本章核心结论 · DD 真相")

    slide_table(prs, "真实 Autonomy Level · 部署评分卡 (1/2)",
        headers=["公司 / 型号", "部署点", "Autonomy", "运行数据", "关键披露"],
        rows=[
            ["Tesla Optimus", "Fremont (R&D) 无客户", "L0-L1", "零规模化部署", "Musk Q4 2025: \"not in usage in our factories in a material way\""],
            ["Figure 02", "BMW Spartanburg", "L3（争议）", "1,250h / 90,000+ 钣金 / 30,000+ X3", "Adcock: \"no teleop\"；Scott Walter 等指演示有遥操作征兆"],
            ["Figure 03", "BMW + Munich 分阶段", "声称 L3（独立验证 [UNDISCLOSED]）", "2026 分阶段铺开", "任何未来发现 teleop 都将打击 $39B 估值"],
            ["1X NEO", "消费者 pre-order", "L0-L2 混合", "{bold}CEO 公开：60-70% 自主率{/bold}", "\"Expert Mode\" 远程操作员；CEO: \"running towards a cliff\""],
            ["Apptronik Apollo", "Mercedes / GXO / Jabil", "L2 (lab-trained)", "仍处 pilot", "商业 scale 目标 2026 H2"],
        ], page_num=25, total=TOTAL, section_name="商业化悬崖",
        subtitle="基于 SAE 改编 L0-L4 分级")

    slide_table(prs, "真实 Autonomy Level · 部署评分卡 (2/2)",
        headers=["公司 / 型号", "部署点", "Autonomy", "运行数据", "关键披露"],
        rows=[
            ["Agility Digit", "GXO / Mercado Libre", "{green}L3 in defined ODD{/green}", "100,000+ 货箱 (2025.11) / 2:1 电池循环", "迄今唯一公开商业 RaaS datapoint"],
            ["智元远征 A2", "中国移动数据采集 PoC", "L0-L1", "data-collection service contract", "实质 R&D 数据采集，不是生产部署"],
            ["宇树 H1/G1", "demo / 教育 / 科研", "L0-L1 工厂", "无 24×7 生产部署", "G1 营收占公司 60%（机器狗+教育为主）"],
            ["BD Atlas (新)", "Hyundai / Google DeepMind", "L2-L3 in part-sorting ODD", "CES 2026 商业揭幕", "2026 整年订单已锁定"],
            ["UBTech Walker S2", "中国政府/工业站点", "L1-L2", "11 月交付节奏", "政府采购投标基本绑定，但端用大多为非生产"],
        ], page_num=26, total=TOTAL, section_name="商业化悬崖",
        note="⚠ MTBF / MTBI 数字全行业缺失 · 任何省略 MTBF 的 PE pitch deck 应视为尽调不完整")

    slide_cards(prs, "订单质量 · Tier A-E 分级框架", [
        {"label": "Tier A", "title": "已收款 + 不可撤销", "accent": Theme.SUCCESS,
         "body": ["PO + 付款 + 单台已交付", "", "取消风险 {bold}< 5%{/bold}"]},
        {"label": "Tier B", "title": "Firm 采购合同", "accent": Theme.SUCCESS,
         "body": ["单价 + 交付时间表锁定", "", "取消风险 {bold}5-15%{/bold}"]},
        {"label": "Tier C", "title": "Master 框架 / MSA", "accent": Theme.WARN,
         "body": ["单台需独立 PO 下单", "", "取消风险 {bold}30-60%{/bold}"]},
        {"label": "Tier D", "title": "付费 PoC", "accent": Theme.DANGER,
         "body": ["无生产承诺", "", "取消风险 {bold}60-80%{/bold}"]},
        {"label": "Tier E", "title": "LOI / MoU", "accent": Theme.DANGER,
         "body": ["无法律约束力", "", "取消风险 {bold}75-95%{/bold}"]},
    ], 27, TOTAL, "商业化悬崖")

    slide_table(prs, "Top 10 公开订单 · 重新分级",
        headers=["#", "交易", "金额", "实际 Tier", "关键 caveat"],
        rows=[
            ["1", "UBTech Walker S 2025 全年", "¥11-13 亿", "B / C mix", "政府采购投标授标即绑定；交付进行中"],
            ["2", "UBTech-广西防城港", "¥2.64 亿", "B", "Walker S2；2025.12 交付"],
            ["3", "UBTech-自贡", "¥1.59 亿", "B", "Walker S2；数据采集中心"],
            ["4", "宇树 2025 工业订单", "¥12 亿", "B / C mix", "G1 占 60% 营收；多为教育+小批量"],
            ["5", "智元+宇树-中移动", "¥1.24 亿", "{red}D (PoC) 含 B{/red}", "数据采集 R&D 服务合同，非生产"],
            ["6", "星动纪元 2025 累计", "¥5 亿", "B/C/D mix", "50% 海外大概率 D-tier"],
            ["7", "Apptronik × Merc/GXO/Jabil", "[UNDISCLOSED]", "D (PoC)", "商业 scale 目标 2026 H2"],
            ["8", "Agility × GXO / Mercado", "[UNDISCLOSED] / 100K+ 货箱", "{green}B (GXO) + C{/green}", "迄今最经实测的人形商业 datapoint"],
            ["9", "Galbot × Baida/CATL/Bosch", "\"1,000+ 计划\"", "C / E", "未公开单价×数量×交付表"],
            ["10", "Figure × BMW", "未披露 / 11 月 PoC", "B + C", "Fortune 报道 BMW 反驳 Adcock"],
            ["—", "Tesla 自厂 1,000 台", "n/a", "{red}剔除（不是订单）{/red}", "Musk Q4 2025: \"not in usage\""],
        ], page_num=28, total=TOTAL, section_name="商业化悬崖",
        note="媒体引述 ~$3B 中：真 firm Tier A+B ≈ $1-1.5B；剩余 $1.5-2B 是 C/D/E")

    slide_table(prs, "替代方案 ROI · 仓储战争已被非人形赢得",
        headers=["方案", "代表厂商", "单台价格", "部署规模", "已公开回本", "关键 KPI"],
        rows=[
            ["AMR", "Locus Robotics", "RaaS-only", "17,000+ AMR / 70 亿+ 拣选", "{green}< 6 个月{/green}", "—"],
            ["Cube ASRS", "AutoStore", "$1-5M/system", "1,950+ 系统全球", "20-26 周安装", "{green}99.8% uptime{/green}"],
            ["Goods-to-person", "Symbotic", "$10M+/系统", "Walmart 42 DC / 承诺 400 系统", "20% opex 降 / 200% 吞吐提升", "FY2025 营收 $2.25B"],
            ["协作臂 6 轴", "Universal Robots", "$11-60K", "UR 累计 > 100,000 cobot", "3 月-2 年（典型 6-12 月）", "—"],
            ["{red}人形机器人{/red}", "Figure/Apptronik/Agility/Tesla", "$25-130K+", "仅几千台（多为 demo）", "{red}Agility naive ~8 年{/red}", "MTBF/uptime 全行业未披露"],
        ], page_num=29, total=TOTAL, section_name="商业化悬崖",
        note="McKinsey 2024：人形 payback 5.3 → 2.8 年（仍是 Locus 3-5×） · Marc Raibert (Lex Fridman): wheeled is far more efficient than bipedal")

    slide_bullets(prs, "安全 / 合规 · 工业部署的法律护城河", [
        "{bold}ISO 10218 / ANSI R15.06{/bold} 主要覆盖工业机器人本体与系统集成；{bold}动态稳定腿式/双足移动机器人风险仍处于新标准化项目{/bold}（A3 / IEC TC 184）与解释适用阶段",
        "A3 已识别 humanoid 属于 \"dynamically stable industrial mobile robot\"，现有 R15.06 / R15.08 并未完全覆盖这一形态的特定危险",
        "协作机器人 (cobots) 有 {bold}ISO/TS 15066{/bold} 力 / 压等级标准；人形机器人在工业用途的安全框架仍在制定中",
        "{bold}责任划分{/bold}：自主水平越高，责任越倾向 OEM + 软件商（与传统\"集成商承担风险\"模型反转）",
        "{bold}家用 / 医疗几乎完全空白{/bold}：1X NEO、Optimus 进家庭的 liability framework 全行业空白",
        "{bold}数据隐私{/bold}：机器人采集的家庭 / 工厂数据如何合规 unclear；远程操作员（1X NEO \"Expert Mode\"）涉及 GDPR / 美国州法律待澄清",
        "{bold}PE 行动{/bold}：工业部署现实是\"caged-cell\" 仍是路径最短选择 — cobots 和 AMR 在工业场景压制人形的法律基础",
    ], 30, TOTAL, "商业化悬崖")

    slide_comparison(prs, "RaaS 资产负债表 · Figure $39B 隐含假设",
        left={"title": "RaaS 模式真实经济", "accent": Theme.WARN, "body": [
            "Agility 公开口径：$30/hr ($10-12/hr opex)",
            "",
            "Figure RaaS 定价 (3 种 inferred，未官方披露)：",
            "(a) $1,000/月 = $12K/年",
            "{bold}(b) $25/h × 8h × 250d = $50K/年 [本报告 base]{/bold}",
            "(c) 2030 量产 blended $25K/台",
            "",
            "Digit 单台 $250K capex → naive payback ~8 年",
            "RaaS 需要 OEM 前置 capex + 维修 + 残值 + 信用风险",
        ]},
        right={"title": "Figure $39B 估值隐含", "accent": Theme.DANGER, "body": [
            "本报告 base case ($50K/年 单班 hourly):",
            "",
            "2030 出货 5 万 × $50K = $2.5B 年收入",
            "$39.5B / $2.5B = {bold}15.8× 2030 收入{/bold}",
            "但 2025 营收估 $15-30M [INFERRED]",
            "2025 PSR 仍 ~790×",
            "",
            "5 万台部署需 capex ~$5-7B：",
            "{bold}capex-light 资产负债表（off-balance-sheet）假设未公开确认{/bold}",
        ]},
        page_num=31, total=TOTAL, section_name="商业化悬崖")

    slide_bullets(prs, "第三章 PE 关键 takeaways", [
        "① 人形机器人当前商业化{bold}远未跑通{/bold} — Musk 亲口承认 Optimus \"无实质使用\"，1X 60-70% 自主率，Figure 与 BMW 关系被 Fortune 质疑",
        "② 媒体宣传 \"~$3B 订单\" 中 {bold}Tier A 几乎为零{/bold}；Tier B 集中在中国政府采购（数据采集 / 安防 / 教育，非 ROI 驱动）",
        "③ 仓储战争{bold}已被 Symbotic / Locus / AutoStore 赢得{/bold}；人形相对成熟方案 ROI 差 5-16×",
        "④ ISO 10218 / ANSI R15.06 {bold}不涵盖双足移动安全{/bold} — 工业部署的法律不在人形侧",
        "⑤ Figure $39B 估值依赖 capex-light 资产负债表假设，{bold}但未公开确认{/bold}",
        "⑥ {bold}\"任何不含 MTBF 的 PE deck 都不算完整尽调\"{/bold} — 全行业未披露的关键指标",
    ], 32, TOTAL, "商业化悬崖")

    # Investment thesis change table - 2 pages
    slide_table(prs, "投资结论与决策映射表 (1/2)",
        headers=["原结论", "新证据", "修订后结论", "对投资动作影响"],
        rows=[
            ["RaaS 5 年现金流是售卖 170×", "Figure $25/h×8h×250d ≈ $50K/年；LTV:CAC 仅 2-5×", "{bold}RaaS 没有数量级优势{/bold}", "美系一线整机权重 ≤ 5%（vs v1 15-25%）"],
            ["2030 出货 150-300 万台，TAM $300-800 亿", "混淆年度新增 vs 累计；中位 ASP $15K", "年度新增 25-80 万 / 累计 60-150 万；TAM $30-144 亿", "{bold}整机估值锚下移 2-4×{/bold}：Figure → sanity $0.87B"],
            ["仓储是人形最先 PMF", "Symbotic $2.25B / Walmart 40% 交付成本降；Locus <6 月", "仓储已被非人形赢得", "Agility / Apptronik 权重 ≤ 3%；优先 cobot+AMR+ASRS"],
            ["Tesla 部署 1,000 台 = 量产前奏", "Musk Q4 2025: not in usage in our factories", "Tesla 未规模化使用", "US-led 概率 15-25%（vs v1 25%）；Bifurcation 40-50%"],
            ["1X NEO = 真自主家用", "CEO 公开承认自主率 60-70%；目标 2028 95%+", "1X 当前是远程操作员+AI 混合", "1X 仓位 ≤ 3%（binary 押 2028 自主率突破）"],
        ], page_num=33, total=TOTAL, section_name="商业化悬崖")

    slide_table(prs, "投资结论与决策映射表 (2/2)",
        headers=["原结论", "新证据", "修订后结论", "对投资动作影响"],
        rows=[
            ["中移动 1.24 亿 = 全球最大订单", "合同实质是 R&D 数据采集服务", "降级 Tier D PoC", "智元/宇树估值剔除；中系订单 Tier A+B 真实 ~$1.5B (vs $3B 引述)"],
            ["中国 550 亿补贴超 EV 早期占 GDP", "年化 / 年 GDP ≈ 0.08%（vs EV 2015 高峰 ~0.51%）", "补贴量级仅 EV 高峰 15-20%", "补贴权重在 China-led 推导中 ≤ 10%；核心论据回归 BOM+场景+数据"],
            ["金力永磁/中研股份独家供应 Tesla", "公司单方披露；Tesla 未公开确认排他性", "改\"披露为供应商，未获 Tesla 确认独家\"", "上游 PE 估值上限剪 30%；按非排他重新定价"],
            ["Figure $39B 估值合理", "100 万台需 capex $10-13B；未公开 capex-light", "$39B 隐含未公开确认；sanity check $0.87B", "美系一线 2027 估值回调 50-70%；当前 ≤ 5%，2027 H2 回调后加至 10-15%"],
            ["humanoid IRA 推出概率 40%", "单点概率仅基于 ITIF 单篇白皮书", "降级 policy watchlist；触发条件而非概率", "改为 watchlist（20-50% × Tesla 兑现度）；组合不依赖单一政策"],
        ], page_num=34, total=TOTAL, section_name="商业化悬崖",
        note="变更元 takeaway：10 处修订 7 处指向美系估值下行风险；中系修订相对中性但板块估值锚下移 2-4×")

    # =============== CHAPTER 4: 6 玩家路径 ===============
    slide_section_divider(prs, "CHAPTER 4", "6 种玩家路径",
        "不是\"中 vs 美\"二分 · 是 6 种本质上不同的打法", 35, TOTAL)

    slide_bullets(prs, "路径 ① · Tesla · 垂直整合 + 自家工厂数据飞轮", [
        "{bold}核心假设{/bold}：垂直整合 + 自家工厂数据 + 自研芯片（D1）+ 100 万台/年 → $20K 一台",
        "{bold}实际进度{/bold}：Musk Q4 2025 \"not in usage in a material way\"，2025 部署 ~1,000 台仅用于 \"learn\"",
        "{bold}关键节点{/bold}：Fremont 2026.7-8 启动 V3 量产线，目标 100 万/年",
        "{bold}退出路径{/bold}：Tesla 内部业务，无独立退出",
        "{bold}失败模式{/bold}：100 万台再次推迟；BMW/Mercedes 等 OEM 模式抢先到 PMF；自研芯片落后",
        "{bold}失败概率{/bold}：40-50%（Musk 历史 8 次跨车型 ramp 中 6 次延期 12-24 月）",
        "{bold}PE 视角{/bold}：直接投资不可能（Tesla 子公司），间接看 Tesla 股价",
    ], 36, TOTAL, "玩家路径")

    slide_bullets(prs, "路径 ② · Figure / Apptronik · VLA + B 端 RaaS", [
        "{bold}核心假设{/bold}：纯软件（Helix VLA）+ B 端 RaaS（$1K/月-$25/h）+ 头部车厂 → 估值 5-10× 中系",
        "{bold}实际进度{/bold}：Figure 11 月 BMW PoC 完成 / BotQ 工厂爬产；Apptronik Mercedes/GXO/Jabil pilot；commercial 目标 2026 H2",
        "{bold}关键节点{/bold}：Adcock \"no teleop in market\" 是 binary credibility event",
        "{bold}退出路径{/bold}：Figure IPO 或被 Tesla / NVIDIA 并购",
        "{bold}失败模式{/bold}：RaaS 资产负债表压力；竞争对手提前 PMF；teleop 争议爆发",
        "{bold}失败概率{/bold}：30-40%（estimate 泡沫风险）",
        "{bold}PE 估值{/bold}：Figure $39.5B / Apptronik $5B (A-X $520M 2026.2)；DCF sanity check $0.87B → 市场 {bold}45× sanity{/bold}",
    ], 37, TOTAL, "玩家路径")

    slide_bullets(prs, "路径 ③ · 1X · 消费级 + 远程操作员补丁", [
        "{bold}核心假设{/bold}：双轨——前期靠 teleop 收集真机数据训练 autonomy；2028 真自主率 95%+",
        "{bold}实际进度{/bold}：NEO 2026.4 投产 Hayward；2026 目标 1 万家用；初始自主率 60-70%（CEO 公开）",
        "{bold}关键节点{/bold}：2028 是否真到 95% 自主率？否则 \"running towards a cliff\"（CEO 自己说）",
        "{bold}退出路径{/bold}：1X IPO 或被巨头并购（OpenAI Fund 是投资人）",
        "{bold}失败模式{/bold}：teleop 不能 scale 到家用（隐私 / 文化阻力）；竞争对手出真自主",
        "{bold}失败概率{/bold}：50-60%（最 ambitious 路径）",
        "{bold}PE 视角{/bold}：估值 $10B+（2025.9 在谈，OpenAI/EqualOcean）；2025 营收 ~$20M / PSR ~500×；高赔率高风险",
    ], 38, TOTAL, "玩家路径")

    slide_bullets(prs, "路径 ④ · Physical Intelligence / Skild · 纯模型层", [
        "{bold}核心假设{/bold}：通用大脑模型（π0/π0.5、Skild Brain）+ 跨形态部署 → 软件公司估值（40-100× ARR）",
        "{bold}实际进度{/bold}：PI π0/π0.5 完成 open-world generalization；Skild $1.4B raise 2026.1，估值 $14B",
        "{bold}关键节点{/bold}：能否真做到 SOTA + 跨硬件商业化（vs 自研整机厂商绕过它们）",
        "{bold}退出路径{/bold}：被 NVIDIA / Google / Meta 收购，或独立 IPO",
        "{bold}失败模式{/bold}：整机厂内化模型层（Tesla、Figure 不会买）；Scaling Law 不来",
        "{bold}失败概率{/bold}：40-50%（被整机厂内化）",
        "{bold}PE 视角{/bold}：PI $11B 在谈、Skild $14B 在谈；纯 software multiples，估值最敏感于 hype 周期",
    ], 39, TOTAL, "玩家路径")

    slide_bullets(prs, "路径 ⑤ · 中国硬件派 · 极致硬件成本 + 走量", [
        "{bold}代表玩家{/bold}：宇树 / 众擎 / 加速进化",
        "{bold}核心假设{/bold}：极致 BOM 成本（$5,900 R1 / ¥99K G1）+ 走量（消费+教育+海外 30% 出口）+ 现金流盈利",
        "{bold}实际进度{/bold}：宇树 2025 出货 5,500 台（全口径含机器狗）/ 营收 ¥17 亿 / {bold}扣非净利 ¥6 亿 / 毛利 60%{/bold} — 唯一公开盈利的人形玩家",
        "{bold}关键节点{/bold}：宇树 A 股 IPO（2026.3.20 受理 / 募 ¥42 亿 / 预计市值 > ¥400 亿）",
        "{bold}退出路径{/bold}：A 股上市",
        "{bold}失败模式{/bold}：R1 价格战 → 毛利下行；纯硬件无 software moat",
        "{bold}失败概率{/bold}：15-25%（已盈利，稳健）",
        "{bold}PE 视角{/bold}：宇树 IPO 是中国具身板块估值锚，跟投上限 ¥250 亿；类比 Apple 而非蔚来",
    ], 40, TOTAL, "玩家路径")

    slide_bullets(prs, "路径 ⑥ · 中国场景派 · 政府订单 + 数据采集 + 仿真", [
        "{bold}代表玩家{/bold}：智元 / 银河通用 / 优必选",
        "{bold}核心假设{/bold}：国资背书 + 政府订单 + 数据采集合同 + 仿真合成数据（GraspVLA 10 亿帧）→ 中国弯道超车",
        "{bold}实际进度{/bold}：智元 2025 出货 5,168 台（IDC）/ AgiBot World 100 万+ 开源；银河通用 2026.3 大基金三期 ¥25 亿首投",
        "{bold}关键节点{/bold}：智元港股 IPO 兑现；银河通用 GraspVLA 真实部署任务完成率",
        "{bold}退出路径{/bold}：港股 IPO / A 股 IPO / 战略并购",
        "{bold}失败模式{/bold}：政府订单可持续性（非 ROI 驱动）；仿真 → 真机 gap 不闭合",
        "{bold}失败概率{/bold}：25-35%（依赖政策窗口）",
        "{bold}PE 视角{/bold}：智元 ¥200亿+ / 银河通用 ¥225亿 / 星海图 ¥200亿 / 星动纪元 ¥100亿+",
    ], 41, TOTAL, "玩家路径")

    slide_table(prs, "6 路径对照 · 核心维度矩阵",
        headers=["路径", "核心 moat", "真实经济模型", "2026-28 关键变量", "失败概率"],
        rows=[
            ["① Tesla", "垂直整合 + 自家场景 + 自研芯片", "内部使用，无外部 ROI 数据", "2026.7-8 V3 量产 + $20K 目标", "{red}40-50%{/red}"],
            ["② Figure / Apptronik", "VLA 软件 + B 端客户绑定", "$1K/月 RaaS - 8 年回本 naive", "BotQ 爬产 + capex-light 验证", "{red}30-40%{/red}"],
            ["③ 1X", "消费品牌 + teleop → 数据飞轮", "$20K 硬件 + $499/月 订阅", "2028 95% 自主率兑现", "{red}50-60%{/red}"],
            ["④ PI / Skild", "通用大脑模型", "软件 license / API", "整机厂买不买 vs 自研", "{red}40-50%{/red}"],
            ["⑤ 中国硬件派", "BOM 成本 + 走量 + 现金流", "毛利 30-60%，现金流为正", "R1 海外持续 + IPO 兑现", "{green}15-25%{/green}"],
            ["⑥ 中国场景派", "国资 + 政府订单 + 仿真", "政府采购毛利低，商业溢出未验", "仿真→真机闭合 + IPO 兑现", "{green}25-35%{/green}"],
        ], page_num=42, total=TOTAL, section_name="玩家路径")

    slide_bullets(prs, "6 路径 · 核心 takeaway", [
        "① {bold}没有一条路径有 < 15% 失败率{/bold} — 这是早期产业的本质；PE 应当押\"路径组合\"而非\"单一公司\"",
        "② {bold}当前唯一公开经盈利验证的是路径 ⑤（中国硬件派）{/bold} — 宇树扣非净利 ¥6 亿；其他 5 条路径都未跑出 ROI 闭环",
        "③ {bold}美系 4 路径下 5 家估值合计 ~$79.5B{/bold}，但 2025 真实营收合计估 ~$105M — PSR ~757× 加权平均",
        "④ \"Tesla 会赢\" 这个论点在 10 人样本中无人支持（样本不含 Tesla bulls；带 bulls 后业内分歧 5:5）",
        "⑤ 路径 ① ② ③ ④ ⑥ 一起押的总暴露 = $90-100B，{bold}全部依赖\"AI Scaling Law 2027-2028 兑现\"{/bold} — 如果不到，集体压力测试",
    ], 43, TOTAL, "玩家路径")

    # =============== CHAPTER 5: 产业链（补充） ===============
    slide_section_divider(prs, "CHAPTER 5", "产业链与硬件（补充）",
        "非核心章节 · BOM 速览 / 玩家 Top10 / 市场规模", 44, TOTAL)

    slide_cards(prs, "上游硬件 · 国产化率三梯队", [
        {"label": "第一梯队 > 60%", "title": "已构成护城河", "accent": Theme.SUCCESS,
         "body": ["谐波减速器（绿的 26% / 国产 60-65%）",
                  "灵巧手（国产价 = 海外 1/20）",
                  "六维力（蓝点 70%）",
                  "激光雷达（90%+）",
                  "深度相机（70%+）",
                  "电池（90%+）"]},
        {"label": "第二梯队 30-50%", "title": "加速替代", "accent": Theme.WARN,
         "body": ["无框力矩电机（步科国内 54%）",
                  "RV 减速器（双环 18%）",
                  "AI 端侧芯片（地平线 30-40%）"]},
        {"label": "第三梯队 < 30%", "title": "卡脖子", "accent": Theme.DANGER,
         "body": ["{bold}行星滚柱丝杠（19-22%）{/bold}",
                  "导航级 IMU（< 30%）",
                  "光电编码器（< 30%）",
                  "AI 训练芯片（昇腾 910C 全球封杀）"]},
    ], 45, TOTAL, "产业链")

    slide_table(prs, "单台 BOM 价值量拆解 · 中 vs 美",
        headers=["子系统", "中国 BOM %", "中国金额", "美国 BOM %", "美国金额", "主导供应商"],
        rows=[
            ["执行器总成", "40-55%", "$14-19K", "50-60%", "$60-80K", "谐波/丝杠/电机"],
            ["├ 谐波减速器", "8-12%", "$3-4K", "5-8%", "$6-10K", "Harmonic Drive / 绿的"],
            ["├ 行星滚柱丝杠", "8-15%", "$3-5K", "10-15%", "$12-18K", "GSA / Rollvis / 贝斯特"],
            ["├ 无框力矩电机", "5-8%", "$2-3K", "6-10%", "$8-12K", "Kollmorgen / 步科"],
            ["灵巧手", "15-32%", "$5-11K", "17.2%", "$9.5K", "Tesla 自研 / 月泉 / 灵心"],
            ["传感器", "8-15%", "$3-5K", "8-12%", "$10-15K", "ATI / 蓝点 / 禾赛 / 奥比"],
            ["AI 芯片", "5-8%", "$2-3K", "5-8%", "$5-10K", "Thor / 寒武纪 / 地平线"],
            ["电池", "3-6%", "$1-2K", "3-5%", "$4-6K", "CATL / 亿纬"],
            ["结构 / 装配", "8-12%", "$3-4K", "12-18%", "$15-25K", "—"],
            ["{bold}合计 BOM{/bold}", "100%", "{bold}$32-46K{/bold}", "100%", "{bold}$90-150K{/bold}", "{bold}剪刀差 2.8×{/bold}"],
        ], page_num=46, total=TOTAL, section_name="产业链")

    slide_table(prs, "中游 AI · VLA 模型对比",
        headers=["模型", "团队", "时间", "架构", "开源", "代表能力"],
        rows=[
            ["RT-2", "Google DeepMind", "2023.7", "VLM + 动作 token", "否", "开山之作"],
            ["OpenVLA-7B", "Stanford+UCB+TRI", "2024.6", "Llama2 7B + 动作头", "{green}是{/green}", "开源破局"],
            ["π0 / π0.5", "Physical Intelligence", "2024.10 / 2025.4", "Flow-matching + PaliGemma", "部分", "高频控制 SOTA"],
            ["Helix", "Figure AI", "2025.2", "S1（慢推理）+ S2（快控制） 双系统", "否", "BMW 工厂闭环"],
            ["GR00T N1", "NVIDIA", "2025.3", "VLA + 仿真训练", "部分（N1.5）", "训练框架完整"],
            ["GO-1", "智元 AgiBot", "2025.3", "VLA + AgiBot World 数据", "{green}是{/green}", "中国 ImageNet 时刻"],
            ["GraspVLA", "银河通用", "2025.6", "VLA + 仿真合成 10 亿帧", "否", "全仿真 0 真机"],
            ["SmolVLA 450M", "Hugging Face", "2025.5", "小模型 + 后训练", "{green}是{/green}", "小模型反超 ACT"],
        ], page_num=47, total=TOTAL, section_name="产业链")

    slide_table(prs, "下游 · 全球估值 Top 10 (2026.05)",
        headers=["#", "公司", "国家", "最新估值", "最近一轮", "状态"],
        rows=[
            ["1", "Figure AI", "美", "$39.5B", "Series C 2025", "✅"],
            ["2", "Skild AI", "美", "$14B", "Series B+ 在谈", "🟡"],
            ["3", "Physical Intelligence", "美", "$11B 在谈", "C 在谈", "🟡"],
            ["4", "1X Technologies", "挪/美", "$10B+ 在谈", "2025.9", "🟡"],
            ["5", "Apptronik", "美", "$5B", "A-X 2026.2 ($520M)", "✅"],
            ["6", "Neura Robotics", "德", "€4B", "Series C 2025.7", "✅"],
            ["7", "银河通用 Galbot", "中", "¥225 亿 ($3B+)", "B+ 2026.3 大基金首投", "✅"],
            ["8", "智元 AgiBot", "中", "¥200 亿+", "Pre-IPO 2025-26", "🟡"],
            ["9", "星海图", "中", "¥200 亿", "A 2026", "✅"],
            ["10", "宇树科技", "中", "¥120 亿 → IPO ~¥400 亿", "2026.3.20 受理", "🟡"],
        ], page_num=48, total=TOTAL, section_name="产业链")

    slide_table(prs, "下游 · 出货 Top 10 (2025)",
        headers=["排名", "厂商", "国家", "2025 出货", "关键产品"],
        rows=[
            ["1", "智元 AgiBot", "中", "5,168 台", "远征 A2 / 灵犀 X1"],
            ["2", "宇树科技", "中", "5,500+ 台 (全口径)", "G1 / H1 / R1"],
            ["3", "优必选", "中", "~1,500 台", "Walker S1/S2"],
            ["4", "众擎 EngineAI", "中", "~1,200 台", "PM01 / SA01"],
            ["5", "加速进化", "中", "~800 台", "T1"],
            ["6", "傅利叶", "中", "~700 台", "GR-1 / GR-2"],
            ["7", "Agility Robotics", "美", "~500 台", "Digit V6"],
            ["8", "1X Technologies", "挪/美", "~300 台", "NEO Gamma"],
            ["9", "Apptronik", "美", "~300 台", "Apollo"],
            ["10", "Figure AI", "美", "~150 台 (Tech360 2026.1)", "Figure 02/03"],
        ], page_num=49, total=TOTAL, section_name="产业链",
        note="宇树 5,500 / 智元 5,168 含机器狗+工业版；纯人形中国份额估 70-85%")

    slide_table(prs, "市场规模 · 投行预测对比",
        headers=["机构", "报告", "2030 市场规模", "关键判断"],
        rows=[
            ["Goldman Sachs", "Humanoid Robot 2025", "—", "2035 $380 亿 / > 25 万台"],
            ["Morgan Stanley", "$5T by 2050", "$150 亿", "2050 美国 $3 万亿"],
            ["Citi", "Rise of AI Robots", "—", "2050 6.5 亿台 / $7 万亿"],
            ["BofA", "Humanoid Robots 101", "$300 亿", "2025 BOM $35K / 100 万台/年"],
            ["McKinsey", "Crossing the Chasm", "—", "通用机器人 2040 $3,700 亿"],
            ["MarketsandMarkets", "2025-2030", "$152.6 亿", "CAGR 39.2%"],
            ["BCC Research", "2025", "$110 亿", "CAGR 42.8%"],
            ["{bold}本报告自算{/bold}", "50 万 × $15K (中位)", "{bold}$30-144 亿 / 中位 $75 亿{/bold}", "出货 CAGR 69-113%"],
        ], page_num=50, total=TOTAL, section_name="产业链",
        note="投行预测口径不一致；本报告 TAM 按出货 × ASP 自洽推导")

    slide_table(prs, "渗透率曲线 · 人形 vs EV",
        headers=["年份", "人形累计保有量", "渗透蓝领%", "类比 EV 时点"],
        rows=[
            ["2024", "~0.3 万台", "0.0001%", "EV 2010（Nissan Leaf）"],
            ["2025", "~2.3 万台", "0.001%", "EV 2012"],
            ["2027E", "25-30 万累计", "0.01%", "EV 2015"],
            ["{bold}2030E{/bold}", "{bold}60-150 万累计{/bold}", "{bold}0.04-0.1%{/bold}", "{bold}EV 2018（渗透 2%）{/bold}"],
            ["2035E", "1,500-3,000 万累计", "1-1.5%", "EV 2023（渗透 18%）"],
            ["2040E", "8,000 万-1.5 亿", "4-8%", "EV 2030（渗透 50%+）"],
            ["2050E", "5-8 亿台", "25-40%", "EV 2050（渗透 90%）"],
        ], page_num=51, total=TOTAL, section_name="产业链",
        note="关键差异：人形 TAM 比 EV 大 5-10×（替代劳动力 vs 替代车辆）；渗透曲线更陡（替代持续人工成本）")

    slide_table(prs, "2024-2026 全球投资项目清单 · Top 10",
        headers=["#", "公司", "项目", "时间", "金额", "投产"],
        rows=[
            ["1", "智元", "上海临港 4,000㎡ 数据工厂", "2024.12", "~¥3 亿", "2025.1"],
            ["2", "宇树", "杭州总部 + 研发基地 (IPO 募投)", "2026.3", "{bold}¥42 亿{/bold}", "2027"],
            ["3", "优必选", "柳州工业人形机器人工厂", "2025.10", "—", "2026 Q1"],
            ["4", "Tesla", "Fremont Optimus V3 量产线", "2026.5 启动", "—", "{bold}2026.7-8{/bold}"],
            ["5", "Figure AI", "BotQ Austin 工厂", "2025.3", "—", "一线 1.2 万台/年"],
            ["6", "Apptronik", "Austin 扩建", "2026.2", "{bold}$520M{/bold}", "2026.H2"],
            ["7", "1X", "Hayward NEO 工厂（垂直一体化）", "2026.4.30", "—", "2026.5 起"],
            ["8", "银河通用", "北京亦庄 + Galbot S1 产线", "2026.3", "¥25 亿（大基金三期）", "2026"],
            ["9", "宁德×智元×千寻", "中州基地\"电池厂+人形\"产线", "2025.12", "—", "已投产"],
            ["10", "NEOM × 沙特 PIF", "一期 10 万台部署", "2024-2025", "—", "2025-2027"],
        ], page_num=52, total=TOTAL, section_name="产业链")

    slide_table(prs, "基础原材料 · 一张表速览",
        headers=["原材料", "单台用量", "单台成本", "占 BOM", "中国地位"],
        rows=[
            ["钕铁硼磁材", "3 kg", "1,800 元", "0.5%", "{bold}磁材 92%+{/bold}"],
            ["铜（绕组+线缆）", "8.2 kg", "600 元", "0.2%", "精炼 45%"],
            ["镁/铝/钢", "26 kg", "650 元", "0.23%", "原镁 87% / 电解铝 58%"],
            ["PEEK", "7 kg", "700 元", "0.2%", "中研 30.6% / 全球第 4"],
            ["碳纤维 T800+", "5 kg", "4,000 元", "1.3%", "中复神鹰 全球第 3"],
            ["{red}铍铜（卡脖子）{/red}", "0.08 kg", "50 元", "0.02%", "{red}美 Materion 70% 主导{/red}"],
            ["{bold}合计{/bold}", "—", "~7,800 元", "{bold}~2.5%{/bold}", "—"],
        ], page_num=53, total=TOTAL, section_name="产业链",
        note="注：上游材料目前不是瓶颈，PE 报告处理为\"主题观察项\"而非核心配置项")

    # =============== CHAPTER 6: 中美格局与政策 ===============
    slide_section_divider(prs, "CHAPTER 6", "中美格局与政策",
        "16 环节能力对比 · 不可能三角 · 三情景 · 政策 · 补贴", 54, TOTAL)

    slide_table(prs, "中美产业链 16 环节能力对比",
        headers=["环节", "中国", "美国", "优势", "关键差异"],
        rows=[
            ["端侧 AI 芯片", "★★★", "★★★★★", "美", "Thor 仍可对华"],
            ["训练算力", "★★", "★★★★★", "美绝对", "H100/B200 受限"],
            ["谐波减速器", "★★★★", "★★★", "中", "绿的 ASP 低 30-40%"],
            ["行星滚柱丝杠", "★★", "★★★★", "美", "中国国产化 19-22%"],
            ["灵巧手", "★★★★", "★★★★", "中（成本）", "国产价 = 海外 1/20"],
            ["力 / 触觉传感", "★★★★", "★★★★", "中略胜", "蓝点 70%"],
            ["稀土永磁", "★★★★★", "★", "中绝对", "92%+ 高性能磁材"],
            ["PEEK / 碳纤维", "★★★★", "★★★", "中略胜", "中研股份 Tesla 供应商"],
            ["铍铜", "★", "★★★★★", "美绝对", "Materion 70%"],
            ["VLA 模型", "★★★", "★★★★★", "美", "Helix/π0/GR00T 引领"],
            ["仿真平台", "★★", "★★★★★", "美绝对", "NVIDIA Isaac 独占"],
            ["训练数据", "★★★★", "★★★★", "中略胜", "AgiBot World 100 万+"],
            ["整机量产", "★★★★★", "★★★", "中绝对", "出货 80-95%，BOM 2.8×"],
            ["客户场景", "★★★★★", "★★★★", "中略胜", "政府/车厂大单密度高"],
            ["资本市场", "★★★", "★★★★★", "美", "Figure $39B vs 智元 $3B"],
        ], page_num=55, total=TOTAL, section_name="中美格局")

    slide_comparison(prs, "双维度 不可能三角",
        left={"title": "维度 A · PE 视角", "accent": Theme.GOLD, "body": [
            "{bold}\"国家安全 × 技术领先 × 经济效益\"{/bold}",
            "",
            "美国偏好：{bold}安全 + 技术{/bold}（出口管制 + 资本人才）",
            "中国偏好：{bold}安全 + 经济{/bold}（自主可控 + 成本敏感）",
            "欧洲偏好：经济 + 部分技术",
            "",
            "类比 Reference 在电池产业链做的\"国家安全 × 环境保护 × 经济效益\"",
        ]},
        right={"title": "维度 B · 技术视角", "accent": Theme.INFO, "body": [
            "{bold}\"数据规模 × 泛化能力 × 商业化成本\"{/bold}",
            "",
            "Tesla / Figure 路径：{bold}海量真机数据 + 高泛化{/bold} → 高成本（$130K）",
            "银河通用路径：{bold}全仿真合成 + 中等泛化{/bold} → 低成本（$30K）",
            "宇树 / 众擎路径：{bold}低成本 + 低泛化（专用场景）{/bold} → 走量",
            "",
            "任何路线只能优化两个",
        ]},
        page_num=56, total=TOTAL, section_name="中美格局")

    slide_table(prs, "三种未来情景 · 详细对照",
        headers=["维度", "China-led (35%)", "US-led (15-25%)", "Bifurcation (40-50%, 基准)"],
        rows=[
            ["触发条件", "中国 2027 ramp + 美无 IRA", "Tesla 100 万兑现 + IRA 推出", "双方推动本土化但未脱钩"],
            ["时间表", "2026-2030 中国持续主导", "2027-2030 美国反超", "2026 起逐步分化"],
            ["2030 累计保有量", "中 100 万 / 美 12 万 / 其他 8 万", "中 30 万 / 美 80 万 / 其他 10 万", "{bold}中 60 万 / 美 25 万 / 其他 5 万 = 90 万{/bold}"],
            ["2030 ASP", "$10K (中) / $35K (美)", "$20K（双方收敛）", "$12K (中) / $40K (美)"],
            ["主要受益方", "中系全链 + 海外硬件依赖中国", "美系一线 + FTA 国家硬件 + NVIDIA", "各国\"双轨\"领导者；NVIDIA+CATL+稀土玩家双赢"],
            ["最大风险", "中国本土洗牌过度", "中国失出口市场 + 算力管制深化", "全球效率损失"],
            ["PE 行动", "重押中系 + 海外硬件", "押美系 + FTA + NVIDIA", "双线对冲"],
        ], page_num=57, total=TOTAL, section_name="中美格局")

    slide_comparison(prs, "政策环境 · 中美对照",
        left={"title": "🇺🇸 美国", "accent": Theme.INFO, "body": [
            "2025.1 Trump EO 14179 取代 Biden 14110",
            "2025.5 撤销 AI Diffusion Rule",
            "{bold}2025.12 H200 转 case-by-case + 25% 关税{/bold}（部分放开）",
            "B200/GB300 维持 presumption of denial",
            "2025.9 CA SB 53 签署；NY RAISE Act 在审",
            "ITIF 2025.7 呼吁推 \"humanoid IRA\"",
            "{bold}推出概率（policy watchlist）：20-50% × Optimus 兑现度{/bold}",
        ]},
        right={"title": "🇨🇳 中国", "accent": Theme.DANGER, "body": [
            "工信部 2023.11 人形机器人创新发展指导意见",
            "2024 工信部具身智能标委会（2026.4 正式成立）",
            "{bold}大基金三期 2026.3 首投具身（银河通用 ¥25 亿）{/bold}",
            "北京/上海/深圳/合肥/武汉 各 ¥100 亿+ 地方基金",
            "{bold}累计承诺 ¥550+ 亿{/bold}",
            "5 年年化 ¥110 亿/年 ≈ 0.08% GDP",
            "{bold}仅为 EV 2015 高峰 (0.51% GDP) 的 15-20%{/bold}",
        ]},
        page_num=58, total=TOTAL, section_name="中美格局")

    slide_table(prs, "中国 550 亿+ 补贴 · 省级分布",
        headers=["城市 / 省份", "政策 / 基金", "时间", "规模"],
        rows=[
            ["中央 · 大基金三期", "首投具身智能", "2026.3", "银河通用 ¥25 亿领投"],
            ["北京（市+海淀+亦庄）", "机器人产业基金", "2024.12", "首期 ¥20 亿，总 {bold}¥100 亿{/bold}"],
            ["北京", "具身智能行动计划 2025-2027", "2025.3", "—"],
            ["上海（市+临港+张江）", "上海具身智能基金", "2025.4", "首期 ¥5.6 亿，目标 ¥10 亿"],
            ["上海", "国投先导 AI 母基金", "2025", "¥225 亿（覆盖部分）"],
            ["深圳（市+南山+福田+龙岗）", "AI 和具身机器人基金", "2025.5", "首期 ¥20 亿，目标 {bold}¥100 亿{/bold}"],
            ["深圳", "三只 AI/机器人基金", "2025.10", "¥30 亿"],
            ["合肥 / 安徽", "智能机器人行动计划", "2024.9", "¥50 亿专项 + {bold}¥100 亿{/bold} 未来基金"],
            ["武汉", "加快人形机器人发展政策", "2025.4", "{bold}¥100 亿{/bold} 招商基金"],
            ["江苏省", "战新产业母基金", "2024-2025", "¥500 亿（覆盖机器人部分）"],
        ], page_num=59, total=TOTAL, section_name="中美格局",
        note="具身专项合计 ¥550+ 亿（不含 AI 母基金间接出资）；2025-2027 三年可调用资金 ¥800-1000 亿")

    slide_cards(prs, "估值方法论 · 4 套", [
        {"label": "① Seed-A", "title": "VC Scorecard", "accent": Theme.SUCCESS,
         "body": ["团队 30% + 市场 25% + 产品 15%", "+ 客户 15% + 竞争 15%",
                  "", "{bold}当前行业{/bold}", "中位 pre-money $50-200M"]},
        {"label": "② Pre-rev", "title": "Berkus", "accent": Theme.INFO,
         "body": ["5 维度 × 各 $0.5-2M", "", "适用于早期 / 种子轮"]},
        {"label": "③ Series B+", "title": "EV / Sales 倍数", "accent": Theme.WARN,
         "body": ["已订单 × 倍数", "", "{bold}当前行业{/bold}",
                  "中系 10-20×", "美系 30-80×"]},
        {"label": "④ Series C+", "title": "DCF on 2030 命中产能", "accent": Theme.DANGER,
         "body": ["年产能 × ASP × 毛利率 × 折现率",
                  "", "{bold}锚定 2030 真实兑现{/bold}"]},
    ], 60, TOTAL, "中美格局")

    slide_table(prs, "估值 sanity check · Figure / 智元 / 宇树",
        headers=["参数", "Figure 🇺🇸", "智元 🇨🇳", "宇树 🇨🇳"],
        rows=[
            ["2030 出货假设", "5 万台（10% 全球）", "12 万台（24% 全球）", "14 万台（含机器狗）"],
            ["2030 ASP", "$25K", "¥85K ($12K)", "¥60K ($8.5K)"],
            ["2030 收入", "$1.25B", "¥104 亿 ($1.44B)", "¥84 亿 ($1.19B)"],
            ["2030 毛利率", "40%", "30%", "35%"],
            ["2030 EBIT", "13% = $163M", "3% = $43M", "8% = $95M"],
            ["WACC", "15%", "18%", "16%"],
            ["Operating value (4y)", "$93M", "$22M", "$52M"],
            ["永续期价值", "$0.78B", "$0.15B", "$0.40B"],
            ["{bold}Sanity check 估值{/bold}", "{bold}$0.87B{/bold}", "{bold}$0.17B{/bold}", "{bold}$0.46B{/bold}"],
            ["当前市场估值", "$39.5B", "$2.8B", "$5.6B IPO"],
            ["{bold}市场 / sanity 比值{/bold}", "{red}45×{/red}", "16×", "12×"],
        ], page_num=61, total=TOTAL, section_name="中美格局",
        note="三家合计 $3.88B / TAM 中位 $7.5B = 52% 集中度合理 · 这不是完整 DCF（未含 2026-30 逐年现金流、capex、税、稀释）；仅作 sanity check")

    slide_table(prs, "Comps · 可比公司估值倍数 Quartile",
        headers=["公司 / 板块", "样本时点", "EV / Sales LTM", "Quartile", "vs 板块中位"],
        rows=[
            ["{bold}Figure AI 🇺🇸{/bold}", "2025.9 Series C", "{red}790×{/red} ($39.5B/$50M)", "{red}Q4 极端高位{/red}", "—"],
            ["{bold}智元 AgiBot 🇨🇳{/bold}", "2025 Pre-IPO", "20× (¥200亿/¥10亿)", "Q3（EV IPO 中位）", "—"],
            ["{bold}宇树 🇨🇳{/bold}", "2026.3 IPO", "23× (¥400亿/¥17.08亿)", "Q3", "—"],
            ["蔚来（IPO 2018.9）", "2018.9", "53×", "Q4", "+141%"],
            ["小鹏（IPO 2020.8）", "2020.8", "24×", "Q3", "+9%"],
            ["ABB", "2025", "2.8×", "Q2", "-20%"],
            ["Symbotic", "2025", "4.5×", "Q3", "+29%"],
            ["Intuitive Surgical", "2025", "14×", "Q4", "+300%"],
            ["OpenAI", "2024.10", "~50×", "Q3", "+25%"],
            ["Anthropic", "2025.3", "~40×", "Q3", "0%"],
        ], page_num=62, total=TOTAL, section_name="中美格局",
        note="关键定位：Figure 790× 是真泡沫；智元/宇树 20-23× 接近 EV IPO 中位（22×），不是泡沫 · 系统性泡沫只在美系一线")

    # 2026 Q1-Q2 events
    slide_bullets(prs, "2026 Q1-Q2 关键事件年表", [
        "{bold}2026.1{/bold} · 优必选 Walker S2 进入比亚迪/吉利 · 累计订单 ¥14-15 亿 · CES 2026: Hyundai AI 战略 · Atlas 接入 Gemini Robotics",
        "{bold}2026.2{/bold} · Apptronik A-X $520M（估值 $5B / 累计 $935M）· 国内月度融资 ¥160+ 亿（新高）· 1X 公布 NEO 量产路线",
        "{bold}2026.3{/bold} · 银河通用 ¥25 亿 B+（大基金三期首投）· 星动纪元 ¥10 亿 A+ · 逐际动力 $2 亿 B · {bold}宇树科创板 IPO 受理（3.20，募 ¥42.02 亿，预计市值 > ¥400 亿）{/bold}",
        "{bold}2026.4{/bold} · Q1 国内具身融资 ¥200+ 亿（YoY +60%）· 工信部具身智能标委会正式成立 · 云深处科技完成上市辅导 · 1X 加州 Hayward NEO 工厂投产（4.30）",
        "{bold}2026.5{/bold} · 优必选 × 空中客车签 Walker S2 采购 · Figure 03 量产爬坡：产线效率 24× 跃升 · Tesla Fremont 终止 Model S/X，7-8 月启动 Optimus V3 量产线",
        "{bold}YTD 累计{/bold}：国内具身累计融资 ¥373 亿 · YTD 日均 ¥2.5 亿 · 2026.4 单月日均 ¥5.8 亿创新高（主要由银河 ¥25 亿 + 星动 ¥10 亿 + 逐际 $2 亿驱动）",
    ], 63, TOTAL, "中美格局")

    # =============== APPENDIX A ===============
    slide_section_divider(prs, "APPENDIX A", "PE 内部分析工具",
        "deal-level 工具 · 与正文行业研究分开", 64, TOTAL)

    slide_cards(prs, "A.1 · 综合判断三段论", [
        {"label": "现状", "title": "非对称稳态格局", "accent": Theme.INFO,
         "body": ["中国主导出货（80-95%）+ 硬件 BOM 优势（2.8×）",
                  "美国主导模型 + 估值",
                  "",
                  "{bold}双向卡脖子{/bold}",
                  "美→中：H100/B200 + 铍铜",
                  "中→美：稀土 92% + 镁 87% + 电池 90% + PEEK 供应"]},
        {"label": "2027-2028", "title": "关键验证窗口 · 5 变量", "accent": Theme.WARN,
         "body": ["1. Tesla Optimus 100 万/$20K 是否兑现？",
                  "2. 中系头部 IPO 兑现 + 单台经济性？",
                  "3. 中国 550+ 亿是否催生半导体级洗牌？",
                  "4. humanoid IRA 是否推出？",
                  "5. VLA Scaling Law 时刻？"]},
        {"label": "2030-2050", "title": "世代级机会", "accent": Theme.SUCCESS,
         "body": ["{bold}全球累计{/bold}：60-150 万（2030）→ 5-8 亿（2050）",
                  "",
                  "{bold}年度 TAM{/bold}：$30-144 亿（2030 中位 $75 亿）→ $5-7 万亿（2050）",
                  "",
                  "蓝领劳动力替代 = EV / 光伏 / AI 同等量级"]},
    ], 65, TOTAL, "附录 A")

    slide_table(prs, "A.2 · PE 投资优先级矩阵",
        headers=["优先级", "投资方向", "代表标的", "估值锚", "时间窗"],
        rows=[
            ["★★★★★", "基础原材料卡点", "金力永磁（Optimus 供应商）、中研股份、中复神鹰、宝武镁业", "已上市 PE 30-60×", "2026 H1 起"],
            ["★★★★★", "中系卡脖子零部件", "行星滚柱丝杠、六维力（蓝点）、视触觉（戴盟）、灵巧手、IMU、无框电机", "未上市 pre-money $50M-$1B", "2026 H1 起"],
            ["★★★★", "仿真 / 数据基础设施", "国内对标 NVIDIA Isaac、数据工厂（智元 4,000㎡）", "早期 + AI 估值溢价", "2026 起"],
            ["★★★★", "中系一线整机（IPO 锚定）", "宇树（2026.3.20 受理）、智元、银河通用、星动纪元", "¥100-400 亿", "2026 IPO 窗口"],
            ["★★★", "美系一线整机（限观察）", "Figure、Apptronik、1X、Skild、π0", "$5-40B", "2027 后调整窗口"],
            ["★★", "美系芯片 / 模型生态", "NVIDIA、Skild、π0、Generalist AI", "高估值", "长期持有"],
            ["★★", "应用场景集成商", "车厂自研子、物流 RaaS、医疗康复", "多元", "2026-2028"],
            ["★", "全栈大而全", "二线整机", "—", "—"],
        ], page_num=66, total=TOTAL, section_name="附录 A")

    slide_cards(prs, "A.3 · PE 组合配置建议", [
        {"label": "人民币基金", "title": "中系重权", "accent": Theme.DANGER,
         "body": ["原材料卡点：15%", "中系零部件：25%", "中系整机：25%",
                  "美系软件/模型：5%", "美系整机：5%",
                  "仿真/数据：15%", "应用层：10%"]},
        {"label": "美元基金", "title": "美系重权", "accent": Theme.INFO,
         "body": ["原材料卡点：5%", "中系零部件：10%", "中系整机：10%",
                  "美系软件/模型：30%", "美系整机：25%",
                  "仿真/数据：15%", "应用层：5%"]},
        {"label": "双货币基金", "title": "平衡配置", "accent": Theme.GOLD,
         "body": ["原材料卡点：10%", "中系零部件：20%", "中系整机：20%",
                  "美系软件/模型：15%", "美系整机：15%",
                  "仿真/数据：15%", "应用层：5%"]},
    ], 67, TOTAL, "附录 A")

    slide_cards(prs, "A.3 · 关键决策树", [
        {"label": "判断 1", "title": "China-led 35%？", "accent": Theme.DANGER,
         "body": ["{bold}是 →{/bold}",
                  "重押中系全产业链",
                  "+ 海外硬件依赖中国",
                  "",
                  "否 → 进入判断 2"]},
        {"label": "判断 2", "title": "Bifurcation 40-50%？", "accent": Theme.GOLD,
         "body": ["{bold}是 →{/bold}",
                  "双线对冲",
                  "中系硬件+原材料",
                  "+ 美系软件 + NVIDIA",
                  "",
                  "否 → 进入判断 3"]},
        {"label": "判断 3", "title": "US-led 15-25%？", "accent": Theme.INFO,
         "body": ["{bold}是 →{/bold}",
                  "重押 Tesla / Figure / Apptronik",
                  "+ NVIDIA / FTA 硬件",
                  "",
                  "否 → 观望情景"]},
        {"label": "默认", "title": "观望情景", "accent": Theme.TEXT_MUTED,
         "body": ["保持现金",
                  "+ 少量种子布局",
                  "（丝杠/灵巧手/IMU/仿真/数据）",
                  "",
                  "等 2027-2028 兑现窗口"]},
    ], 68, TOTAL, "附录 A")

    slide_bullets(prs, "A.4 · 投委会必问 Q&A · Top 5 (1/2)", [
        "{bold}Q1{/bold}：出货 85% 是教育/数据采集，工业落地 < 5%，是补贴催生的 GDP 数字？",
        "答：① 教育采购量大但单价低；数据采集是 VLA 训练真实需求 ② 工业 PoC 已闭环（比亚迪 2026 计划 20K 台、空客签约）③ 美国 BMW-Figure / GXO-Agility 真实 RaaS 收入；但承认中国工业 PoC → 规模化仍需 12-18 月验证",
        "",
        "{bold}Q2{/bold}：Tesla 部署 1,000 台 useful work 受质疑，2026 V3 100 万产能不兑现，US-led 概率降到多少？",
        "答：base case Tesla 2030 = 30/200/100 万；如 2026 Q3 < 10K 台/年，US-led 从 25% → 15%，Bifurcation 升至 50%",
        "",
        "{bold}Q3{/bold}：LP 大半在美国，投中系标的政治退出风险如何对冲？humanoid IRA + 中概股退市怎么办？",
        "答：① 美元基金聚焦\"中国制造但海外销售\"零部件 ② QFII/沪深港通 ③ 二级对冲（做空 Figure / 做多中国零部件）④ 跨境 SPV ⑤ humanoid IRA 出台后美系硬件供应商相对受益",
    ], 69, TOTAL, "附录 A")

    slide_bullets(prs, "A.4 · 投委会必问 Q&A · Top 5 (2/2)", [
        "{bold}Q4{/bold}：宇树/智元 IPO 估值 ¥120-400 亿 是合理还是泡沫？vs 蔚来 IPO？",
        "答：宇树 2025 营收 ¥17 亿 + 扣非净利 ¥6 亿 + 毛利 60.27%，是同期蔚来（亏损）没有的；毛利结构类似 Apple 而非蔚来；DCF sanity $0.46B vs IPO $5.6B = 12× 比值，跟投上限 ¥250 亿为安全锚",
        "",
        "{bold}Q5{/bold}：如果让你只投 3 个标的，押 5 年 IRR > 30%，是哪 3 个？",
        "答：",
        "① {bold}宇树{/bold}（A 股 IPO 受理，毛利 60%，DCF 隐含 PSR 锚定）— IRR 5y 35-45%",
        "② {bold}金力永磁{/bold}（Tesla Optimus 稀土永磁供应商，3-5 年窗口期，Tesla 锁定保护）— IRR 5y 30-40%",
        "③ {bold}仿真/数据基础设施种子轮{/bold}（国内对标 NVIDIA Isaac、AgiBot World 类开源数据工厂）— IRR 5y 50%+（10× potential）",
    ], 70, TOTAL, "附录 A")

    slide_bullets(prs, "A.5 · IC Memo 样本 · 智元机器人 (1/2)", [
        "{bold}公司{/bold}：智元机器人 AgiBot（北京 / 上海）",
        "{bold}建议{/bold}：{green}CONDITIONAL PROCEED{/green} ── Lead 早期布局，限 Pre-IPO 轮持仓 ≤ 基金规模 3%",
        "",
        "{bold}1. Executive Summary{/bold}：中国 2026 估值最高人形整机厂（约 $2.8B+，IPO 锚 $4-7B），创始人彭志辉。2025 首批量产 5,168 台（IDC）。Top 3 risks：① 量产爬坡未经第三方审计 ② 估值已 frothy（但 PSR 20× 接近 EV IPO 中位，相对合理）③ 软件栈相对 Helix/π0 优势不明显",
        "",
        "{bold}2. Company Overview{/bold}：远征 A2 (¥620K) / A2 青春版 (¥168K) / 灵犀 X1 (¥109K)；投资方：腾讯/上汽/京东/比亚迪/红杉/北京机器人基金/上海国资",
        "{bold}3. Industry & Market{/bold}：占中国 2025 出货 40%、全球 ~30%，Top 1 by IDC",
        "{bold}4. Financial Analysis{/bold}：[UNSOURCED]——智元未披露完整财务；DD 阶段必拿数据室项",
    ], 71, TOTAL, "附录 A")

    slide_table(prs, "A.5 · IC Memo 样本 · 智元 Returns Analysis (2/2)",
        headers=["Returns Analysis (5y hold)", "触发", "Exit 估值", "Multiple", "IRR", "MOIC"],
        rows=[
            ["Bull (China-led 35%)", "2027 量产 5 万 + 工业商业化兑现", "$35B IPO", "5.0×", "{green}38%{/green}", "5.0×"],
            ["Base (Bifurcation 40%)", "2027 量产 2 万 + Tesla 部分兑现", "$18B", "2.6×", "{green}21%{/green}", "2.6×"],
            ["Bear (US-led 25%)", "2027 量产 < 10K + 美 IRA 出台", "$5B", "0.7×", "{red}-6%{/red}", "0.7×"],
            ["{bold}概率加权（35/40/25）{/bold}", "—", "—", "—", "{bold}~20%{/bold}", "{bold}3.0×{/bold}"],
        ], page_num=72, total=TOTAL, section_name="附录 A",
        note="风险（排序+mitigants）：① 量产 yield — 季度 yield report + 第三方审计权 + ratchet ② 估值泡沫 — IPO < $5B 触发反稀释 ③ 关键人风险 — 4 年 vesting ④ 政策反复 — 海外 GR + FTA hedging · 建议：CONDITIONAL PROCEED 跟投 $20-30M")

    # =============== APPENDIX B: References ===============
    slide_section_divider(prs, "APPENDIX B", "数据来源与引用追溯",
        "206 条引用 · Tier 1/2/3 分级 · 置信度声明", 73, TOTAL)

    slide_cards(prs, "数据来源 · 206 条引用按 Tier 分级", [
        {"label": "Tier 1 · 39%", "title": "一手 / 权威机构 · 81 条", "accent": Theme.SUCCESS,
         "body": ["Goldman Sachs · Morgan Stanley · Citi · BofA",
                  "McKinsey · Bain · IDC · IFR",
                  "USGS · 工信部 · 中国信通院 · GGII",
                  "公司公告 · 招股说明书",
                  "arXiv 论文 · BIS 文件"]},
        {"label": "Tier 2 · 41%", "title": "券商 / Bloomberg · 84 条", "accent": Theme.INFO,
         "body": ["中信建投 / 东吴 / 华泰 / 招商",
                  "东方 / 安信 / 国信 / 华宝 / 中泰 / 国金 / 民生",
                  "Bloomberg · Reuters · The Information",
                  "第一财经 · 21 经济网 · 华尔街见闻 · 钛媒体"]},
        {"label": "Tier 3 · 20%", "title": "行业媒体 · 41 条", "accent": Theme.WARN,
         "body": ["36 氪 · 量子位 · 机器之心 · 智东西",
                  "艾邦机器人 · 机器人大讲堂",
                  "知乎专栏 · 行业网站",
                  "",
                  "210 条数据点 / 81 公司档案",
                  "完整 JSON 数据见 data/references.json"]},
    ], 74, TOTAL, "附录 B")

    slide_table(prs, "数据置信度声明 · 已知 limitations",
        headers=["数据点", "报告引用值", "置信度", "备注"],
        rows=[
            ["2025 全球出货", "1.8 万台", "{green}High{/green}", "IDC、信通院、GGII 三方交叉"],
            ["中国出货占全球", "80-95%（按口径区间）", "{orange}Medium{/orange}", "IDC ~95% / GGII ~85%（口径差异）"],
            ["中美 BOM 剪刀差", "$46K vs $130K", "{orange}Medium{/orange}", "MS Humanoid 100；BofA $35K 为压制后口径"],
            ["Figure AI 估值", "$39.5B", "{green}High{/green}", "Figure 官方 + The Information"],
            ["Skild AI $14B", "$14B 在谈", "{orange}Medium{/orange}", "The Information 报道；未落定"],
            ["智元 2025 出货 5,168", "5,168 台", "{orange}Medium{/orange}", "IDC 转引；IDC 原报告未公开"],
            ["三情景概率", "35 / 15-25 / 40-50", "{red}Analyst Judgment{/red}", "非市场共识；需 reviewer 校准"],
            ["humanoid IRA", "Policy watchlist", "{red}Watchlist{/red}", "不给单点概率"],
            ["稀土高性能磁材 92%", "92%+", "{green}High{/green}", "USGS + 中国稀土行业协会"],
            ["中国 550 亿+ 补贴", "累计承诺", "{orange}Medium{/orange}", "\"承诺\" vs \"实际投放\" 有差距"],
        ], page_num=75, total=TOTAL, section_name="附录 B")

    # =============== CLOSING ===============
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, Theme.NAVY_DARK)
    add_rect(s, Inches(0), Inches(0), Inches(0.3), Theme.H, Theme.GOLD)
    add_text(s, Inches(1), Inches(2.5), Inches(11.5), Inches(1.5),
             "本报告对外严格保密", font_size=40, bold=True,
             color=Theme.TEXT_INVERT, align=PP_ALIGN.LEFT)
    add_text(s, Inches(1), Inches(3.5), Inches(11.5), Inches(0.8),
             "仅供 PE 内部决策使用", font_size=22, color=Theme.GOLD,
             align=PP_ALIGN.LEFT)
    add_text(s, Inches(1), Inches(5.0), Inches(11.5), Inches(0.5),
             "数据截止：2026 年 5 月  |  下次更新：随关键 catalyst 触发", font_size=14,
             color=Theme.ICE_BLUE, align=PP_ALIGN.LEFT)
    add_text(s, Inches(1), Inches(5.5), Inches(11.5), Inches(0.5),
             "（Tesla V3 量产 / 中系 IPO / VLA Scaling Law）", font_size=12,
             color=Theme.ICE_BLUE, align=PP_ALIGN.LEFT, italic=True)
    add_text(s, Inches(1), Inches(6.5), Inches(11.5), Inches(0.4),
             "总产出：6 万中文字 · 75 张图表 · 206 引用 · 81 家公司档案", font_size=11,
             color=Theme.GOLD, align=PP_ALIGN.LEFT)

    return prs

if __name__ == "__main__":
    prs = build_deck()
    output = "/Users/bytedance/Downloads/claude-cowork/embodied_ai_report/embodied_ai_deck.pptx"
    prs.save(output)
    print(f"✅ Saved: {output}")
    print(f"Total slides: {len(prs.slides)}")
